from __future__ import annotations
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any

import fitz  # PyMuPDF
from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger
from dataflow_agent.state import KBMindMapState, MainState
from dataflow_agent.agentroles import create_simple_agent
from dataflow_agent.utils import get_project_root
from langchain_core.messages import HumanMessage

log = get_logger(__name__)

# Try importing office libraries
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

@register("kb_mindmap")
def create_kb_mindmap_graph() -> GenericGraphBuilder:
    """
    Workflow for Knowledge Base MindMap Generation
    Steps:
    1. Parse uploaded files (PDF/Office)
    2. Analyze content structure using LLM
    3. Generate Mermaid mindmap syntax using LLM
    """
    builder = GenericGraphBuilder(state_model=KBMindMapState, entry_point="_start_")

    def _start_(state: KBMindMapState) -> KBMindMapState:
        # Ensure request fields
        if not state.request.files:
            state.request.files = []

        # Initialize output directory
        if not state.result_path:
            project_root = get_project_root()
            import time
            ts = int(time.time())
            email = getattr(state.request, 'email', 'default')
            output_dir = project_root / "outputs" / "kb_outputs" / email / f"{ts}_mindmap"
            output_dir.mkdir(parents=True, exist_ok=True)
            state.result_path = str(output_dir)

        state.file_contents = []
        state.content_structure = ""
        state.mermaid_code = ""
        state.mindmap_svg_path = ""
        return state

    async def parse_files_node(state: KBMindMapState) -> KBMindMapState:
        """
        Parse all files and extract content
        """
        files = state.request.files
        if not files:
            state.file_contents = []
            return state

        async def process_file(file_path: str) -> Dict[str, Any]:
            file_path_obj = Path(file_path)
            filename = file_path_obj.name

            if not file_path_obj.exists():
                return {
                    "filename": filename,
                    "content": f"[Error: File not found {file_path}]"
                }

            suffix = file_path_obj.suffix.lower()
            raw_content = ""

            try:
                # PDF
                if suffix == ".pdf":
                    try:
                        doc = fitz.open(file_path)
                        text = ""
                        for page in doc:
                            text += page.get_text() + "\n"
                        raw_content = text
                    except Exception as e:
                        raw_content = f"[Error parsing PDF: {e}]"

                # Word
                elif suffix in [".docx", ".doc"]:
                    if Document is None:
                         raw_content = "[Error: python-docx not installed]"
                    else:
                        try:
                            doc = Document(file_path)
                            raw_content = "\n".join([p.text for p in doc.paragraphs])
                        except Exception as e:
                             raw_content = f"[Error parsing Docx: {e}]"

                # PPT
                elif suffix in [".pptx", ".ppt"]:
                    if Presentation is None:
                        raw_content = "[Error: python-pptx not installed]"
                    else:
                        try:
                            prs = Presentation(file_path)
                            text = ""
                            for i, slide in enumerate(prs.slides):
                                text += f"--- Slide {i+1} ---\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text += shape.text + "\n"
                            raw_content = text
                        except Exception as e:
                            raw_content = f"[Error parsing PPT: {e}]"

                else:
                    try:
                        with open(file_path, "r", encoding="utf-8") as f:
                            raw_content = f.read()
                    except:
                        raw_content = "[Unsupported file type]"

            except Exception as e:
                 raw_content = f"[Parse Error: {e}]"

            # Truncate content
            truncated_content = raw_content[:50000] if len(raw_content) > 50000 else raw_content

            return {
                "filename": filename,
                "content": truncated_content
            }

        # Run in parallel
        tasks = [process_file(f) for f in files]
        results = await asyncio.gather(*tasks)

        state.file_contents = results
        return state

    async def analyze_structure_node(state: KBMindMapState) -> KBMindMapState:
        """
        Analyze content structure using LLM
        """
        if not state.file_contents:
            state.content_structure = "No content available for analysis."
            return state

        # Format file contents
        contents_str = ""
        for item in state.file_contents:
            contents_str += f"=== {item['filename']} ===\n{item['content']}\n\n"

        # Structure analysis prompt
        language = state.request.language
        max_depth = state.request.max_depth
        prompt = f"""你是一位专业的知识结构分析师。请分析以下文档内容，提取出层级化的知识结构。

要求：
1. 识别主要主题和子主题
2. 提取关键概念和要点
3. 建立清晰的层级关系（最多{max_depth}层）
4. 使用{language}语言
5. 输出格式为层级化的文本结构，使用缩进表示层级

文档内容：
{contents_str}

请输出层级化的知识结构："""

        try:
            agent = create_simple_agent(
                name="structure_analyzer",
                model_name=state.request.model,
                chat_api_url=state.request.chat_api_url,
                temperature=0.3
            )

            temp_state = MainState(request=state.request)
            temp_state.messages = [HumanMessage(content=prompt)]

            res_state = await agent.execute(temp_state)

            if res_state.messages and res_state.messages[-1].type == "ai":
                state.content_structure = res_state.messages[-1].content
            else:
                state.content_structure = "[Structure analysis failed]"
        except Exception as e:
            log.error(f"Structure analysis failed: {e}")
            state.content_structure = f"[Structure analysis error: {e}]"

        return state

    async def generate_mermaid_node(state: KBMindMapState) -> KBMindMapState:
        """
        Generate Mermaid mindmap syntax using LLM
        """
        if not state.content_structure or state.content_structure.startswith("["):
            state.mermaid_code = "mindmap\n  root((Error))\n    No content structure available"
            return state

        # Mermaid generation prompt
        style = state.request.mindmap_style
        prompt = f"""你是一位专业的Mermaid图表生成专家。请根据以下知识结构，生成Mermaid mindmap语法。

知识结构：
{state.content_structure}

要求：
1. 使用Mermaid mindmap语法
2. 风格：{style}
3. 保持层级关系清晰
4. 节点名称简洁明了
5. 只输出Mermaid代码，不要有其他解释

Mermaid mindmap语法示例：
```
mindmap
  root((中心主题))
    主题1
      子主题1.1
      子主题1.2
    主题2
      子主题2.1
```

请生成Mermaid mindmap代码："""

        try:
            agent = create_simple_agent(
                name="mermaid_generator",
                model_name=state.request.model,
                chat_api_url=state.request.chat_api_url,
                temperature=0.5
            )

            temp_state = MainState(request=state.request)
            temp_state.messages = [HumanMessage(content=prompt)]

            res_state = await agent.execute(temp_state)

            if res_state.messages and res_state.messages[-1].type == "ai":
                mermaid_raw = res_state.messages[-1].content
                # Extract mermaid code from markdown code blocks if present
                if "```" in mermaid_raw:
                    lines = mermaid_raw.split("\n")
                    in_code_block = False
                    code_lines = []
                    for line in lines:
                        if line.strip().startswith("```"):
                            in_code_block = not in_code_block
                            continue
                        if in_code_block:
                            code_lines.append(line)
                    state.mermaid_code = "\n".join(code_lines)
                else:
                    state.mermaid_code = mermaid_raw
            else:
                state.mermaid_code = "mindmap\n  root((Error))\n    Generation failed"
        except Exception as e:
            log.error(f"Mermaid generation failed: {e}")
            state.mermaid_code = f"mindmap\n  root((Error))\n    {str(e)}"

        # Save mermaid code to file
        try:
            mermaid_path = Path(state.result_path) / "mindmap.mmd"
            mermaid_path.write_text(state.mermaid_code, encoding="utf-8")
            log.info(f"Mermaid code saved to: {mermaid_path}")
        except Exception as e:
            log.error(f"Failed to save mermaid code: {e}")

        return state

    nodes = {
        "_start_": _start_,
        "parse_files": parse_files_node,
        "analyze_structure": analyze_structure_node,
        "generate_mermaid": generate_mermaid_node,
        "_end_": lambda s: s
    }

    edges = [
        ("_start_", "parse_files"),
        ("parse_files", "analyze_structure"),
        ("analyze_structure", "generate_mermaid"),
        ("generate_mermaid", "_end_")
    ]

    builder.add_nodes(nodes).add_edges(edges)
    return builder
