from __future__ import annotations
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any

import fitz  # PyMuPDF
from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger
from dataflow_agent.state import KBPodcastState, MainState
from dataflow_agent.agentroles import create_simple_agent
from dataflow_agent.utils import get_project_root
from dataflow_agent.toolkits.multimodaltool.req_tts import generate_speech_and_save_async
from langchain_core.messages import HumanMessage

log = get_logger(__name__)

# Try importing office libraries
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

@register("kb_podcast")
def create_kb_podcast_graph() -> GenericGraphBuilder:
    """
    Workflow for Knowledge Base Podcast Generation
    Steps:
    1. Parse uploaded files (PDF/Office)
    2. Generate podcast script using LLM
    3. Generate audio using TTS
    """
    builder = GenericGraphBuilder(state_model=KBPodcastState, entry_point="_start_")

    def _start_(state: KBPodcastState) -> KBPodcastState:
        # Ensure request fields
        if not state.request.files:
            state.request.files = []

        # Initialize output directory
        if not state.result_path:
            project_root = get_project_root()
            import time
            ts = int(time.time())
            email = getattr(state.request, 'email', 'default')
            output_dir = project_root / "outputs" / "kb_outputs" / email / f"{ts}_podcast"
            output_dir.mkdir(parents=True, exist_ok=True)
            state.result_path = str(output_dir)

        state.file_contents = []
        state.podcast_script = ""
        state.audio_path = ""
        return state

    async def parse_files_node(state: KBPodcastState) -> KBPodcastState:
        """
        Parse all files and extract content
        """
        files = state.request.files
        if not files:
            state.file_contents = []
            return state

        async def process_file(file_path: str) -> Dict[str, Any]:
            file_path_obj = Path(file_path)
            filename = file_path_obj.name

            if not file_path_obj.exists():
                return {
                    "filename": filename,
                    "content": f"[Error: File not found {file_path}]"
                }

            suffix = file_path_obj.suffix.lower()
            raw_content = ""

            try:
                # PDF
                if suffix == ".pdf":
                    try:
                        doc = fitz.open(file_path)
                        text = ""
                        for page in doc:
                            text += page.get_text() + "\n"
                        raw_content = text
                    except Exception as e:
                        raw_content = f"[Error parsing PDF: {e}]"

                # Word
                elif suffix in [".docx", ".doc"]:
                    if Document is None:
                         raw_content = "[Error: python-docx not installed]"
                    else:
                        try:
                            doc = Document(file_path)
                            raw_content = "\n".join([p.text for p in doc.paragraphs])
                        except Exception as e:
                             raw_content = f"[Error parsing Docx: {e}]"

                # PPT
                elif suffix in [".pptx", ".ppt"]:
                    if Presentation is None:
                        raw_content = "[Error: python-pptx not installed]"
                    else:
                        try:
                            prs = Presentation(file_path)
                            text = ""
                            for i, slide in enumerate(prs.slides):
                                text += f"--- Slide {i+1} ---\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text += shape.text + "\n"
                            raw_content = text
                        except Exception as e:
                            raw_content = f"[Error parsing PPT: {e}]"

                else:
                    try:
                        with open(file_path, "r", encoding="utf-8") as f:
                            raw_content = f.read()
                    except:
                        raw_content = "[Unsupported file type]"

            except Exception as e:
                 raw_content = f"[Parse Error: {e}]"

            # Truncate content
            truncated_content = raw_content[:50000] if len(raw_content) > 50000 else raw_content

            return {
                "filename": filename,
                "content": truncated_content
            }

        # Run in parallel
        tasks = [process_file(f) for f in files]
        results = await asyncio.gather(*tasks)

        state.file_contents = results
        return state

    async def generate_script_node(state: KBPodcastState) -> KBPodcastState:
        """
        Generate podcast script using LLM
        """
        if not state.file_contents:
            state.podcast_script = "No content available for podcast generation."
            return state

        # Format file contents
        contents_str = ""
        for item in state.file_contents:
            contents_str += f"=== {item['filename']} ===\n{item['content']}\n\n"

        # Podcast script prompt
        language = state.request.language
        prompt = f"""你是一位专业的知识播客主播。基于以下资料，生成一段5-10分钟的知识播客脚本。

要求：
1. 口语化、生动有趣，避免书面语
2. 结构清晰：开场白 → 核心内容 → 总结
3. 使用类比和例子帮助理解
4. 适当加入互动性语言（"你可能会想..."）
5. 使用{language}语言

资料内容：
{contents_str}

请生成播客脚本："""

        try:
            agent = create_simple_agent(
                name="podcast_script_generator",
                model_name=state.request.model,
                chat_api_url=state.request.chat_api_url,
                temperature=0.7
            )

            temp_state = MainState(request=state.request)
            temp_state.messages = [HumanMessage(content=prompt)]

            res_state = await agent.execute(temp_state)

            if res_state.messages and res_state.messages[-1].type == "ai":
                state.podcast_script = res_state.messages[-1].content
            else:
                state.podcast_script = "[Script generation failed]"
        except Exception as e:
            log.error(f"Script generation failed: {e}")
            state.podcast_script = f"[Script generation error: {e}]"

        # Save script to file
        try:
            script_path = Path(state.result_path) / "script.txt"
            script_path.write_text(state.podcast_script, encoding="utf-8")
        except Exception as e:
            log.error(f"Failed to save script: {e}")

        return state

    async def generate_audio_node(state: KBPodcastState) -> KBPodcastState:
        """
        Generate audio using TTS
        """
        if not state.podcast_script or state.podcast_script.startswith("["):
            state.audio_path = ""
            return state

        try:
            audio_path = str(Path(state.result_path) / "podcast.wav")

            result_path = await generate_speech_and_save_async(
                text=state.podcast_script,
                save_path=audio_path,
                api_url=state.request.chat_api_url,
                api_key=state.request.api_key,
                model=state.request.tts_model,
                voice_name=state.request.voice_name
            )

            state.audio_path = result_path
            log.info(f"Audio generated successfully: {result_path}")
        except Exception as e:
            log.error(f"Audio generation failed: {e}")
            state.audio_path = f"[Audio generation error: {e}]"

        return state

    nodes = {
        "_start_": _start_,
        "parse_files": parse_files_node,
        "generate_script": generate_script_node,
        "generate_audio": generate_audio_node,
        "_end_": lambda s: s
    }

    edges = [
        ("_start_", "parse_files"),
        ("parse_files", "generate_script"),
        ("generate_script", "generate_audio"),
        ("generate_audio", "_end_")
    ]

    builder.add_nodes(nodes).add_edges(edges)
    return builder
