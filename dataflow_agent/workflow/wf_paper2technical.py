"""
paper2technical workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
生成时间: 2025-12-07 23:36:51

1. 在 **TOOLS** 区域定义需要暴露给 Prompt 的前置工具
2. 在 **NODES**  区域实现异步节点函数 (await-able)
3. 在 **EDGES**  区域声明有向边
4. 最后返回 builder.compile() 或 GenericGraphBuilder
"""

from __future__ import annotations
import json
import time
from pathlib import Path
import re

from dataflow_agent.state import Paper2FigureState
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.workflow.registry import register
from dataflow_agent.agentroles import create_simple_agent
from dataflow_agent.toolkits.tool_manager import get_tool_manager
from dataflow_agent.toolkits.multimodaltool.bg_tool import (
    local_tool_for_svg_render,
    local_tool_for_raster_to_svg,
)
from dataflow_agent.toolkits.multimodaltool.mineru_tool import svg_to_emf
from dataflow_agent.utils import get_project_root
from dataflow_agent.logger import get_logger
log = get_logger(__name__)


def _ensure_result_path(state: Paper2FigureState) -> str:
    """
    统一本次 workflow 的根输出目录：
    - 如果 state.result_path 已存在（通常由调用方传入，形如 时间戳+编码），直接使用；
    - 否则：使用 get_project_root() / "outputs" / "paper2tec" / <timestamp>，
      并回写到 state.result_path，确保后续节点共享同一目录，避免数据串台。
    """
    raw = getattr(state, "result_path", None)
    if raw:
        return raw

    root = get_project_root()
    ts = int(time.time())
    base_dir = (root / "outputs" / "paper2tec" / str(ts)).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)
    state.result_path = str(base_dir)
    return state.result_path


def _extract_svg_from_react_md(md_path: Path) -> str:
    """
    从 React 组件的 .md 文件中提取纯 SVG 代码。

    输入的 .md 文件包含 React 组件代码,其中嵌入了 SVG。
    此函数提取 <svg>...</svg> 部分,并将 React 语法转换为纯 SVG。
    """
    if not md_path.exists():
        log.warning(f"模板文件不存在: {md_path}")
        return ""

    try:
        content = md_path.read_text(encoding="utf-8")

        # 查找 <svg 开始标签
        svg_start = content.find("<svg")
        if svg_start == -1:
            log.warning(f"未在文件中找到 <svg 标签: {md_path}")
            return ""

        # 查找对应的 </svg> 结束标签
        svg_end = content.rfind("</svg>")
        if svg_end == -1:
            log.warning(f"未在文件中找到 </svg> 标签: {md_path}")
            return ""

        svg_end += len("</svg>")
        svg_code = content[svg_start:svg_end]

        # 清理 React 特有语法:
        # 1. 将 className 替换为 class
        svg_code = svg_code.replace('className="', 'class="')

        # 2. 将 JSX 驼峰命名属性转换为 SVG 连字符命名
        # 注意：某些属性在 SVG 中必须保持驼峰命名（如 markerWidth, markerHeight, viewBox 等）
        jsx_to_svg_attrs = {
            'strokeWidth': 'stroke-width',
            'strokeDasharray': 'stroke-dasharray',
            'strokeLinecap': 'stroke-linecap',
            'strokeLinejoin': 'stroke-linejoin',
            'strokeOpacity': 'stroke-opacity',
            'fillOpacity': 'fill-opacity',
            'textAnchor': 'text-anchor',
            'fontWeight': 'font-weight',
            'fontSize': 'font-size',
            'fontFamily': 'font-family',
            # markerWidth 和 markerHeight 应该保持驼峰命名，不转换
            'markerEnd': 'marker-end',
            'markerStart': 'marker-start',
            'markerMid': 'marker-mid',
            'clipPath': 'clip-path',
        }
        for jsx_attr, svg_attr in jsx_to_svg_attrs.items():
            svg_code = svg_code.replace(f'{jsx_attr}=', f'{svg_attr}=')

        # 3. 将 {colors.xxx} 这样的变量引用替换为实际颜色值
        colors_match = re.search(r'const colors = \{([^}]+)\}', content, re.DOTALL)
        if colors_match:
            colors_def = colors_match.group(1)
            # 解析颜色定义
            color_map = {}
            for line in colors_def.split('\n'):
                match = re.search(r'(\w+):\s*"([^"]+)"', line)
                if match:
                    color_map[match.group(1)] = match.group(2)

            # 替换 {colors.xxx} 为实际颜色值
            for key, value in color_map.items():
                svg_code = svg_code.replace(f'{{colors.{key}}}', value)

        # 4. 移除 React 注释 {/* ... */}
        svg_code = re.sub(r'\{/\*.*?\*/\}', '', svg_code, flags=re.DOTALL)

        # 5. 转义 XML 特殊字符（在文本内容中）
        # 注意：只转义 text 元素内的 &，不转义已经是实体引用的部分
        # 使用负向前瞻确保不会重复转义已经转义的内容
        svg_code = re.sub(r'&(?!amp;|lt;|gt;|quot;|apos;|#)', '&amp;', svg_code)

        return svg_code.strip()

    except Exception as e:
        log.error(f"提取 SVG 代码失败: {e}")
        return ""


def _get_template_svg_code(state: Paper2FigureState) -> str:
    """
    根据语言选择合适的 SVG 模板代码。

    - 中文: 使用 SVG_template_ZN_gray.md
    - 英文: 使用 SVG_template_EN_gray.md

    返回纯 SVG 代码字符串。
    """
    root = get_project_root()
    lang = getattr(getattr(state, "request", None), "language", "EN")

    # 根据语言选择模板文件
    if lang.upper() in ["ZH", "CN", "CHINESE", "中文"]:
        template_file = root / "SVG_template_ZN_gray.md"
    else:
        template_file = root / "SVG_template_EN_gray.md"

    svg_code = _extract_svg_from_react_md(template_file)

    if not svg_code:
        log.warning(f"无法从模板文件提取 SVG 代码,使用空字符串")

    return svg_code


def _get_palette_config(state: Paper2FigureState) -> dict | None:
    """
    根据 request.tech_route_palette 返回色卡配置；未选择则返回 None。
    """
    palette_name = getattr(getattr(state, "request", None), "tech_route_palette", "") or ""
    if not palette_name:
        return None

    palettes = {
        "academic_blue": {
            "name": "academic_blue",
            "colors": ["#1F6FEB", "#60A5FA", "#A7C7FF", "#0B3D91"],
            "level_colors": ["#A7C7FF", "#60A5FA", "#1F6FEB", "#0B3D91"],
            "arrow_color": "#0B3D91",
            "text_color": "#0B3D91",
        },
        "teal_orange": {
            "name": "teal_orange",
            "colors": ["#0F766E", "#14B8A6", "#F59E0B", "#FB923C"],
            "level_colors": ["#14B8A6", "#0F766E", "#F59E0B", "#FB923C"],
            "arrow_color": "#0F766E",
            "text_color": "#0F766E",
        },
        "slate_rose": {
            "name": "slate_rose",
            "colors": ["#334155", "#64748B", "#F43F5E", "#FCA5A5"],
            "level_colors": ["#64748B", "#334155", "#FCA5A5", "#F43F5E"],
            "arrow_color": "#334155",
            "text_color": "#334155",
        },
        "indigo_amber": {
            "name": "indigo_amber",
            "colors": ["#4338CA", "#6366F1", "#F59E0B", "#FCD34D"],
            "level_colors": ["#6366F1", "#4338CA", "#FCD34D", "#F59E0B"],
            "arrow_color": "#4338CA",
            "text_color": "#4338CA",
        },
    }

    return palettes.get(palette_name)


@register("paper2technical")
def create_paper2technical_graph() -> GenericGraphBuilder:  # noqa: N802
    """
    Workflow factory: dfa run --wf paper2technical
    """
    # 使用 Paper2FigureState，复用其中的 paper_file / paper_idea / fig_desc 等字段，
    # 这里不做图像生成和抠图，只负责“技术路线图”的 SVG + PPT 逻辑。
    builder = GenericGraphBuilder(
        state_model=Paper2FigureState,
        entry_point="_start_",        # 入口统一为 _start_，再由路由函数分发
    )

    # ----------------------------------------------------------------------
    # TOOLS (pre_tool definitions)
    # ----------------------------------------------------------------------
    # 1) 提供给 paper_idea_extractor 的 PDF 内容（标题 + 前几页正文）
    @builder.pre_tool("paper_content", "paper_idea_extractor")
    def _get_paper_content(state: Paper2FigureState):
        """
        前置工具: 读取论文 PDF 的标题和前若干页内容，供 paper_idea_extractor 节点使用。

        - 作用: 为大模型提供足够的上下文，让其抽取论文中的技术路线/实验流程关键信息。
        - 输出: 一个字符串，包含论文标题 + 前若干页文本。
        """
        import fitz  # PyMuPDF
        import PyPDF2

        pdf_path = state.paper_file
        if not pdf_path:
            log.warning("paper_file 为空，无法读取 PDF 内容")
            return ""

        try:
            with open(pdf_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                paper_title = reader.metadata.get("/Title", "Unknown Title")
        except Exception:
            paper_title = "Unknown Title"

        try:
            doc = fitz.open(pdf_path)
        except Exception as e:
            log.error(f"打开 PDF 失败: {e}")
            return f"The title of the paper is {paper_title}"

        text_parts: list[str] = []
        # 读取前 10 页内容，通常技术路线、整体框架会在前几页出现
        for page_idx in range(min(10, len(doc))):
            page = doc.load_page(page_idx)
            text_parts.append(page.get_text("text") or "")

        content = "\n".join(text_parts).strip()
        final_text = (
            f"The title of the paper is {paper_title}\n\n"
            f"Here are the first 10 pages of the paper:\n{content}"
        )
        log.info("paper_content 提取完成")
        return final_text

    @builder.pre_tool("paper_idea", "technical_route_bw_svg_generator")
    def _get_bw_paper_idea(state: Paper2FigureState):
        return state.paper_idea or ""

    @builder.pre_tool("template_svg_code", "technical_route_bw_svg_generator")
    def _get_template_svg(state: Paper2FigureState):
        """
        前置工具: 提供 SVG 模板代码给黑白技术路线图生成器。

        - 作用: 根据语言选择合适的 SVG 模板,供 agent 参考其结构和布局。
        - 输出: 纯 SVG 代码字符串。
        """
        return _get_template_svg_code(state)

    @builder.pre_tool("validation_feedback", "technical_route_bw_svg_generator")
    def _get_bw_feedback(state: Paper2FigureState):
        return state.temp_data.get("validation_feedback", "") if hasattr(state, "temp_data") else ""

    @builder.pre_tool("validation_feedback", "technical_route_colorize_svg")
    def _get_color_feedback(state: Paper2FigureState):
        return state.temp_data.get("validation_feedback", "") if hasattr(state, "temp_data") else ""

    @builder.pre_tool("bw_svg_code", "technical_route_colorize_svg")
    def _get_bw_svg_code(state: Paper2FigureState):
        return state.figure_tec_svg_bw_content or ""

    @builder.pre_tool("palette_json", "technical_route_colorize_svg")
    def _get_palette_json(state: Paper2FigureState):
        return state.temp_data.get("palette_json", "") if hasattr(state, "temp_data") else ""

    # ----------------------------------------------------------------------

    # ==============================================================
    # NODES
    # ==============================================================
    async def paper_idea_extractor_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 1: 从 PDF 中抽取论文的核心思想 / 技术路线相关信息

        - 只在 input_type == "PDF" 时作为入口节点被调用。
        - 基于 pre_tool("paper_content") 提供的标题 + 前若干页内容，
          调用专门的 agent（例如 paper_idea_extractor）生成摘要。
        - 该摘要用于后续技术路线图描述生成。

        输入:
            state.paper_file : 论文 PDF 路径
        输出:
            state.paper_idea : 论文核心思想 / 技术路线要点摘要
            state.agent_results["paper_idea_extractor"] : agent 原始输出
        """
        agent = create_simple_agent("paper_idea_extractor")
        state = await agent.execute(state=state)
        return state

    def _svg_has_cjk(text: str) -> bool:
        """简单判断 SVG 中是否包含中文字符，用于日志和调试。"""
        return bool(re.search(r"[\u4e00-\u9fff]", text))


    def _inject_chinese_font(svg_code: str) -> str:
        """
        如果 SVG 中没有设定中文友好的 font-family，则注入一段全局样式，
        指定一组 CJK 字体作为优先字体。

        注意：字体名请根据实际安装的字体调整。
        """
        if "font-family" in svg_code:
            return svg_code

        idx = svg_code.find(">")
        if idx == -1:
            return svg_code

        style_block = """
  <style type="text/css">
    text, tspan {
      font-family: "Noto Sans CJK SC", "Microsoft YaHei", "SimHei", "SimSun", "WenQuanYi Zen Hei", sans-serif;
    }
  </style>
"""
        return svg_code[: idx + 1] + style_block + svg_code[idx + 1 :]


    def _extract_svg_fragment(svg_code: str) -> str:
        if not svg_code:
            return ""
        text = svg_code.strip()
        if text.startswith("```"):
            lines = [line for line in text.splitlines() if line.strip("`").strip()]
            for i, line in enumerate(lines):
                if "<svg" in line or "<SVG" in line or "<Svg" in line:
                    text = "\n".join(lines[i:])
                    break
        start = text.find("<svg")
        end = text.rfind("</svg>")
        if start == -1 or end == -1:
            return text
        end += len("</svg>")
        return text[start:end].strip()

    def _validate_svg_renderable(svg_code: str, output_path: str) -> tuple[bool, str]:
        import xml.etree.ElementTree as ET

        if not svg_code:
            return False, "svg_code 为空"
        fragment = _extract_svg_fragment(svg_code)
        if "<svg" not in fragment and "<SVG" not in fragment:
            return False, "未检测到 <svg> 根标签"
        if "viewBox" not in fragment and "viewbox" not in fragment:
            return False, "缺少 viewBox 属性"
        try:
            ET.fromstring(fragment)
        except Exception as e:
            return False, f"SVG XML 解析失败: {e}"

        try:
            local_tool_for_svg_render(
                {
                    "svg_code": fragment,
                    "output_path": output_path,
                }
            )
        except Exception as e:
            return False, f"SVG 渲染失败: {e}"

        return True, ""

    async def technical_route_bw_svg_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        黑白技术路线图生成（参考模板 SVG 代码）。
        不再使用 VLM 图片输入,而是将 SVG 模板代码作为文本输入。
        失败时在节点内做渲染校验重试。
        """
        from dataflow_agent.agentroles.paper2any_agents.technical_route_bw_svg_generator import (
            create_technical_route_bw_svg_generator,
        )

        base_dir = Path(_ensure_result_path(state))
        base_dir.mkdir(parents=True, exist_ok=True)

        max_retries = 3
        last_error = ""

        for attempt in range(max_retries):
            if not hasattr(state, "temp_data"):
                state.temp_data = {}
            state.temp_data["validation_feedback"] = last_error

            agent = create_technical_route_bw_svg_generator(
                model_name="gpt-5.2-medium",
                temperature=0.0,
                max_tokens=16384,  # 增加到16K以支持完整的SVG代码输出
                parser_type="json",
            )

            state = await agent.execute(state=state)
            svg_code = getattr(state, "figure_tec_svg_bw_content", None)
            if not svg_code:
                last_error = "缺少 svg_code"
                continue

            svg_code = _inject_chinese_font(svg_code)
            validate_path = str((base_dir / f"technical_route_bw_validate_{attempt}.png").resolve())
            ok, err = _validate_svg_renderable(svg_code, validate_path)
            if ok:
                # 输出落盘
                timestamp = int(time.time())
                svg_output_path = str((base_dir / f"technical_route_bw_{timestamp}.svg").resolve())
                png_output_path = str((base_dir / f"technical_route_bw_{timestamp}.png").resolve())
                try:
                    Path(svg_output_path).write_text(svg_code, encoding="utf-8")
                    png_path = local_tool_for_svg_render(
                        {
                            "svg_code": svg_code,
                            "output_path": png_output_path,
                        }
                    )
                    state.svg_bw_file_path = svg_output_path
                    state.svg_bw_img_path = png_path
                    state.svg_file_path = svg_output_path
                    state.svg_img_path = png_path
                    state.figure_tec_svg_content = svg_code
                    log.critical(f"[state.svg_bw_img_path]: {state.svg_bw_img_path}")
                    log.critical(f"[state.svg_bw_file_path]: {state.svg_bw_file_path}")
                    return state
                except Exception as e:
                    last_error = f"SVG 落盘/渲染失败: {e}"
                    continue

            last_error = err

        log.error(f"technical_route_bw_svg_generator_node: 重试失败: {last_error}")
        return state

    async def technical_route_colorize_svg_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        彩色技术路线图生成：仅基于黑白 SVG + 色卡。
        """
        from dataflow_agent.agentroles.paper2any_agents.technical_route_colorize_svg_agent import (
            create_technical_route_colorize_svg_agent,
        )

        base_dir = Path(_ensure_result_path(state))
        base_dir.mkdir(parents=True, exist_ok=True)

        palette_cfg = _get_palette_config(state)
        if not palette_cfg:
            return state

        if not hasattr(state, "temp_data"):
            state.temp_data = {}
        state.temp_data["palette_json"] = json.dumps(palette_cfg, ensure_ascii=False)

        max_retries = 3
        last_error = ""
        for attempt in range(max_retries):
            state.temp_data["validation_feedback"] = last_error

            agent = create_technical_route_colorize_svg_agent(
                model_name="gpt-5.2-medium",
                temperature=0.0,
                max_tokens=16384,  # 增加到16K以支持完整的SVG代码输出
                parser_type="json",
            )
            state = await agent.execute(state=state)

            svg_code = getattr(state, "figure_tec_svg_color_content", None)
            if not svg_code:
                last_error = "缺少 svg_code"
                continue

            svg_code = _inject_chinese_font(svg_code)
            validate_path = str((base_dir / f"technical_route_color_validate_{attempt}.png").resolve())
            ok, err = _validate_svg_renderable(svg_code, validate_path)
            if ok:
                timestamp = int(time.time())
                svg_output_path = str((base_dir / f"technical_route_color_{timestamp}.svg").resolve())
                png_output_path = str((base_dir / f"technical_route_color_{timestamp}.png").resolve())
                try:
                    Path(svg_output_path).write_text(svg_code, encoding="utf-8")
                    png_path = local_tool_for_svg_render(
                        {
                            "svg_code": svg_code,
                            "output_path": png_output_path,
                        }
                    )
                    state.svg_color_file_path = svg_output_path
                    state.svg_color_img_path = png_path
                    log.critical(f"[state.svg_color_img_path]: {state.svg_color_img_path}")
                    log.critical(f"[state.svg_color_file_path]: {state.svg_color_file_path}")
                    return state
                except Exception as e:
                    last_error = f"SVG 落盘/渲染失败: {e}"
                    continue

            last_error = err

        log.error(f"technical_route_colorize_svg_node: 重试失败: {last_error}")
        return state

    async def technical_ppt_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 4: 基于技术路线图 PNG 生成 PPT

        - 根据前面步骤生成的 PNG 整图（state.svg_img_path），
          生成用于展示技术路线图的 PPT。
        - 仅使用 SVG -> PNG 路径插入位图，不再尝试 EMF。
        """
        from pptx import Presentation
        from PIL import Image

        # ✅ 临时提高 PIL 图像大小限制，防止 decompression bomb 错误
        original_max_pixels = Image.MAX_IMAGE_PIXELS
        Image.MAX_IMAGE_PIXELS = None  # 或设置为更大的值，如 500_000_000

        try:
            # 输出目录：统一使用本次 workflow 的根输出目录
            run_root = Path(_ensure_result_path(state))
            output_dir = run_root
            output_dir.mkdir(parents=True, exist_ok=True)

            timestamp = int(time.time())
            ppt_path = output_dir / f"technical_route_{timestamp}.pptx"

            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]

            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # ------------------------------------------------------------------
            # 第 1 页：技术路线图（仅使用 PNG）
            # ------------------------------------------------------------------
            slide = prs.slides.add_slide(blank_slide_layout)

            png_path = getattr(state, "svg_img_path", None)
            color_png_path = getattr(state, "svg_color_img_path", None)
            palette_selected = bool(getattr(getattr(state, "request", None), "tech_route_palette", ""))

            def _insert_picture(pic_path: str) -> bool:
                """通用插图函数：按 80% 宽度缩放并居中。"""
                try:
                    pic = slide.shapes.add_picture(pic_path, 0, 0)
                except Exception as e:
                    log.error(f"technical_ppt_generator_node: 插入图片失败: {e}")
                    return False

                if pic.width and pic.width > 0:
                    scale = (slide_width * 0.8) / pic.width
                else:
                    scale = 1.0

                pic.width = int(pic.width * scale)
                pic.height = int(pic.height * scale)

                pic.left = int((slide_width - pic.width) / 2)
                pic.top = int((slide_height - pic.height) / 2)
                return True

            # 直接使用 PNG（位图）插入 PPT
            inserted = False

            final_png = color_png_path if palette_selected and color_png_path else png_path
            if final_png:
                try:
                    ok = _insert_picture(final_png)
                    if ok:
                        log.info(
                            "technical_ppt_generator_node: 使用 PNG 插入技术路线图成功: %s",
                            final_png,
                        )
                    else:
                        log.error(
                            "technical_ppt_generator_node: PNG 插入失败，第一页可能为空白"
                        )
                except Exception as e:
                    log.error(
                        "technical_ppt_generator_node: PNG 插入失败，第一页将为空白: %s",
                        e,
                    )

            if (not inserted) and (not final_png):
                log.warning(
                    "technical_ppt_generator_node: svg_file_path / svg_img_path 均为空，"
                    "第一页将为空白"
                )

            # ------------------------------------------------------------------
            # 第 2 页：操作提示页（写上“右键转换成形状”）
            # ------------------------------------------------------------------
            slide2 = prs.slides.add_slide(blank_slide_layout)
            left = int(slide_width * 0.1)
            top = int(slide_height * 0.3)
            width = int(slide_width * 0.8)
            height = int(slide_height * 0.4)

            textbox = slide2.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = "右键转换成形状"

            prs.save(str(ppt_path))
            state.ppt_path = str(ppt_path)
            log.info(f"technical_ppt_generator_node: PPT 已生成: {ppt_path}")
        finally:
            # ✅ 恢复原始限制
            Image.MAX_IMAGE_PIXELS = original_max_pixels

        return state

    # ==============================================================
    # 注册 nodes / edges
    # ==============================================================

    def set_entry_node(state: Paper2FigureState) -> str:
        """
        路由函数: 根据输入类型选择技术路线工作流的入口节点。

        - input_type == "PDF"  : 从 PDF 中抽取论文想法，先走 paper_idea_extractor
        - input_type == "TEXT" : 直接使用调用方提供的文本描述，跳过 PDF 抽取，
                                 从 technical_route_bw_svg_generator 开始
        其他值:
        - 认为是不合法输入，直接结束工作流。
        """
        input_type = getattr(state.request, "input_type", "PDF")
        if input_type == "PDF":
            log.critical("paper2technical: 进入 PDF 流程 (paper_idea_extractor)")
            return "paper_idea_extractor"
        elif input_type == "TEXT":
            log.critical("paper2technical: 进入 TEXT 流程 (technical_route_bw_svg_generator)")
            return "technical_route_bw_svg_generator"
        else:
            log.error(f"paper2technical: Invalid input type: {input_type}")
            return "_end_"

    def _init_result_path(state: Paper2FigureState) -> Paper2FigureState:
        """
        _start_ 节点：确保本次 workflow 有一个统一的 result_path 根目录。
        - 若用户已在 state.result_path 传入自定义目录，则直接使用该目录；
        - 若未传入，则初始化为 get_project_root()/outputs/paper2tec/<timestamp>。
        """
        _ensure_result_path(state)
        return state

    nodes = {
        "_start_": _init_result_path,
        "paper_idea_extractor": paper_idea_extractor_node,
        "technical_route_bw_svg_generator": technical_route_bw_svg_generator_node,
        "technical_route_colorize_svg": technical_route_colorize_svg_node,
        "technical_ppt_generator": technical_ppt_generator_node,
        "_end_": lambda state: state,  # 终止节点
    }

    # ------------------------------------------------------------------
    # EDGES  (从节点 A 指向节点 B)
    # ------------------------------------------------------------------
    edges = [
        # PDF 流程: 先抽想法，再生成黑白技术路线
        ("paper_idea_extractor", "technical_route_bw_svg_generator"),
        # 生成黑白后，如有配色则上色，否则直接进 PPT
        ("technical_route_colorize_svg", "technical_ppt_generator"),
        ("technical_ppt_generator", "_end_"),
    ]

    def _route_after_bw(state: Paper2FigureState) -> str:
        palette = getattr(getattr(state, "request", None), "tech_route_palette", "") or ""
        if palette:
            return "technical_route_colorize_svg"
        return "technical_ppt_generator"

    builder.add_nodes(nodes).add_edges(edges).add_conditional_edge("_start_", set_entry_node)
    builder.add_conditional_edge("technical_route_bw_svg_generator", _route_after_bw)
    return builder
