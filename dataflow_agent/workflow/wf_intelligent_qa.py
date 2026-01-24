from __future__ import annotations
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any, Optional

import fitz  # PyMuPDF
from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger
from dataflow_agent.state import IntelligentQAState, MainState
from dataflow_agent.agentroles import create_vlm_agent, create_simple_agent
from dataflow_agent.utils import get_project_root
from dataflow_agent.promptstemplates.resources.pt_qa_agent_repo import QaAgent as QaAgentPrompts
from langchain_core.messages import HumanMessage

log = get_logger(__name__)

# Try importing office libraries
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

@register("intelligent_qa")
def create_intelligent_qa_graph() -> GenericGraphBuilder:
    """
    Workflow for Intelligent Q&A with Parallel Processing
    Steps:
    1. Parse uploaded files (PDF/Office/Image/Video)
    2. Parallel Analysis: Call LLM/VLM for each file individually
    3. Aggregate context and Final QA
    """
    builder = GenericGraphBuilder(state_model=IntelligentQAState, entry_point="_start_")

    def _start_(state: IntelligentQAState) -> IntelligentQAState:
        # Ensure request fields
        if not state.request.files:
            state.request.files = []
        if not state.request.query:
            state.request.query = ""
        # Initialize file analyses
        state.file_analyses = []
        return state

    async def parallel_parse_node(state: IntelligentQAState) -> IntelligentQAState:
        """
        Parallel parsing AND analysis of all files
        """
        files = state.request.files
        if not files:
            state.context_content = ""
            return state

        async def process_file(file_path: str) -> Dict[str, Any]:
            file_path_obj = Path(file_path)
            filename = file_path_obj.name
            
            if not file_path_obj.exists():
                return {
                    "filename": filename,
                    "analysis": f"[Error: File not found {file_path}]",
                    "content": ""
                }
            
            suffix = file_path_obj.suffix.lower()
            raw_content = ""
            analysis_result = ""
            file_type = "unknown"
            
            try:
                # ==========================
                # 1. Extraction Phase
                # ==========================
                
                # PDF
                if suffix == ".pdf":
                    file_type = "document"
                    try:
                        doc = fitz.open(file_path)
                        text = ""
                        for page in doc:
                            text += page.get_text() + "\n"
                        raw_content = text
                    except Exception as e:
                        raw_content = f"[Error parsing PDF: {e}]"

                # Word
                elif suffix in [".docx", ".doc"]:
                    file_type = "document"
                    if Document is None:
                         raw_content = "[Error: python-docx not installed]"
                    else:
                        try:
                            doc = Document(file_path)
                            raw_content = "\n".join([p.text for p in doc.paragraphs])
                        except Exception as e:
                             raw_content = f"[Error parsing Docx: {e}]"

                # PPT (Simulate PPT -> Text)
                elif suffix in [".pptx", ".ppt"]:
                    file_type = "presentation"
                    if Presentation is None:
                        raw_content = "[Error: python-pptx not installed]"
                    else:
                        try:
                            prs = Presentation(file_path)
                            text = ""
                            for i, slide in enumerate(prs.slides):
                                text += f"--- Slide {i+1} ---\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text += shape.text + "\n"
                            raw_content = text
                        except Exception as e:
                            raw_content = f"[Error parsing PPT: {e}]"
                
                # Image / Video
                elif suffix in [".jpg", ".jpeg", ".png", ".mp4", ".mov", ".avi"]:
                    file_type = "media"
                    # For media, raw_content comes from VLM understanding
                    # We will do analysis directly here using VLM Agent
                    raw_content = "[Media file - will be analyzed by VLM]"
                
                else:
                    file_type = "text"
                    try:
                        with open(file_path, "r", encoding="utf-8") as f:
                            raw_content = f.read()
                    except:
                        raw_content = "[Unsupported file type]"

                # ==========================
                # 2. Analysis Phase (Parallel LLM Call)
                # ==========================

                if file_type == "media":
                    # Use VLM Agent for Media
                    vlm_mode = "understanding"
                    input_key = "input_image"
                    if suffix in [".mp4", ".mov", ".avi"]:
                        vlm_mode = "video_understanding"
                        input_key = "input_video"

                    try:
                        # Create VLM Agent
                        agent = create_vlm_agent(
                            name=f"vlm_analyzer_{filename}",
                            vlm_mode=vlm_mode,
                            model_name="gemini-2.5-flash",
                            chat_api_url=state.request.chat_api_url,
                            additional_params={input_key: file_path}
                        )

                        # Construct Prompt for VLM
                        # We ask it to describe AND relate to user query
                        vlm_prompt = f"Please analyze this media content. User Question: {state.request.query}. Describe relevant details."

                        temp_state = MainState(request=state.request)
                        temp_state.messages = [HumanMessage(content=vlm_prompt)]

                        res_state = await agent.execute(temp_state)

                        if res_state.messages and res_state.messages[-1].type == "ai":
                            analysis_result = res_state.messages[-1].content
                            raw_content = "[Media Content Processed by VLM]"
                        else:
                            analysis_result = "[VLM returned no content]"
                    except Exception as e:
                        log.error(f"VLM analysis failed for {filename}: {e}")
                        analysis_result = f"[VLM Analysis Error: {e}]"

                else:
                    # Use Simple Agent for Text
                    # Only if content is not empty or error
                    if raw_content and not raw_content.startswith("[Error"):
                        try:
                            # Prepare Prompt
                            # Limit raw_content size to avoid context overflow if huge (simple truncation)
                            truncated_content = raw_content[:50000]  # 50k char limit rough guard

                            analysis_prompt = QaAgentPrompts.file_analysis_prompt.format(
                                filename=filename,
                                file_type=file_type,
                                content=truncated_content,
                                query=state.request.query
                            )

                            agent = create_simple_agent(
                                name=f"text_analyzer_{filename}",
                                model_name=state.request.model,
                                chat_api_url=state.request.chat_api_url,
                                temperature=0.3
                            )

                            temp_state = MainState(request=state.request)
                            temp_state.messages = [HumanMessage(content=analysis_prompt)]

                            res_state = await agent.execute(temp_state)

                            if res_state.messages and res_state.messages[-1].type == "ai":
                                analysis_result = res_state.messages[-1].content
                            else:
                                analysis_result = "[LLM Analysis Failed]"
                        except Exception as e:
                            log.error(f"Text analysis failed for {filename}: {e}")
                            analysis_result = f"[Text Analysis Error: {e}]"
                    else:
                        analysis_result = raw_content  # Pass through error or empty

            except Exception as e:
                 analysis_result = f"[Analysis Error: {e}]"
            
            return {
                "filename": filename,
                "analysis": analysis_result,
                "content": raw_content[:1000] + "..." if len(raw_content) > 1000 else raw_content # Store brief raw content for debug
            }

        # Run in parallel
        tasks = [process_file(f) for f in files]
        results = await asyncio.gather(*tasks)
        
        state.file_analyses = results
        return state

    async def chat_node(state: IntelligentQAState) -> IntelligentQAState:
        """
        Final synthesis using aggregated analyses
        """
        # Construct history string
        history_str = ""
        for msg in state.request.history:
            role = msg.get("role", "user")
            content = msg.get("content", "")
            history_str += f"{role}: {content}\n"
        
        # Format analyses for the prompt
        analyses_str = ""
        for item in state.file_analyses:
            analyses_str += f"--- Analysis of {item['filename']} ---\n{item['analysis']}\n\n"
            
        final_prompt = QaAgentPrompts.final_qa_prompt.format(
            query=state.request.query,
            file_analyses=analyses_str,
            history=history_str
        )
        
        # Use QA Agent for Final Answer
        agent = create_qa_agent(
            model_name=state.request.model,
            chat_api_url=state.request.chat_api_url,
            temperature=0.7
        )
        
        from langchain_core.messages import HumanMessage
        # Override messages
        state.messages = [HumanMessage(content=final_prompt)]
        
        new_state = await agent.execute(state)
        
        if new_state.messages and new_state.messages[-1].type == "ai":
            state.answer = new_state.messages[-1].content
        else:
            state.answer = "Sorry, I couldn't generate an answer."
            
        return state

    nodes = {
        "_start_": _start_,
        "parallel_parse": parallel_parse_node,
        "chat": chat_node,
        "_end_": lambda s: s
    }

    edges = [
        ("_start_", "parallel_parse"),
        ("parallel_parse", "chat"),
        ("chat", "_end_")
    ]

    builder.add_nodes(nodes).add_edges(edges)
    return builder
