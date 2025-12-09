"""
icongen workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
生成时间: 2025-10-27 11:11:56

1. 在 **TOOLS** 区域定义需要暴露给 Prompt 的前置工具
2. 在 **NODES**  区域实现异步节点函数 (await-able)
3. 在 **EDGES**  区域声明有向边
4. 最后返回 builder.compile() 或 GenericGraphBuilder
"""

from __future__ import annotations
import asyncio
import json
import os
from dataflow_agent.state import MainState, Paper2FigureState
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder


from dataflow_agent.workflow.registry import register
# from dataflow_agent.agentroles import get_agent_cls, create_agent

from dataflow_agent.toolkits.tool_manager import get_tool_manager
from langchain.tools import tool
from langgraph.graph import StateGraph
from langgraph.prebuilt import ToolNode, tools_condition

from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger

from dataflow_agent.toolkits.imtool.req_img import generate_or_edit_and_save_image_async
from dataflow_agent.toolkits.imtool.bg_tool import local_tool_for_bg_remove, local_tool_for_raster_to_svg
from dataflow_agent.toolkits.imtool.sam_tool import segment_layout_boxes
from dataflow_agent.toolkits.imtool.mineru_tool import (
    svg_to_emf,
    recursive_mineru_layout,
)
from dataflow_agent.agentroles import create_graph_agent

import re, pdfplumber, PyPDF2, time, shutil, fitz
import numpy as np
from PIL import Image

from dataflow_agent.utils import (
    build_output_directory,
    add_image_element,
    add_text_element,
    setup_presentation_size,
    get_project_root,
)

from pathlib import Path
import time, random
from pptx import Presentation
from pptx.dml.color import RGBColor 


log = get_logger(__name__)

def _ensure_result_path(state: Paper2FigureState) -> str:
    """
    统一本次 paper2figure_with_sam workflow 的根输出目录：
    - 如果 state.result_path 已存在（通常由调用方传入），直接使用；
    - 否则：使用 get_project_root() / "outputs" / "paper2figure" / <timestamp>，
      并写回 state.result_path，后续节点共享同一目录。
    """
    raw = getattr(state, "result_path", None)
    if raw:
        return raw

    root = get_project_root()
    ts = int(time.time())
    base_dir = (root / "outputs" / "paper2figure" / str(ts)).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)
    state.result_path = str(base_dir)
    return state.result_path

def _ts_name(stem: str, ext: str = ".png") -> str:
    timestamp = int(time.time())  # 获取当前时间戳（秒）
    return f"./{stem}{timestamp}{ext}"

@register("paper2fig_with_sam")
def create_p2fig_graph() -> GenericGraphBuilder:  # noqa: N802
    """
    Workflow factory: dfa run --wf paper2fig
    """
    builder = GenericGraphBuilder(state_model=Paper2FigureState,
                                  entry_point="_start_")  # 自行修改入口

    # ----------------------------------------------------------------------
    # TOOLS (pre_tool definitions)
    # ----------------------------------------------------------------------
    @builder.pre_tool("paper_content", "paper_idea_extractor")
    def _get_abstract_intro(state: Paper2FigureState):
        """
        Robustly extract Abstract + Introduction from PDF.
        """

        # 1. Read metadata title
        try:
            with open(state.paper_file, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                paper_title = reader.metadata.get('/Title', 'Unknown Title')
        except Exception:
            paper_title = "Unknown Title"

        # Open the PDF file using the path from state
        file_path = state.paper_file
        pdf_document = fitz.open(file_path)

        # Extract text from the first 10 pages
        text = ""
        for page_num in range(min(10, len(pdf_document))):
            page = pdf_document.load_page(page_num)
            text += page.get_text("text")

        content = text.strip()

        final_text = (
            f"The title of the paper is {paper_title}\n\n"
            f"Here's first ten page content: {content}"
        )

        log.info(f"{final_text}")
        return final_text
    
    @builder.pre_tool("paper_idea", "figure_desc_generator")
    def _get_paper_idea(state: Paper2FigureState):
        """
        Return paper ideas summary.
        """
        return state.paper_idea

    # ==============================================================
    # NODES
    # ==============================================================
    async def paper_idea_extractor_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        提取论文的关键贡献点
        """
        paper_idea_extractor = create_graph_agent("paper_idea_extractor", tool_manager=get_tool_manager())
        state = await paper_idea_extractor.execute(state, use_agent=True)
        return state
    
    async def figure_desc_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        图标提示词生成器节点
        """
        figure_desc_generator = create_graph_agent("figure_desc_generator", tool_manager=get_tool_manager())
        state = await figure_desc_generator.execute(state, use_agent=True)
        return state

    async def figure_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        图像生成或编辑节点：
        1) 先生成带内容的图 (fig_draft_path)
        2) 再基于该图进行一次固定提示词的二次编辑，生成空框模板图 (fig_layout_path)
        """
        prompt = state.agent_results.get("figure_desc_generator").get("results").get("fig_desc", {})
        safe_prompt = json.dumps(prompt, ensure_ascii=False)  # 确保中文字符正常显示

        edit_prompt = state.request.get("edit_prompt")
        image_path = state.request.get("prev_image")

        # 如果是二次编辑，prompt可以为空
        final_prompt = edit_prompt if image_path else safe_prompt

        log.info(f'final_prompt{final_prompt} - edit_prompt：{edit_prompt} - image_path：{image_path} - prompt：{safe_prompt}')

        # 统一输出根目录（outputs/paper2figure/<ts>）
        result_root = Path(_ensure_result_path(state))
        result_root.mkdir(parents=True, exist_ok=True)

        # 1) 生成带内容的图，直接存到 result_root
        fig_name = f"fig_{int(time.time())}.jpg"
        save_path = str(result_root / fig_name)

        await generate_or_edit_and_save_image_async(
            prompt=final_prompt,
            save_path=save_path,
            aspect_ratio=state.aspect_ratio,
            api_url=state.request.chat_api_url,
            api_key=os.getenv("DF_API_KEY") or state.request.chat_api_key,
            model=state.request.gen_fig_model,
            image_path=image_path,
            use_edit=True if image_path else False
        )
        state.agent_results["gen_img"] = {"path": save_path}
        state.fig_draft_path = save_path

        # 2) 基于第一次生成的图，做一次“空模板”二次编辑，也放在 result_root
        TEMPLATE_EDIT_PROMPT = (
            "Keep only the outermost rectangles and arrows(if any in the original box).\n"
            "Remove all inner content including title, subtitles, icons, explainary texts and all that.\n"
            "Keep the layout exactly the same.\n"
            "Output a description of an empty template composed of these boxes."
        )

        layout_name = f"layout_{int(time.time())}.jpg"
        layout_save_path = str(result_root / layout_name)
        await generate_or_edit_and_save_image_async(
            prompt=TEMPLATE_EDIT_PROMPT,
            save_path=layout_save_path,
            aspect_ratio=state.aspect_ratio,
            api_url=state.request.chat_api_url,
            api_key=os.getenv("DF_API_KEY") or state.request.chat_api_key,
            model=state.request.gen_fig_model,
            image_path=save_path,
            use_edit=True,
        )
        state.fig_layout_path = layout_save_path
        state.agent_results["gen_img_template"] = {"path": layout_save_path}

        return state

    async def figure_layout_sam_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        针对二次编辑后的空框模板图 (fig_layout_path) 进行:
        SAM 自动分割 -> 过滤 -> 裁剪子图 -> PNG->SVG->EMF，
        结果写入 state.layout_items，仅作为 PPT 背景框架层。
        """
        img_path = Path(state.fig_layout_path)
        if not img_path.exists():
            log.error(f"[figure_layout_sam] fig_layout_path 不存在: {img_path}")
            return state

        base_dir = Path(_ensure_result_path(state))
        out_dir = base_dir / "layout_items"
        out_dir.mkdir(parents=True, exist_ok=True)

        # 1. SAM 分割 + 过滤 + 裁剪子图
        layout_items = segment_layout_boxes(
            image_path=str(img_path),
            output_dir=str(out_dir),
            # 这里的参数可以根据 mask_detail_level 调整
            min_area=200,
            min_score=0.0,
            iou_threshold=0.3,
            top_k=None,
            nms_by="bbox",
        )

        # 2. 每个 layout PNG 转 SVG -> EMF
        for idx, it in enumerate(layout_items):
            png_path = it.get("png_path")
            if not png_path:
                continue

            svg_path = out_dir / f"layout_{idx}.svg"
            svg_abs = local_tool_for_raster_to_svg(
                {
                    "image_path": png_path,
                    "output_svg": str(svg_path),
                    "colormode": "color",
                    "hierarchical": "stacked",
                    "mode": "spline",
                }
            )
            it["svg_path"] = svg_abs

            emf_path = out_dir / f"layout_{idx}.emf"
            try:
                emf_abs = svg_to_emf(svg_abs, str(emf_path))
                it["emf_path"] = emf_abs
            except Exception as e:
                log.error(f"[figure_layout_sam] svg_to_emf failed for {svg_abs}: {e}")
                it["emf_path"] = None

        state.layout_items = layout_items
        log.info(f"[figure_layout_sam] 共生成 {len(layout_items)} 个布局元素")
        return state

    async def figure_mask_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        生成Figure进行元素切割，并提取 bbox + image_path 信息，递归处理子图。
        使用 MinerU HTTP 对原始带内容的图 (fig_draft_path) 做解析，得到内容层元素。
        规则：
        - 标题块(type == 'title') 保留为 text；
        - 其它所有块一律从顶层图裁剪出子图，当作 image，用于 icon / 局部视觉元素。
        """

        img_path = Path(state.fig_draft_path)
        if not img_path.exists():
            log.error(f"[figure_mask] fig_draft_path 不存在: {img_path}")
            return state

        # MinerU 所有中间结果统一放在本次 outputs 下
        base_dir = Path(_ensure_result_path(state))
        out_dir = base_dir / "mineru_recursive"
        out_dir.mkdir(parents=True, exist_ok=True)
        log.info(f"[figure_mask] MinerU 输出目录: {out_dir}")

        # MinerU 端口：优先从 state.request.mineru_port 读取，默认 8001
        port = getattr(state.request, "mineru_port", 8001)
        max_depth = getattr(state, "mask_detail_level", 2)

        log.critical(f"mask detail level : {max_depth} ")
        log.critical(f'[img_path]: {img_path}')
        log.critical(f'[mineru_port]: {port}')

        # 1. 调用新的 HTTP MinerU 递归处理，获取元素列表（归一化坐标）
        mineru_items = await recursive_mineru_layout(
            image_path=str(img_path),
            port=port,
            max_depth=max_depth,
            output_dir=out_dir,
        )

        # 顶层图像尺寸，用于 norm->pixel 映射与裁剪
        top_img = Image.open(state.fig_draft_path)
        top_w, top_h = top_img.size

        # 图标原图输出目录
        icons_raw_dir = base_dir / "icons_raw"
        icons_raw_dir.mkdir(parents=True, exist_ok=True)

        fig_mask = []
        icon_count = 0
        text_count = 0

        for idx, it in enumerate(mineru_items):
            elem_type_raw = it.get("type") or ""
            elem_type = elem_type_raw.lower()
            bbox = it.get("bbox")
            text = it.get("text") or ""

            if not bbox or len(bbox) != 4:
                continue

            # 归一化 -> 像素坐标
            x1n, y1n, x2n, y2n = bbox
            x1 = int(round(x1n * top_w))
            y1 = int(round(y1n * top_h))
            x2 = int(round(x2n * top_w))
            y2 = int(round(y2n * top_h))

            if x2 <= x1 or y2 <= y1:
                continue

            px_bbox = [x1, y1, x2, y2]

            # 1) 标题类块：仍然作为文本
            if elem_type in ["title"]:
                fig_mask.append(
                    {
                        "type": "text",
                        "bbox": px_bbox,
                        "text": text,
                        "text_level": 1,
                        "page_idx": 0,
                    }
                )
                text_count += 1
                continue

            # 2) 其它所有块：一律裁图，当作 image，用于 icon / 元素图层
            try:
                crop = top_img.crop((x1, y1, x2, y2))
                icon_path = icons_raw_dir / f"blk_{idx}.png"
                crop.save(icon_path)
                icon_abs = str(icon_path)
                fig_mask.append(
                    {
                        "type": "image",
                        "bbox": px_bbox,
                        "img_path": icon_abs,
                        "page_idx": 0,
                    }
                )
                icon_count += 1
            except Exception as e:
                log.error(f"[figure_mask] 裁剪子图失败 idx={idx}, bbox={px_bbox}: {e}")
                # 兜底：作为普通文本
                fig_mask.append(
                    {
                        "type": "text",
                        "bbox": px_bbox,
                        "text": text,
                        "text_level": None,
                        "page_idx": 0,
                    }
                )
                text_count += 1

        type_counter = {}
        for e in fig_mask:
            t = e.get("type")
            type_counter[t] = type_counter.get(t, 0) + 1

        log.info(
            f"[figure_mask] fig_mask size = {len(fig_mask)}, "
            f"type distribution = {type_counter}, "
            f"title_text={text_count}, icons(raw)={icon_count}"
        )

        # 更新 state 的 fig_mask 信息
        state.fig_mask = fig_mask
        log.info(f"[figure_mask] 共解析出 {len(fig_mask)} 个元素 (via MinerU HTTP, pixel bbox + raw icons)")

        return state
    
    async def figure_icon_bg_remover_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        把Mask里面的图标去除背景
        """
        base_dir = Path(_ensure_result_path(state))
        icons_dir = base_dir / "icons"
        icons_dir.mkdir(parents=True, exist_ok=True)

        img_cnt = 0
        for item in state.fig_mask:
            if item.get('type') in ['image', 'table']:
                img_cnt += 1
                output_path = local_tool_for_bg_remove({
                    "image_path": item.get('img_path'),
                    "model_path": state.request.bg_rm_model,
                    "output_dir": str(icons_dir)
                })
                if output_path:
                    item['img_path'] = output_path
                    log.info(f"[figure_icon_bg_remover] background removed: {output_path}")
                else:
                    log.warning(f"[figure_icon_bg_remover] bg remove failed for {item.get('img_path')}")
        log.info(f"[figure_icon_bg_remover] processed image/table elements: {img_cnt}")

        return state

    async def figure_ppt_generation_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        生成单页 PPT：
        - 页面尺寸与原始带内容图 fig_draft_path 一致
        - 背景为白色
        - 底层：根据 layout_items 的 EMF 元素按 bbox 放置
        - 上层：根据 MinerU 的 fig_mask 元素按原逻辑放置文本/图片/表格
        """
        try:
            from pptx.util import Emu

            # 从state获取输出目录（若未设置则自动初始化 outputs/paper2figure/<timestamp>）
            output_dir = Path(_ensure_result_path(state))
            output_dir.mkdir(parents=True, exist_ok=True)

            # 生成唯一文件名
            timestamp = int(time.time())
            ppt_filename = f"presentation_{timestamp}.pptx"
            ppt_path = output_dir / ppt_filename

            # 创建Presentation对象
            prs = Presentation()

            # 设置PPT尺寸，依据原始带内容图
            img = Image.open(state.fig_draft_path)
            width_px, height_px = img.size
            slide_width_px, slide_height_px = setup_presentation_size(prs, width_px, height_px)

            # 创建单页幻灯片
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # 白色背景
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)

            # 1) 先渲染 layout_items (SAM + SVG + EMF 背景层)
            layout_drawn = 0
            for item in state.layout_items or []:
                emf_path = item.get("emf_path")
                bbox = item.get("bbox")
                if not emf_path or not bbox:
                    continue

                if not os.path.exists(emf_path):
                    log.warning(f"[figure_ppt_generation] emf_path 不存在: {emf_path}")
                    continue

                x1, y1, x2, y2 = bbox
                left_px = x1 * slide_width_px
                top_px = y1 * slide_height_px
                width_box_px = (x2 - x1) * slide_width_px
                height_box_px = (y2 - y1) * slide_height_px

                left = Emu(left_px)
                top = Emu(top_px)
                width = Emu(width_box_px)
                height = Emu(height_box_px)

                try:
                    slide.shapes.add_picture(emf_path, left, top, width, height)
                    layout_drawn += 1
                except Exception as e:
                    log.error(f"[figure_ppt_generation] add_picture EMF 失败: {emf_path}, {e}")

            # 2) 再渲染 MinerU fig_mask（内容层）
            img_drawn = 0
            text_drawn = 0
            for element in state.fig_mask or []:
                elem_type = element.get('type', '')

                if elem_type == 'text':
                    add_text_element(slide, element)
                    text_drawn += 1
                elif elem_type in ['image', 'table']:
                    add_image_element(slide, element)
                    img_drawn += 1

            # 保存PPT
            prs.save(str(ppt_path))
            state.ppt_path = ppt_path
            print(f"PPT generated successfully: {ppt_path}")
            print(f"Slide size: {slide_width_px}x{slide_height_px} pixels")
            print(f"Total layout items: {len(state.layout_items)}, drawn: {layout_drawn}")
            print(f"Total content elements added: {len(state.fig_mask)}, text_drawn={text_drawn}, img_drawn={img_drawn}")

        except Exception as e:
            print(f"Error generating PPT: {e}")

        return state

    # ==============================================================
    # 注册 nodes / edges
    # ==============================================================
    def set_entry_node(state: Paper2FigureState) -> str:
        if(state.request.input_type == "PDF"):
            log.critical(f'进入PDF node ......')
            return "paper_idea_extractor"
        elif(state.request.input_type == "TEXT"):
            log.critical(f'进入TEXT node ......')
            return "figure_desc_generator"
        elif(state.request.input_type == "FIGURE"):
            log.critical(f'进入FIGURE node ......')
            return "figure_mask_generator"
        else:
            log.error(f"Invalid input type: {state.request.input_type}")
            return "_end_"

    def _init_result_path(state: Paper2FigureState) -> Paper2FigureState:
        """
        _start_ 节点：确保本次 workflow 有一个统一的 result_path 根目录。
        - 若用户已在 state.result_path 传入自定义目录，则直接使用该目录；
        - 若未传入，则初始化为 get_project_root()/outputs/paper2figure/<timestamp>。
        """
        _ensure_result_path(state)
        return state

    nodes = {
        '_start_': _init_result_path,
        "paper_idea_extractor": paper_idea_extractor_node,
        "figure_desc_generator": figure_desc_generator_node,
        "figure_generator": figure_generator_node,
        "figure_layout_sam": figure_layout_sam_node,
        "figure_mask_generator": figure_mask_generator_node,
        "figure_icon_bg_remover": figure_icon_bg_remover_node,
        "figure_ppt_generator": figure_ppt_generation_node,
        '_end_': lambda state: state,  # 终止节点
    }

    # ------------------------------------------------------------------
    # EDGES  (从节点 A 指向节点 B)
    # ------------------------------------------------------------------
    edges = [
        ("paper_idea_extractor", "figure_desc_generator"),
        ("figure_desc_generator", "figure_generator"),
        ("figure_generator", "figure_layout_sam"),
        ("figure_layout_sam", "figure_mask_generator"),
        ("figure_mask_generator", "figure_icon_bg_remover"),
        ("figure_icon_bg_remover", "figure_ppt_generator"),
        ("figure_ppt_generator", "_end_"),
    ]

    builder.add_nodes(nodes).add_edges(edges).add_conditional_edge("_start_", set_entry_node)
    return builder
