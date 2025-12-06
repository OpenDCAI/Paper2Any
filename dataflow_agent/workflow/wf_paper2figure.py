"""
icongen workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
生成时间: 2025-10-27 11:11:56

1. 在 **TOOLS** 区域定义需要暴露给 Prompt 的前置工具
2. 在 **NODES**  区域实现异步节点函数 (await-able)
3. 在 **EDGES**  区域声明有向边
4. 最后返回 builder.compile() 或 GenericGraphBuilder
"""

from __future__ import annotations
import asyncio
import json
import os
from dataflow_agent.state import MainState, Paper2FigureState
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder


from dataflow_agent.workflow.registry import register
# from dataflow_agent.agentroles import get_agent_cls, create_agent

from dataflow_agent.toolkits.tool_manager import get_tool_manager
from langchain.tools import tool
from langgraph.graph import StateGraph
from langgraph.prebuilt import ToolNode, tools_condition

from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger

from dataflow_agent.toolkits.imtool.req_img import generate_or_edit_and_save_image_async
from dataflow_agent.toolkits.imtool.bg_tool import local_tool_for_bg_remove
from dataflow_agent.agentroles import create_graph_agent

import re, pdfplumber, PyPDF2, time, shutil, fitz
import numpy as np
from PIL import Image

from dataflow_agent.utils import (
    build_output_directory,
    recursive_run_mineru,
    add_image_element,
    add_text_element,
    setup_presentation_size,
)

from pathlib import Path
import time, random
from pptx import Presentation
from pptx.dml.color import RGBColor 


log = get_logger(__name__)

def _ts_name(stem: str, ext: str = ".png") -> str:
    timestamp = int(time.time())  # 获取当前时间戳（秒）
    return f"./{stem}{timestamp}{ext}"

@register("paper2fig")
def create_p2fig_graph() -> GenericGraphBuilder:  # noqa: N802
    """
    Workflow factory: dfa run --wf paper2fig
    """
    builder = GenericGraphBuilder(state_model=Paper2FigureState,
                                  entry_point="_start_")  # 自行修改入口

    # ----------------------------------------------------------------------
    # TOOLS (pre_tool definitions)
    # ----------------------------------------------------------------------
    # 例:
    # @builder.pre_tool("purpose", "step1")
    # def _purpose(state: MainState):
    #     return "这里放入字符串 / 数值 / 列表 / 字典等供 prompt 使用"

    # @builder.post_tool('','')
    # def _post_tool1():
    # ----------------------------------------------------------------------
    @builder.pre_tool("paper_content", "paper_idea_extractor")
    def _get_abstract_intro(state: Paper2FigureState):
        """
        Robustly extract Abstract + Introduction from PDF.
        """

        # 1. Read metadata title
        try:
            with open(state.paper_file, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                paper_title = reader.metadata.get('/Title', 'Unknown Title')
        except Exception:
            paper_title = "Unknown Title"

        # 2. Read PDF text, collect lines
        # abs_and_intro = read_intro_from_paper(state.paper_file)

        # Open the PDF file using the path from state
        file_path = state.paper_file
        pdf_document = fitz.open(file_path)

        # Extract text from the first 5 pages
        text = ""
        for page_num in range(min(10, len(pdf_document))):  # Limit to first 5 pages
            page = pdf_document.load_page(page_num)
            text += page.get_text("text")  # Extract text content from the page

        # Store the extracted content in state
        content = text.strip()  # Strip any leading/trailing whitespace

        final_text = (
            f"The title of the paper is {paper_title}\n\n"
            f"Here's first ten page content: {content}"
        )

        log.info(f"{final_text}")
        return final_text
    
    @builder.pre_tool("paper_idea", "figure_desc_generator")
    def _get_paper_idea(state: Paper2FigureState):
        """
        Return paper ideas summary.
        """
        return state.paper_idea

    # ==============================================================
    # NODES
    # ==============================================================
    async def paper_idea_extractor_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        提取论文的关键贡献点
        """
        paper_idea_extractor = create_graph_agent("paper_idea_extractor", tool_manager=get_tool_manager())
        state = await paper_idea_extractor.execute(state, use_agent=True)
        return state
    
    async def figure_desc_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        图标提示词生成器节点
        """
        figure_desc_generator = create_graph_agent("figure_desc_generator", tool_manager=get_tool_manager())
        state = await figure_desc_generator.execute(state, use_agent=True)
        return state

    async def figure_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        图像生成或编辑节点
        """
        prompt = state.agent_results.get("figure_desc_generator").get("results").get("fig_desc", {})
        safe_prompt = json.dumps(prompt, ensure_ascii=False)  # 确保中文字符正常显示
        
        # prompt = "a cat in a tree."
        
        edit_prompt = state.request.get("edit_prompt")
        image_path = state.request.get("prev_image")

        # 如果是二次编辑，prompt可以为空
        final_prompt = edit_prompt if image_path else safe_prompt

        log.info(f'final_prompt{final_prompt} - edit_prompt：{edit_prompt} - image_path：{image_path} - prompt：{safe_prompt}')

        save_path = _ts_name("tmps/", ".jpg")

        # log.critical(f'use_edit: {False if image_path == "" else True}')

        await generate_or_edit_and_save_image_async(
            prompt=final_prompt,
            save_path=save_path,
            aspect_ratio = state.aspect_ratio,
            api_url=state.request.chat_api_url,
            api_key=os.getenv("DF_API_KEY"), 
            model=state.request.gen_fig_model,
            image_path=image_path,
            use_edit= True if image_path else False
            # edit_prompt=edit_prompt,
        )
        state.agent_results["gen_img"] = {"path": save_path}
        state.fig_draft_path = save_path
        shutil.copy(save_path, state.result_path)
        return state

    async def figure_mask_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        生成Figure进行元素切割，并提取 bbox + image_path 信息，递归处理子图。
        """

        img_path = Path(state.fig_draft_path)
        if not img_path.exists():
            log.error(f"[figure_mask] fig_draft_path 不存在: {img_path}")
            return state

        out_dir = build_output_directory(img_path)
        log.info(f"[figure_mask] MinerU 输出目录: {out_dir}")

        # 1. 调用递归的 mineru 处理，获取元素列表
        print("mask detail level", state.mask_detail_level)
        items = await recursive_run_mineru(img_path, out_dir, state.mask_detail_level)

        # 更新 state 的 fig_mask 信息
        state.fig_mask = items
        log.info(f"[figure_mask] 共解析出 {len(items)} 个元素")

        return state
    
    async def figure_icon_bg_remover_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        把Mask里面的图标去除背景
        """
        for item in state.fig_mask:
            if item.get('type') in ['image', 'table']:
                output_path = local_tool_for_bg_remove({
                    "image_path": item.get('img_path'),
                    "model_path": state.request.bg_rm_model,
                    "output_dir": state.result_path + "/icons"
                })
                if output_path:
                    item['img_path'] = output_path
                log.info(f"{item.get('img_path')} background removed.")

    async def figure_ppt_generation_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        基于图片的mask信息生成五页PPT，每一页使用不同的背景色
        """
        try:
            # 从state获取输出目录
            output_dir = Path(state.result_path)
            output_dir.mkdir(parents=True, exist_ok=True)
            
            # 生成唯一文件名
            timestamp = int(time.time())
            ppt_filename = f"presentation_{timestamp}.pptx"
            ppt_path = output_dir / ppt_filename
            state.ppt_path = ppt_path

            # 创建Presentation对象
            prs = Presentation()
            
            # 设置PPT尺寸
            img = Image.open(state.fig_draft_path)
            width_px, height_px = img.size
            slide_width_px, slide_height_px = setup_presentation_size(prs, width_px, height_px)
            
            # 预定义的五个背景色
            background_colors = ['#BCE0FE', '#E2F0D9', '#F2F2F2', '#FFF2CC', '#F2DCDB']
            
            # 创建五张幻灯片，每张幻灯片使用不同的背景色
            for i, selected_color in enumerate(background_colors):
                # 创建单页幻灯片
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # 设置背景色
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(
                    int(selected_color[1:3], 16), 
                    int(selected_color[3:5], 16), 
                    int(selected_color[5:7], 16)
                )
                
                # 添加所有元素到单页幻灯片
                for element in state.fig_mask:
                    elem_type = element.get('type', '')
                    
                    if elem_type == 'text':
                        add_text_element(slide, element)
                    elif elem_type in ['image', 'table']:
                        add_image_element(slide, element)
            
            # 保存PPT
            prs.save(str(ppt_path))
            
            print(f"PPT generated successfully: {ppt_path}")
            print(f"Slide size: {slide_width_px}x{slide_height_px} pixels")
            print(f"Total elements added: {len(state.fig_mask)}")
        
        except Exception as e:
            print(f"Error generating PPT: {e}")
        
        return state

    # ==============================================================
    # 注册 nodes / edges
    # ==============================================================
    def set_entry_node(state: Paper2FigureState) -> str:
        if(state.input_type == "PDF"):
            return "paper_idea_extractor"
        elif(state.input_type == "TEXT"):
            return "figure_desc_generator"
        elif(state.input_type == "FIGURE"):
            return "figure_mask_generator"
        else:
            log.error(f"Invalid input type: {state.input_type}")
            return "_end_"

    nodes = {
        '_start_': lambda state: state,
        "paper_idea_extractor": paper_idea_extractor_node,
        "figure_desc_generator": figure_desc_generator_node,
        "figure_generator": figure_generator_node,
        "figure_mask_generator": figure_mask_generator_node,
        "figure_icon_bg_remover": figure_icon_bg_remover_node,
        "figure_ppt_generator": figure_ppt_generation_node,
        '_end_': lambda state: state,  # 终止节点
    }

    # ------------------------------------------------------------------
    # EDGES  (从节点 A 指向节点 B)
    # ------------------------------------------------------------------
    edges = [
        ("paper_idea_extractor", "figure_desc_generator"),
        ("figure_desc_generator", "figure_generator"),
        ("figure_generator", "figure_mask_generator"),
        ("figure_mask_generator", "figure_icon_bg_remover"),
        ("figure_icon_bg_remover", "figure_ppt_generator"),
        ("figure_ppt_generator", "_end_"),
    ]

    builder.add_nodes(nodes).add_edges(edges).add_conditional_edge("_start_", set_entry_node)
    return builder


    # async def figure_mask_generator_node(state: Paper2FigureState) -> Paper2FigureState:
    #     """
    #     生成Figure进行元素切割，并提取 bbox + image_path 信息。
    #     """

    #     img_path = Path(state.fig_draft_path)
    #     if not img_path.exists():
    #         log.error(f"[figure_mask] fig_draft_path 不存在: {img_path}")
    #         return state

    #     out_dir = build_output_directory(img_path)
    #     log.info(f"[figure_mask] MinerU 输出目录: {out_dir}")

    #     # --- 1. 调用 mineru ---
    #     ok = await run_mineru(img_path, out_dir)
    #     if not ok:
    #         return state

    #     # --- 2. 找 JSON ---
    #     content_json = locate_content_json(out_dir)
    #     if content_json is None:
    #         return state

    #     # --- 3. 读取内容并修复路径 ---
    #     items = load_and_fix_items(content_json, out_dir)
    #     log.info(f"Layout Detection Info:  {items}")
    #     state.fig_mask = items

    #     log.info(f"[figure_mask] 共解析出 {len(items)} 个元素")
    #     return state

    # 假设你的模型路径现在存储在 state.request.sam2_model_path
    # async def figure_mask_generator_node(state: Paper2FigureState) -> Paper2FigureState:
    #     """
    #     生成Figure进行元素切割，保留坐标
    #     """

    #     # 从state中获取模型路径
    #     model_path = state.request.sam2_model

    #     # 加载模型
    #     generator = pipeline("mask-generation", model=model_path, device=0)

    #     # 获取原图路径
    #     original_image_path = state.fig_draft_path

    #     # 生成掩码
    #     outputs = generator(original_image_path, points_per_batch=64)

    #     # 加载原图
    #     original_image = Image.open(original_image_path)

    #     # 创建子图保存的目录
    #     base_name = os.path.splitext(os.path.basename(original_image_path))[0]  # 去掉文件后缀
    #     save_dir = os.path.join(os.path.dirname(original_image_path), f"{base_name}_sub_images")
    #     os.makedirs(save_dir, exist_ok=True)

    #     # 初始化一个空的mask信息列表
    #     mask_info = []
    #     valid_mask_count = 0

    #     # 遍历每个掩码，裁剪并保存子图
    #     for i, mask in enumerate(outputs["masks"]):
    #         # 转换mask为numpy数组 (binary: 0 and 1)
    #         mask_array = mask.numpy().astype(np.uint8)

    #         # 获取mask的bounding box（坐标范围）
    #         y_coords, x_coords = np.where(mask_array == 1)  # 获取掩码区域的所有坐标
    #         if len(y_coords) == 0 or len(x_coords) == 0:
    #             continue  # 如果没有找到有效的掩码区域，跳过

    #         # 计算bounding box
    #         top_left = (x_coords.min(), y_coords.min())  # (x_min, y_min)
    #         bottom_right = (x_coords.max(), y_coords.max())  # (x_max, y_max)

    #         # 裁剪原图得到子图
    #         sub_image = original_image.crop((top_left[0], top_left[1], bottom_right[0], bottom_right[1]))

    #         # 保存子图到指定目录
    #         sub_image_path = os.path.join(save_dir, f"sub_image_{i}.png")
    #         sub_image.save(sub_image_path)

    #         # 将子图的路径和坐标保存到mask_info
    #         mask_info.append({
    #             "sub_image_path": sub_image_path,
    #             "box_coord": [top_left, bottom_right]
    #         })
    #         valid_mask_count += 1

    #     # 将生成的mask信息保存到state.mask_info
    #     state.mask_info = mask_info
        
    #     # 只在关键结果处添加提示信息
    #     log.info(f"✅ 图像掩码生成完成，共处理 {valid_mask_count}/{len(outputs['masks'])} 个有效掩码")
    #     log.info(f"📁 子图保存目录: {save_dir}")
    #     log.info(f"📊 生成的掩码信息数量: {len(mask_info)}")

    #     return state