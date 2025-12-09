"""
paper2technical workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
生成时间: 2025-12-07 23:36:51

1. 在 **TOOLS** 区域定义需要暴露给 Prompt 的前置工具
2. 在 **NODES**  区域实现异步节点函数 (await-able)
3. 在 **EDGES**  区域声明有向边
4. 最后返回 builder.compile() 或 GenericGraphBuilder
"""

from __future__ import annotations
import json
import time
from pathlib import Path

from dataflow_agent.state import Paper2FigureState
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.workflow.registry import register
from dataflow_agent.agentroles import create_graph_agent, create_react_agent, create_simple_agent
from dataflow_agent.toolkits.tool_manager import get_tool_manager
from dataflow_agent.toolkits.imtool.bg_tool import (
    local_tool_for_svg_render,
    local_tool_for_raster_to_svg,
)
from dataflow_agent.toolkits.imtool.mineru_tool import (
    run_aio_two_step_extract,
    crop_mineru_blocks_by_type,
    crop_mineru_blocks_with_meta,
    svg_to_emf,
    recursive_mineru_layout,
)
from dataflow_agent.logger import get_logger
from dataflow_agent.utils import get_project_root

from pptx.util import Pt

log = get_logger(__name__)


def _strip_svg_text_nodes(svg_code: str) -> str:
    """
    兜底版 SVG 去文字函数，用于在 svg_bg_cleaner agent 失败时硬编码移除文本。

    处理策略（尽量“只删字不动图形”）：
    - 删除所有 <text>...</text> 块（支持跨行、多属性写法）；
    - 删除所有 <tspan>...</tspan> 块；
    - 删除 <title>...</title> 块；
    - 删除自闭合 text / tspan 节点（如 <text .../>，<tspan .../>）。
    其它图形元素（rect/circle/path/...）保留不动。
    """
    import re

    if not svg_code:
        return svg_code

    cleaned = svg_code

    # 1) 删除 <title>...</title>（SVG 文档标题/图标题）
    cleaned = re.sub(
        r"<title[^>]*?>.*?</title>",
        "",
        cleaned,
        flags=re.IGNORECASE | re.DOTALL,
    )

    # 2) 删除 <tspan>...</tspan>，避免保留行内文本
    cleaned = re.sub(
        r"<tspan[^>]*?>.*?</tspan>",
        "",
        cleaned,
        flags=re.IGNORECASE | re.DOTALL,
    )

    # 3) 删除 <text>...</text>，主要文字节点
    cleaned = re.sub(
        r"<text[^>]*?>.*?</text>",
        "",
        cleaned,
        flags=re.IGNORECASE | re.DOTALL,
    )

    # 4) 删除自闭合的 text / tspan 节点（无内容的标签）
    cleaned = re.sub(
        r"<text[^>]*/>",
        "",
        cleaned,
        flags=re.IGNORECASE,
    )
    cleaned = re.sub(
        r"<tspan[^>]*/>",
        "",
        cleaned,
        flags=re.IGNORECASE,
    )

    # 5) 简单压缩多余空行，避免产生大片空白
    lines = cleaned.splitlines()
    cleaned_lines = []
    for line in lines:
        if line.strip() == "":
            if cleaned_lines and cleaned_lines[-1].strip() == "":
                continue
        cleaned_lines.append(line)
    return "\n".join(cleaned_lines)


def _ensure_result_path(state: Paper2FigureState) -> str:
    """
    统一本次 workflow 的根输出目录：
    - 如果 state.result_path 已存在（通常由调用方传入，形如 时间戳+编码），直接使用；
    - 否则：使用 get_project_root() / "outputs" / "paper2tec" / <timestamp>，
      并回写到 state.result_path，确保后续节点共享同一目录，避免数据串台。
    """
    raw = getattr(state, "result_path", None)
    if raw:
        return raw

    root = get_project_root()
    ts = int(time.time())
    base_dir = (root / "outputs" / "paper2tec" / str(ts)).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)
    state.result_path = str(base_dir)
    return state.result_path


@register("paper2technical_bg_remove")
def create_paper2technical_graph() -> GenericGraphBuilder:  # noqa: N802
    """
    Workflow factory: dfa run --wf paper2technical
    """
    # 使用 Paper2FigureState，复用其中的 paper_file / paper_idea / fig_desc 等字段，
    # 这里不做图像生成和抠图，只负责“技术路线图”的 SVG + PPT 逻辑。
    builder = GenericGraphBuilder(
        state_model=Paper2FigureState,
        entry_point="_start_",        # 入口统一为 _start_，再由路由函数分发
    )

    # ----------------------------------------------------------------------
    # TOOLS (pre_tool definitions)
    # ----------------------------------------------------------------------
    # 1) 提供给 paper_idea_extractor 的 PDF 内容（标题 + 前几页正文）
    @builder.pre_tool("paper_content", "paper_idea_extractor")
    def _get_paper_content(state: Paper2FigureState):
        """
        前置工具: 读取论文 PDF 的标题和前若干页内容，供 paper_idea_extractor 节点使用。

        - 作用: 为大模型提供足够的上下文，让其抽取论文中的技术路线/实验流程关键信息。
        - 输出: 一个字符串，包含论文标题 + 前若干页文本。
        """
        import fitz  # PyMuPDF
        import PyPDF2

        pdf_path = state.paper_file
        if not pdf_path:
            log.warning("paper_file 为空，无法读取 PDF 内容")
            return ""

        try:
            with open(pdf_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                paper_title = reader.metadata.get("/Title", "Unknown Title")
        except Exception:
            paper_title = "Unknown Title"

        try:
            doc = fitz.open(pdf_path)
        except Exception as e:
            log.error(f"打开 PDF 失败: {e}")
            return f"The title of the paper is {paper_title}"

        text_parts: list[str] = []
        # 读取前 10 页内容，通常技术路线、整体框架会在前几页出现
        for page_idx in range(min(10, len(doc))):
            page = doc.load_page(page_idx)
            text_parts.append(page.get_text("text") or "")

        content = "\n".join(text_parts).strip()
        final_text = (
            f"The title of the paper is {paper_title}\n\n"
            f"Here are the first 10 pages of the paper:\n{content}"
        )
        log.info("paper_content 提取完成")
        return final_text

    # 2) 提供给技术路线图描述生成器的“论文核心想法/摘要”
    @builder.pre_tool("paper_idea", "technical_route_desc_generator")
    def _get_paper_idea(state: Paper2FigureState):
        """
        前置工具: 为 technical_route_desc_generator 节点暴露论文的核心想法摘要。

        - 在 PDF 模式下，该摘要由 paper_idea_extractor 节点写入 state.paper_idea。
        - 在 TEXT 模式下，可以直接由调用方事先把概要写入 state.paper_idea。
        """
        return state.paper_idea or ""

    # ----------------------------------------------------------------------

    # ==============================================================
    # NODES
    # ==============================================================
    async def paper_idea_extractor_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 1: 从 PDF 中抽取论文的核心思想 / 技术路线相关信息

        - 只在 input_type == "PDF" 时作为入口节点被调用。
        - 基于 pre_tool("paper_content") 提供的标题 + 前若干页内容，
          调用专门的 agent（例如 paper_idea_extractor）生成摘要。
        - 该摘要用于后续技术路线图描述生成。

        输入:
            state.paper_file : 论文 PDF 路径
        输出:
            state.paper_idea : 论文核心思想 / 技术路线要点摘要
            state.agent_results["paper_idea_extractor"] : agent 原始输出
        """
        agent = create_simple_agent("paper_idea_extractor")
        state = await agent.execute(state=state)
        return state

    async def technical_route_desc_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 2: 技术路线图描述生成器

        - 根据论文摘要（PDF 模式）或用户直接提供的文本描述（TEXT 模式），
          生成“技术路线/实验流程”的结构化自然语言描述或 JSON。
        - 典型内容包括: 各阶段实验步骤、模块之间的依赖关系、输入输出数据流等。

        输入:
            - PDF 模式: state.paper_idea 由 paper_idea_extractor 填充
            - TEXT 模式: 可以事先把文本写入 state.paper_idea 或其他字段
        输出:
            - 建议: 在 agent 内把结果存到 state.fig_desc 或 state.agent_results["technical_route_desc_generator"]
        """
        agent = create_react_agent(
            name="technical_route_desc_generator",
            max_retries=4,
            model_name="claude-haiku-4-5-20251001",
        )
        state = await agent.execute(state=state)

        # --------------------------------------------------------------
        # 将 LLM 生成的 SVG 源码渲染为实际图像文件，并写入统一的 result_path 目录
        # --------------------------------------------------------------
        svg_code = getattr(state, "figure_tec_svg_content", None)
        if svg_code:
            # 确保本次 workflow 的根输出目录已确定
            base_dir = Path(_ensure_result_path(state))
            base_dir.mkdir(parents=True, exist_ok=True)

            timestamp = int(time.time())
            # 同时输出 SVG 源码文件和 PNG 位图
            svg_output_path = str((base_dir / f"technical_route_{timestamp}.svg").resolve())
            svg_bg_output_path = str((base_dir / f"technical_route_{timestamp}_bg.svg").resolve())
            png_output_path_bg = str((base_dir / f"technical_route_{timestamp}_bg.png").resolve())
            png_output_path_full = str((base_dir / f"technical_route_{timestamp}_full.png").resolve())

            try:
                # 1) 保存原始 SVG 源码（含文字）
                Path(svg_output_path).write_text(svg_code, encoding="utf-8")
                state.svg_file_path = svg_output_path

                # 2) 调用 svg_bg_cleaner agent 生成“去文字版” SVG，失败时回退到本地函数
                svg_bg_code = None
                try:
                    cleaner_agent = create_react_agent(
                        name="svg_bg_cleaner",
                        max_retries=4,
                        model_name="claude-haiku-4-5-20251001",
                    )
                    # 将原始 SVG 挂到 state，方便 agent 读取
                    cleaner_state = await cleaner_agent.execute(state=state)
                    svg_bg_code = cleaner_state.svg_bg_code

                except Exception as e:
                    log.warning(f"svg_bg_cleaner agent 执行失败，回退到本地去文字函数: {e}")

                if not svg_bg_code:
                    svg_bg_code = _strip_svg_text_nodes(svg_code)

                Path(svg_bg_output_path).write_text(svg_bg_code, encoding="utf-8")
                state.svg_bg_file_path = svg_bg_output_path

                # 3) 用“纯背景 SVG”渲染 PNG 供 MinerU 使用（背景通路）
                png_bg_path = local_tool_for_svg_render(
                    {
                        "svg_code": svg_bg_code,
                        "output_path": png_output_path_bg,
                    }
                )
                state.svg_img_path = png_bg_path

                # 4) 额外渲染一份“带文字版” PNG 供 MinerU 抽取文本
                png_full_path = local_tool_for_svg_render(
                    {
                        "svg_path": svg_output_path,
                        "output_path": png_output_path_full,
                    }
                )
                state.svg_full_img_path = png_full_path

                log.critical(f"[state.svg_img_path]: {state.svg_img_path}")
                log.critical(f"[state.svg_full_img_path]: {state.svg_full_img_path}")
                log.critical(f"[state.svg_file_path]: {state.svg_file_path}")
                log.critical(f"[state.svg_bg_file_path]: {state.svg_bg_file_path}")
            except Exception as e:
                # 渲染或写文件失败时仅记录日志，避免打断整体 workflow
                log.error(f"technical_route_desc_generator_node: SVG 落盘/渲染失败: {e}")

        return state

    async def svg_fragment_miner_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 4: SVG 结构切分 / 小图块生成 (MinerU 接入)

        新策略：
        - 仅对“带文字版” PNG (full_image_path) 调用 MinerU，获取全局 layout (mineru_full)；
        - 使用 mineru_full 的 bbox 在“去文字版” PNG (bg_image_path) 上裁剪出纯背景小图；
        - 这些背景小图再进行 PNG->SVG->EMF 转换，用于 PPT 背景层；
        - 文本仍由 mineru_full 中的 text/title/image_caption block 提取，用于 PPT 文本框 overlay。
        """
        bg_image_path = getattr(state, "svg_img_path", None)
        full_image_path = getattr(state, "svg_full_img_path", None)

        if not bg_image_path or not full_image_path:
            log.error(
                "svg_fragment_miner_node: svg_img_path 或 svg_full_img_path 为空，"
                "无法调用 MinerU 进行双通路布局与裁剪"
            )
            return state

        mineru_port = 8001  # MinerU 服务端口

        if getattr(state, "agent_results", None) is None:
            state.agent_results = {}

        try:
            # -----------------------------
            # 1) layout 通路: 带文字 PNG -> MinerU -> mineru_full
            # -----------------------------
            log.info(
                f"svg_fragment_miner_node[full]: 调用 MinerU recursive_mineru_layout, "
                f"image_path={full_image_path}, port={mineru_port}"
            )
            mineru_full = await recursive_mineru_layout(
                image_path=full_image_path,
                port=mineru_port,
                max_depth=3,
            )
            state.agent_results["mineru_svg_fragment_full"] = mineru_full

            try:
                log.warning(
                    "svg_fragment_miner_node[full]: MinerU 返回结果 (JSON): "
                    + json.dumps(mineru_full, ensure_ascii=False, indent=2)
                )
            except Exception:
                log.warning(
                    "svg_fragment_miner_node[full]: MinerU 返回结果 (repr): "
                    + repr(mineru_full)
                )

            # -----------------------------
            # 2) 背景裁剪：用 mineru_full 的 bbox 在“去文字版” PNG 上裁剪纯背景小图
            # -----------------------------
            run_root = Path(_ensure_result_path(state))
            crop_output_dir = run_root / "crops" / "bg"

            # 根据需要选择哪些 block 用于生成背景原子块
            allowed_bg_types = {"image", "img", "figure", "table"}
            bg_blocks_for_crop = [
                b
                for b in mineru_full
                # if b.get("bbox") is not None
                # and len(b.get("bbox")) == 4
                # and (b.get("type") in allowed_bg_types)
            ]

            crops_with_meta = crop_mineru_blocks_with_meta(
                image_path=bg_image_path,          # 去文字 PNG
                blocks=mineru_full,         # 布局来自 full 图
                output_dir=str(crop_output_dir / "meta"),
                prefix="paper2technical_bg_",
            )

            svg_output_dir = crop_output_dir / "svgs"
            svg_output_dir.mkdir(parents=True, exist_ok=True)
            blocks_for_ppt: list[dict] = []

            for item in crops_with_meta:
                png_path = item.get("png_path")
                if not png_path:
                    continue

                try:
                    png_p = Path(png_path)
                    svg_path = str((svg_output_dir / f"{png_p.stem}.svg").resolve())
                    out_svg = local_tool_for_raster_to_svg(
                        {
                            "image_path": str(png_p),
                            "output_svg": svg_path,
                            "colormode": "color",
                        }
                    )
                    blocks_for_ppt.append(
                        {
                            "block_index": item.get("block_index"),
                            "type": item.get("type"),
                            "bbox": item.get("bbox"),  # 直接沿用 mineru_full 的归一化 bbox
                            "png_path": png_path,
                            "svg_path": out_svg,
                        }
                    )
                except Exception as e:
                    log.error(f"svg_fragment_miner_node[bg]: PNG->SVG 转换失败 {png_path}: {e}")

            state.agent_results["mineru_blocks_for_ppt"] = blocks_for_ppt

            # -----------------------------
            # 3) 文本通路：从 mineru_full 提取文本块 (text/title/image_caption)
            # -----------------------------
            # text_types = {"title", "text", "image_caption"}
            text_blocks = [
                {
                    "type": b.get("type"),
                    "bbox": b.get("bbox"),
                    "text": b.get("text"),
                    "depth": b.get("depth"),
                }
                for b in mineru_full
                # if b.get("bbox") is not None
                # and b.get("text") not in (None, "")
                # and b.get("type") in text_types
            ]

            state.agent_results["mineru_text_blocks"] = text_blocks
            # 若 Paper2FigureState 定义了 mineru_text_blocks 字段，可同步一份
            try:
                state.mineru_text_blocks = text_blocks  # type: ignore[attr-defined]
            except Exception:
                pass

        except Exception as e:
            log.error(f"svg_fragment_miner_node: MinerU 调用失败: {e}", exc_info=True)

        return state

    async def technical_ppt_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 5: 基于技术路线 SVG / 片段生成 PPT

        - 根据前面步骤生成的 SVG 代码或 svg_fragments，
          生成一份或多份 PPT 幻灯片，用于展示技术路线图。
        - 与 paper2figure 的 PPT 生成不同:
          - 这里不依赖位图图片和抠图，不需要图像背景去除模型；
          - 完全围绕“技术路线图”的结构信息进行排版。

        现在的策略：
        - 优先使用 state.agent_results["mineru_blocks_for_ppt"] 中的
          (svg_path, bbox) 信息，在一页 PPT 中根据 bbox 还原整体布局；
        - 若该字段不存在，则回退到旧逻辑：逐个 SVG/EMF 单独占一页。
        """
        from pptx import Presentation

        # 输出目录：统一使用本次 workflow 的根输出目录
        run_root = Path(_ensure_result_path(state))
        output_dir = run_root
        output_dir.mkdir(parents=True, exist_ok=True)

        timestamp = int(time.time())
        ppt_path = output_dir / f"technical_route_{timestamp}.pptx"

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # 优先尝试按 bbox 还原布局
        blocks_for_ppt: list[dict] = []
        if getattr(state, "agent_results", None):
            blocks_for_ppt = state.agent_results.get("mineru_blocks_for_ppt", []) or []

        # 预先将完整 SVG / 去文字 SVG 转为 EMF，用于后面单页展示
        full_svg_path = getattr(state, "svg_file_path", None)
        bg_svg_path = getattr(state, "svg_bg_file_path", None)

        full_emf = None
        bg_emf = None

        # 统一 EMF 输出目录
        emf_output_dir = output_dir / "ppt_emf"
        emf_output_dir.mkdir(parents=True, exist_ok=True)

        if full_svg_path:
            try:
                full_emf = svg_to_emf(
                    full_svg_path,
                    str((emf_output_dir / "technical_route_full.emf").resolve()),
                )
            except Exception as e:
                log.error(f"technical_ppt_generator_node: full SVG -> EMF 失败 {full_svg_path}: {e}")

        if bg_svg_path:
            try:
                bg_emf = svg_to_emf(
                    bg_svg_path,
                    str((emf_output_dir / "technical_route_bg.emf").resolve()),
                )
            except Exception as e:
                log.error(f"technical_ppt_generator_node: bg SVG -> EMF 失败 {bg_svg_path}: {e}")

        if blocks_for_ppt:
            # 在同一页 slide 上按 MinerU bbox 摆放所有 EMF 图块，并叠加文本框
            slide = prs.slides.add_slide(blank_slide_layout)

            # 1) 背景原子图块（EMF）
            for blk in blocks_for_ppt:
                svg_path = blk.get("svg_path")
                bbox = blk.get("bbox") or [0, 0, 1, 1]
                x1, y1, x2, y2 = bbox

                if not svg_path:
                    continue

                try:
                    svg_p = Path(svg_path)
                    emf_path = str((emf_output_dir / f"{svg_p.stem}.emf").resolve())
                    emf_abs = svg_to_emf(str(svg_p), emf_path)
                except Exception as e:
                    log.error(f"technical_ppt_generator_node: SVG -> EMF 失败 {svg_path}: {e}")
                    continue

                left = int(slide_width * x1)
                top = int(slide_height * y1)
                width = int(slide_width * (x2 - x1))
                height = int(slide_height * (y2 - y1))

                try:
                    slide.shapes.add_picture(
                        emf_abs,
                        left,
                        top,
                        width=width,
                        height=height,
                    )
                except Exception as e:
                    log.error(
                        f"technical_ppt_generator_node: 将 EMF 按 bbox 插入 PPT 失败 {emf_abs}: {e}"
                    )
                    continue

            # 2) 文本块 overlay：将 MinerU 抽取的文本按 bbox 作为文本框叠加到同一页
            text_blocks: list[dict] = []
            if getattr(state, "agent_results", None):
                text_blocks = state.agent_results.get("mineru_text_blocks", []) or []

            for tb in text_blocks:
                bbox = tb.get("bbox") or [0, 0, 1, 1]
                if len(bbox) != 4:
                    continue
                x1, y1, x2, y2 = bbox
                raw_text = tb.get("text") or ""
                text = raw_text.strip()
                if not text:
                    continue

                left = int(slide_width * x1)
                top = int(slide_height * y1)
                width = int(slide_width * (x2 - x1))
                height = int(slide_height * (y2 - y1))

                try:
                    tx_box = slide.shapes.add_textbox(left, top, width, height)
                    tf = tx_box.text_frame
                    tf.text = text

                    for p in tf.paragraphs:
                        p.font.size = Pt(12)
                        if tb.get("type") == "title":
                            p.font.bold = True
                            p.font.size = Pt(14)
                except Exception as e:
                    log.error(
                        f"technical_ppt_generator_node: 文本块插入失败 "
                        f"{text[:30]}...: {e}"
                    )

            # 额外页 1：完整带文字版 SVG 的 EMF
            if full_emf:
                slide_full = prs.slides.add_slide(blank_slide_layout)
                try:
                    slide_full.shapes.add_picture(
                        full_emf,
                        0,
                        0,
                        width=slide_width,
                        height=slide_height,
                    )
                except Exception as e:
                    log.error(
                        f"technical_ppt_generator_node: 将 full EMF 插入 PPT 失败 {full_emf}: {e}"
                    )

            # 额外页 2：完整去文字版 SVG 的 EMF
            if bg_emf:
                slide_bg = prs.slides.add_slide(blank_slide_layout)
                try:
                    slide_bg.shapes.add_picture(
                        bg_emf,
                        0,
                        0,
                        width=slide_width,
                        height=slide_height,
                    )
                except Exception as e:
                    log.error(
                        f"technical_ppt_generator_node: 将 bg EMF 插入 PPT 失败 {bg_emf}: {e}"
                    )

            prs.save(str(ppt_path))
            state.ppt_path = str(ppt_path)
            log.info(
                "technical_ppt_generator_node: PPT 已按 MinerU bbox 还原整体布局并叠加文本生成，"
                "并追加 full/bg 整图页面: "
                f"{ppt_path}"
            )
            return state

        # 若没有 blocks_for_ppt，退回旧逻辑：每个 SVG 单独一页居中缩放
        svg_paths: list[str] = []
        if getattr(state, "agent_results", None):
            svg_paths = state.agent_results.get("mineru_crops_svg", []) or []

        if not svg_paths:
            # 没有碎片 SVG，保留一页空白，用于兼容/调试
            prs.slides.add_slide(blank_slide_layout)
            log.warning(
                "technical_ppt_generator_node: 未找到 mineru_crops_svg，生成占位空白 PPT"
            )
        else:
            # 临时 EMF 输出目录
            emf_output_dir = output_dir / "ppt_emf_fallback"
            emf_output_dir.mkdir(parents=True, exist_ok=True)

            for svg_path in svg_paths:
                slide = prs.slides.add_slide(blank_slide_layout)

                try:
                    svg_p = Path(svg_path)
                    emf_path = str((emf_output_dir / f"{svg_p.stem}.emf").resolve())

                    # SVG -> EMF（保持矢量特性）
                    emf_abs = svg_to_emf(str(svg_p), emf_path)
                except Exception as e:
                    log.error(f"technical_ppt_generator_node: SVG -> EMF 失败 {svg_path}: {e}")
                    continue

                # 将 EMF 插入 PPT，先插入再按 80% 宽度缩放并居中
                try:
                    pic = slide.shapes.add_picture(emf_abs, 0, 0)

                    # 缩放到宽度 80%，保持纵横比
                    if pic.width and pic.width > 0:
                        scale = (slide_width * 0.8) / pic.width
                    else:
                        scale = 1.0

                    pic.width = int(pic.width * scale)
                    pic.height = int(pic.height * scale)

                    # 居中
                    pic.left = int((slide_width - pic.width) / 2)
                    pic.top = int((slide_height - pic.height) / 2)
                except Exception as e:
                    log.error(
                        f"technical_ppt_generator_node: 将 EMF 插入 PPT 失败 {emf_abs}: {e}"
                    )
                    continue

        prs.save(str(ppt_path))
        state.ppt_path = str(ppt_path)
        log.info(f"technical_ppt_generator_node: PPT 已生成: {ppt_path}")

        return state

    # ==============================================================
    # 注册 nodes / edges
    # ==============================================================

    def set_entry_node(state: Paper2FigureState) -> str:
        """
        路由函数: 根据输入类型选择技术路线工作流的入口节点。

        - input_type == "PDF"  : 从 PDF 中抽取论文想法，先走 paper_idea_extractor
        - input_type == "TEXT" : 直接使用调用方提供的文本描述，跳过 PDF 抽取，
                                 从 technical_route_desc_generator 开始
        其他值:
        - 认为是不合法输入，直接结束工作流。
        """
        input_type = getattr(state.request, "input_type", "PDF")
        if input_type == "PDF":
            log.critical("paper2technical: 进入 PDF 流程 (paper_idea_extractor)")
            return "paper_idea_extractor"
        elif input_type == "TEXT":
            log.critical("paper2technical: 进入 TEXT 流程 (technical_route_desc_generator)")
            return "technical_route_desc_generator"
        else:
            log.error(f"paper2technical: Invalid input type: {input_type}")
            return "_end_"

    def _init_result_path(state: Paper2FigureState) -> Paper2FigureState:
        """
        _start_ 节点：确保本次 workflow 有一个统一的 result_path 根目录。
        - 若用户已在 state.result_path 传入自定义目录，则直接使用该目录；
        - 若未传入，则初始化为 get_project_root()/outputs/paper2tec/<timestamp>。
        """
        _ensure_result_path(state)
        return state

    nodes = {
        "_start_": _init_result_path,
        "paper_idea_extractor": paper_idea_extractor_node,
        "technical_route_desc_generator": technical_route_desc_generator_node,
        "svg_fragment_miner": svg_fragment_miner_node,
        "technical_ppt_generator": technical_ppt_generator_node,
        "_end_": lambda state: state,  # 终止节点
    }

    # ------------------------------------------------------------------
    # EDGES  (从节点 A 指向节点 B)
    # ------------------------------------------------------------------
    edges = [
        # PDF 流程: 先抽想法，再生成技术路线描述
        ("paper_idea_extractor", "technical_route_desc_generator"),
        # PDF/TEXT 后续流程共用: 描述 -> 结构切分 -> PPT
        ("technical_route_desc_generator", "svg_fragment_miner"),
        ("svg_fragment_miner", "technical_ppt_generator"),
        ("technical_ppt_generator", "_end_"),
    ]

    builder.add_nodes(nodes).add_edges(edges).add_conditional_edge("_start_", set_entry_node)
    return builder
