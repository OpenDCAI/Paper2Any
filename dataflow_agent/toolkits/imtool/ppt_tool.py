# -*- coding: utf-8 -*-
"""
ppt_tool

本模块提供将一组页面图片转换为 PDF 与可编辑 PPTX 的工具函数。

功能概述：
- 从指定目录按自然顺序读取图片
- 生成包含所有图片页的 PDF 文件
- 使用 PaddleOCR 对页面进行 OCR，识别文本行
- 基于识别结果自动估计字体大小和颜色，将文本叠加到 PPTX 中
- 可选地对原始页面进行 inpaint，生成“去文字的干净底图”作为 PPT 背景

典型用法：
- 在 DataFlow-Agent 的图像处理流程中，作为从图片页到可编辑 PPT 文稿的后处理工具
- 也可在其它组件或脚本中通过对外函数直接调用
"""

import os
import re
from typing import Sequence, Optional, Dict, Any, List, Tuple

import numpy as np
from PIL import Image
import cv2
from paddleocr import PaddleOCR

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from dataflow_agent.utils import get_project_root

# ----------------------------
# Config (默认配置，可通过对外函数参数进行部分覆盖)
# ----------------------------
ADD_BACKGROUND_IMAGE = True
CLEAN_BACKGROUND = True  # 是否尝试抠掉文字、生成无字底图再叠加OCR文字
EXTRACT_TEXT_COLOR = True  # 是否提取原图文字颜色
INPAINT_METHOD = cv2.INPAINT_TELEA  # or cv2.INPAINT_NS
INPAINT_RADIUS = 7  # 增大修复半径（从3提高到7）
SIMPLE_BG_VAR_THRESH = 50.0  # 放宽阈值（从12提高到50）
MASK_DILATE_ITER = 2  # 增加膨胀次数（从1提高到2）
USE_ADAPTIVE_MASK = True  # 使用自适应mask生成

# 输出PPT比例（16:9）
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Debug：落盘你送入 OCR 的图片，方便肉眼确认内容/分辨率/是否被处理坏
DEBUG_DUMP_FIRST_N = 2
DEBUG_DIR = f"{get_project_root()}/tests/debug_frames"

# ---------- 核心修复：低分辨率页面的 OCR 前增强 ----------
UPSCALE_LONG_SIDE_TO = 2200   # 建议 2000~3200，越大越慢
UPSCALE_INTERP = cv2.INTER_CUBIC
ENABLE_SHARPEN = True         # 轻度锐化，提升边缘对比度
SHARPEN_AMOUNT = 0.8          # 0.0~1.5 之间

# 识别过滤阈值
DROP_SCORE = 30  # Tesseract的confidence是0-100的整数

# 字号优化配置（放大整体字号，明显拉开标题对正文的差距）
BASE_BODY_PT = 16.0  # 正文基准字号
FONT_SCALE_FACTOR = 1.0  # 全局字号缩放因子
TITLE_RATIO_MIN = 2.0  # 标题最小倍率
TITLE_RATIO_MAX = 3.5  # 标题最大倍率
SUBTITLE_RATIO_MIN = 1.4  # 副标题最小倍率
SUBTITLE_RATIO_MAX = 2.0  # 副标题最大倍率
BODY_RATIO_MIN = 0.9  # 正文最小倍率
BODY_RATIO_MAX = 1.1  # 正文最大倍率

# PaddleOCR 配置（全局只初始化一次）
PADDLE_OCR = PaddleOCR(
    use_angle_cls=True,  # 角度分类，处理横竖混排
    lang="ch",  # 中文 + 英文
)

# ----------------------------
# IO helpers
# ----------------------------


def natural_key(s: str):
    return [int(t) if t.isdigit() else t.lower() for t in re.split(r"(\d+)", s)]


def list_images_in_dir(d: str) -> List[str]:
    """
    按自然顺序列出目录中所有图片文件路径。
    """
    exts = (".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff")
    files = [f for f in os.listdir(d) if f.lower().endswith(exts)]
    files.sort(key=natural_key)
    return [os.path.join(d, f) for f in files]


def read_bgr(path: str) -> np.ndarray:
    """
    Robust image reader:
    - supports non-ascii paths (np.fromfile + imdecode)
    - returns BGR uint8 HxWx3
    """
    data = np.fromfile(path, dtype=np.uint8)
    img = cv2.imdecode(data, cv2.IMREAD_UNCHANGED)
    if img is None:
        img = cv2.imread(path, cv2.IMREAD_UNCHANGED)
    if img is None:
        raise ValueError(f"Failed to read image: {path}")

    # Normalize to BGR uint8
    if img.ndim == 2:
        img = cv2.cvtColor(img, cv2.COLOR_GRAY2BGR)
    elif img.ndim == 3 and img.shape[2] == 4:
        img = cv2.cvtColor(img, cv2.COLOR_BGRA2BGR)

    if img.dtype != np.uint8:
        img = np.clip(img, 0, 255).astype(np.uint8)

    return img


def debug_dump(img: np.ndarray, tag: str = "dbg") -> None:
    """
    将中间图像写入 DEBUG_DIR 方便调试。
    """
    os.makedirs(DEBUG_DIR, exist_ok=True)

    print(tag, "type:", type(img))
    if isinstance(img, np.ndarray):
        print(
            tag,
            "shape:",
            img.shape,
            "dtype:",
            img.dtype,
            "min/max:",
            int(img.min()),
            int(img.max()),
        )

    out_path = os.path.join(DEBUG_DIR, f"{tag}.png")
    ok = cv2.imwrite(out_path, img)
    print(tag, "saved:", out_path, "ok:", ok)


# ----------------------------
# PDF
# ----------------------------


def images_to_pdf(image_paths: Sequence[str], output_pdf_path: str) -> str:
    """
    将一组图片导出为单个 PDF 文件。
    """
    imgs: List[Image.Image] = []
    for p in image_paths:
        im = Image.open(p)
        if im.mode != "RGB":
            im = im.convert("RGB")
        imgs.append(im)
    if not imgs:
        raise ValueError("No images for PDF.")
    imgs[0].save(output_pdf_path, save_all=True, append_images=imgs[1:])
    return output_pdf_path


# ----------------------------
# Preprocess
# ----------------------------


def upscale_if_needed(
    bgr: np.ndarray,
    long_side_to: int = UPSCALE_LONG_SIDE_TO,
    interp: int = UPSCALE_INTERP,
):
    h, w = bgr.shape[:2]
    long_side = max(h, w)
    if long_side >= long_side_to:
        return bgr, 1.0

    scale = long_side_to / float(long_side)
    new_w = int(round(w * scale))
    new_h = int(round(h * scale))
    up = cv2.resize(bgr, (new_w, new_h), interpolation=interp)
    return up, scale


def sharpen(bgr: np.ndarray, amount: float = SHARPEN_AMOUNT) -> np.ndarray:
    """
    Unsharp mask style: sharpen = img*(1+a) - blur*a
    """
    if amount <= 0:
        return bgr
    blur = cv2.GaussianBlur(bgr, (0, 0), sigmaX=1.2, sigmaY=1.2)
    out = cv2.addWeighted(bgr, 1.0 + amount, blur, -amount, 0)
    return np.clip(out, 0, 255).astype(np.uint8)


def preprocess_for_ocr(bgr: np.ndarray):
    """
    Make a "det-friendly" version of the page:
    - upscale to a reasonable working resolution
    - optional sharpen
    """
    up, scale = upscale_if_needed(bgr)
    if ENABLE_SHARPEN:
        up = sharpen(up, amount=SHARPEN_AMOUNT)
    return up, scale


# ----------------------------
# OCR helpers
# ----------------------------


def is_cjk(s: str) -> bool:
    return any("\u4e00" <= ch <= "\u9fff" for ch in s)


def iou(a, b) -> float:
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    ix1, iy1 = max(ax1, bx1), max(ay1, by1)
    ix2, iy2 = min(ax2, bx2), min(ay2, by2)
    iw, ih = max(0, ix2 - ix1), max(0, iy2 - iy1)
    inter = iw * ih
    if inter <= 0:
        return 0.0
    area_a = (ax2 - ax1) * (ay2 - ay1)
    area_b = (bx2 - bx1) * (by2 - by1)
    return inter / (area_a + area_b - inter + 1e-6)


def merge_lines(
    lines: Sequence[Tuple[Sequence[float], str, float]], y_tol: int = 12, x_gap: int = 18
):
    """
    将OCR的word/短行合并成句子级别的行
    """
    if not lines:
        return []
    lines = sorted(lines, key=lambda x: (x[0][1], x[0][0]))

    def union(b1, b2):
        return [
            min(b1[0], b2[0]),
            min(b1[1], b2[1]),
            max(b1[2], b2[2]),
            max(b1[3], b2[3]),
        ]

    merged = []
    cur_bbox, cur_text, cur_conf_sum, cur_n = (
        lines[0][0],
        lines[0][1],
        lines[0][2],
        1,
    )

    for bbox, text, conf in lines[1:]:
        cy1 = (cur_bbox[1] + cur_bbox[3]) / 2
        cy2 = (bbox[1] + bbox[3]) / 2
        same_line = abs(cy1 - cy2) <= y_tol
        near_x = (bbox[0] - cur_bbox[2]) <= x_gap

        if same_line and near_x:
            cur_bbox = union(cur_bbox, bbox)
            if (not is_cjk(cur_text)) and (not is_cjk(text)):
                cur_text = (cur_text + " " + text).strip()
            else:
                cur_text = (cur_text + text).strip()
            cur_conf_sum += conf
            cur_n += 1
        else:
            merged.append((cur_bbox, cur_text, cur_conf_sum / cur_n))
            cur_bbox, cur_text, cur_conf_sum, cur_n = bbox, text, conf, 1

    merged.append((cur_bbox, cur_text, cur_conf_sum / cur_n))
    return merged


def text_score(lines) -> float:
    if not lines:
        return 0.0
    total_chars = sum(len(t) for (_, t, _) in lines)
    avg_conf = sum(conf for (_, _, conf) in lines) / max(1, len(lines))
    cjk_bonus = 1.1 if any(is_cjk(t) for (_, t, _) in lines) else 1.0
    return total_chars * (avg_conf / 100.0) * cjk_bonus  # normalize confidence to 0-1


def paddle_ocr(bgr: np.ndarray, drop_score: int = DROP_SCORE):
    """
    使用 PaddleOCR 识别整页图片
    返回格式：[(bbox, text, confidence), ...]
    bbox: [x1, y1, x2, y2]
    注意：这里直接在 BGR 图上跑，PaddleOCR 内部会处理颜色空间。
    """
    h, w = bgr.shape[:2]

    # ocr_result: List[List[ [box, (text, score)], ... ]]
    ocr_result = PADDLE_OCR.ocr(bgr, cls=True)
    lines = []

    if not ocr_result:
        return lines

    # 通常一页对应 ocr_result[0]
    for line in ocr_result[0]:
        box, (text, score) = line
        if not text:
            continue
        if score * 100.0 < drop_score:
            continue

        # box 是四点多边形：[ [x1,y1], [x2,y2], [x3,y3], [x4,y4] ]
        xs = [p[0] for p in box]
        ys = [p[1] for p in box]
        x1, y1, x2, y2 = min(xs), min(ys), max(xs), max(ys)

        # 边界裁剪
        x1 = max(0, min(w - 1, x1))
        x2 = max(0, min(w, x2))
        y1 = max(0, min(h - 1, y1))
        y2 = max(0, min(h, y2))
        if x2 <= x1 or y2 <= y1:
            continue

        bbox = [float(x1), float(y1), float(x2), float(y2)]
        # 保持 0-100 置信度范围，和原 Tesseract 逻辑兼容
        lines.append((bbox, text.strip(), float(score * 100.0)))

    return lines


# ----------------------------
# Color extraction
# ----------------------------


def extract_text_color(
    bgr: np.ndarray, bbox, bg_color=None
) -> Tuple[int, int, int]:
    """
    从文字区域提取主色调
    返回 (r, g, b) 元组
    """
    x1, y1, x2, y2 = [int(round(v)) for v in bbox]
    h, w = bgr.shape[:2]

    # 边界检查
    x1 = max(0, min(w - 1, x1))
    x2 = max(0, min(w, x2))
    y1 = max(0, min(h - 1, y1))
    y2 = max(0, min(h, y2))

    if x2 <= x1 or y2 <= y1:
        return (0, 0, 0)  # 默认黑色

    # 提取文字区域
    region = bgr[y1:y2, x1:x2]
    if region.size == 0:
        return (0, 0, 0)

    # 转换为RGB
    region_rgb = cv2.cvtColor(region, cv2.COLOR_BGR2RGB)
    pixels = region_rgb.reshape(-1, 3)

    # 如果像素太少，直接返回中位数颜色
    if len(pixels) < 10:
        median_color = np.median(pixels, axis=0).astype(int)
        return tuple(int(x) for x in median_color)

    # 使用K-means聚类找出主色调（2-3个聚类）
    try:
        from sklearn.cluster import KMeans

        n_clusters = min(3, len(pixels))
        kmeans = KMeans(n_clusters=n_clusters, random_state=42, n_init=10)
        kmeans.fit(pixels)

        # 获取聚类中心和每个聚类的像素数
        centers = kmeans.cluster_centers_
        labels = kmeans.labels_
        counts = np.bincount(labels)

        # 如果提供了背景色，排除接近背景色的聚类
        if bg_color is not None:
            bg_array = np.array(bg_color)
            valid_centers = []
            valid_counts = []

            for i, center in enumerate(centers):
                # 计算与背景色的距离
                dist = np.linalg.norm(center - bg_array)
                if dist > 30:  # 距离阈值
                    valid_centers.append(center)
                    valid_counts.append(counts[i])

            if valid_centers:
                centers = np.array(valid_centers)
                counts = np.array(valid_counts)

        # 选择出现频率最高的颜色
        dominant_idx = np.argmax(counts)
        dominant_color = centers[dominant_idx].astype(int)

        return tuple(int(x) for x in dominant_color)
    except Exception:
        # 如果sklearn不可用或出错，使用简单的中位数方法
        median_color = np.median(pixels, axis=0).astype(int)
        return tuple(int(x) for x in median_color)


def estimate_background_color(bgr: np.ndarray, lines):
    """
    估计背景主色调，用于颜色提取时排除背景
    """
    h, w = bgr.shape[:2]

    # 创建文字mask
    mask = np.ones((h, w), dtype=np.uint8) * 255
    for bbox, _, _ in lines:
        x1, y1, x2, y2 = [int(round(v)) for v in bbox]
        x1 = max(0, min(w - 1, x1))
        x2 = max(0, min(w, x2))
        y1 = max(0, min(h - 1, y1))
        y2 = max(0, min(h, y2))
        if x2 > x1 and y2 > y1:
            mask[y1:y2, x1:x2] = 0

    # 提取背景区域像素
    bg_pixels = bgr[mask > 0]
    if bg_pixels.size == 0:
        return None

    # 转换为RGB并计算中位数
    bg_rgb = cv2.cvtColor(
        bg_pixels.reshape(-1, 1, 3), cv2.COLOR_BGR2RGB
    ).reshape(-1, 3)
    median_bg = np.median(bg_rgb, axis=0).astype(int)

    return tuple(int(x) for x in median_bg)


# ----------------------------
# PPT helpers
# ----------------------------


def px_to_emu(px: float, emu_per_px: float) -> int:
    return int(px * emu_per_px)


def analyze_line_heights(lines) -> Optional[float]:
    """
    统计行高分布，估计"正文行高"
    """
    if not lines:
        return None
    hs = [max(1, b[3] - b[1]) for (b, _, _) in lines]
    return float(np.median(hs))


def classify_line_role(bbox, img_h_px: int, body_h_px: Optional[float]) -> str:
    """
    大致区分：title / subtitle / body
    """
    x1, y1, x2, y2 = bbox
    h = max(1, y2 - y1)
    if body_h_px is None or body_h_px <= 0:
        return "body"
    ratio = h / float(body_h_px)

    # 位置辅助：靠近页面顶部 + 较高行
    y_center = (y1 + y2) / 2.0
    top_region = img_h_px * 0.3

    if ratio > 1.7 and y_center < top_region:
        return "title"
    if ratio > 1.3:
        return "subtitle"
    return "body"


def estimate_font_pt(
    bbox, img_h_px: int, body_h_px: Optional[float], slide_h_in: float = SLIDE_H_IN
):
    """
    根据行高和所属角色估计字号（优化版）
    """
    role = classify_line_role(bbox, img_h_px, body_h_px)

    # 使用优化后的基准字号
    base_body_pt = BASE_BODY_PT * FONT_SCALE_FACTOR

    if body_h_px and body_h_px > 0:
        x1, y1, x2, y2 = bbox
        h_px = max(1, y2 - y1)
        ratio = h_px / body_h_px
    else:
        ratio = 1.0

    # 使用优化后的倍率范围
    if role == "title":
        pt = base_body_pt * max(TITLE_RATIO_MIN, min(TITLE_RATIO_MAX, ratio * 1.5))
        # 顶部大标题额外放大
        y_center = (bbox[1] + bbox[3]) / 2.0
        if y_center < img_h_px * 0.2:
            pt *= 1.2
    elif role == "subtitle":
        pt = base_body_pt * max(
            SUBTITLE_RATIO_MIN, min(SUBTITLE_RATIO_MAX, ratio * 1.2)
        )
    else:
        pt = base_body_pt * max(BODY_RATIO_MIN, min(BODY_RATIO_MAX, ratio))

    return Pt(max(8, min(72, pt)))


def add_background(
    slide, bgr: np.ndarray, slide_w_emu: int, slide_h_emu: int, tmp_path: str
) -> None:
    rgb = cv2.cvtColor(bgr, cv2.COLOR_BGR2RGB)
    pil = Image.fromarray(rgb)
    pil.save(tmp_path)
    slide.shapes.add_picture(tmp_path, 0, 0, width=slide_w_emu, height=slide_h_emu)
    os.remove(tmp_path)


def build_text_mask_from_lines(bgr: np.ndarray, lines) -> np.ndarray:
    """
    根据OCR行框生成初始mask（粗略矩形）
    """
    h, w = bgr.shape[:2]
    mask = np.zeros((h, w), dtype=np.uint8)

    for bbox, text, conf in lines:
        x1, y1, x2, y2 = [int(round(v)) for v in bbox]
        x1 = max(0, min(w - 1, x1))
        x2 = max(0, min(w, x2))
        y1 = max(0, min(h - 1, y1))
        y2 = max(0, min(h, y2))
        if x2 <= x1 or y2 <= y1:
            continue
        mask[y1:y2, x1:x2] = 255

    if MASK_DILATE_ITER > 0:
        kernel = np.ones((3, 3), np.uint8)
        mask = cv2.dilate(mask, kernel, iterations=MASK_DILATE_ITER)

    return mask


def build_adaptive_mask(bgr: np.ndarray, lines) -> np.ndarray:
    """
    使用自适应方法生成更精细的文字主mask
    结合OCR bbox和实际文字形状（内部边缘 + 阈值）
    """
    h, w = bgr.shape[:2]
    mask = np.zeros((h, w), dtype=np.uint8)

    gray = cv2.cvtColor(bgr, cv2.COLOR_BGR2GRAY)

    for bbox, text, conf in lines:
        x1, y1, x2, y2 = [int(round(v)) for v in bbox]
        x1 = max(0, min(w - 1, x1))
        x2 = max(0, min(w, x2))
        y1 = max(0, min(h - 1, y1))
        y2 = max(0, min(h, y2))
        if x2 <= x1 or y2 <= y1:
            continue

        region = gray[y1:y2, x1:x2]
        if region.size == 0:
            continue

        try:
            # 先看局部对比度
            if np.var(region) < 100:
                # 对比度很低，优先用 Canny 边缘找笔画
                edges = cv2.Canny(region, 50, 150)
                kernel = np.ones((2, 2), np.uint8)
                binary = cv2.dilate(edges, kernel, iterations=1)
            else:
                # 对比度正常，用阈值法
                if region.shape[0] < 20 or region.shape[1] < 20:
                    _, binary = cv2.threshold(
                        region, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU
                    )
                    binary = 255 - binary  # 反色：文字为白
                else:
                    binary = cv2.adaptiveThreshold(
                        region,
                        255,
                        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                        cv2.THRESH_BINARY_INV,
                        11,
                        2,
                    )
                kernel = np.ones((2, 2), np.uint8)
                binary = cv2.morphologyEx(binary, cv2.MORPH_CLOSE, kernel)

            mask[y1:y2, x1:x2] = cv2.bitwise_or(mask[y1:y2, x1:x2], binary)
        except Exception:
            mask[y1:y2, x1:x2] = 255

    if MASK_DILATE_ITER > 0:
        kernel = np.ones((3, 3), np.uint8)
        mask = cv2.dilate(mask, kernel, iterations=MASK_DILATE_ITER)

    return mask


def is_simple_background_region(bgr: np.ndarray, mask: np.ndarray) -> bool:
    """
    简单判定：mask 区域附近背景是否接近纯色（方差较小）
    """
    gray = cv2.cvtColor(bgr, cv2.COLOR_BGR2GRAY)
    # 扩大一点范围取邻域
    dilated = cv2.dilate((mask > 0).astype(np.uint8), np.ones((5, 5), np.uint8), iterations=1)
    region = gray[dilated > 0]
    if region.size == 0:
        return False
    var = float(np.var(region))
    return var < SIMPLE_BG_VAR_THRESH


def fill_with_neighbor(bgr: np.ndarray, mask: np.ndarray) -> np.ndarray:
    """
    对复杂背景时，优先用邻域像素粗略填充，再交给 inpaint 做平滑，
    避免 NS/TELEA 在大块区域产生奇怪纹理。
    """
    result = bgr.copy()
    h, w = mask.shape
    for y in range(h):
        xs = np.where(mask[y] > 0)[0]
        if len(xs) == 0:
            continue
        x_min, x_max = xs[0], xs[-1]
        left_src = max(0, x_min - 3)
        right_src = min(w - 1, x_max + 3)
        fill_color = (
            (bgr[y, left_src].astype(np.int32) + bgr[y, right_src].astype(np.int32))
            // 2
        ).astype(np.uint8)
        result[y, x_min : x_max + 1] = fill_color
    return result


def make_clean_background(bgr: np.ndarray, lines) -> np.ndarray:
    """
    使用改进的 inpaint 生成“无字版底图”：
    - 自适应主文字 mask
    - 扩展阴影/发光区域 mask 只用于 inpaint
    - 简单背景直接 inpaint，复杂背景先邻域填充再小半径 inpaint
    """
    if not lines:
        return bgr

    # 使用自适应或简单mask（主文字区域）
    if USE_ADAPTIVE_MASK:
        main_mask = build_adaptive_mask(bgr, lines)
    else:
        main_mask = build_text_mask_from_lines(bgr, lines)

    # 扩展阴影/发光区域，inpaint 时用这个大 mask
    shadow_mask = cv2.dilate(main_mask, np.ones((7, 7), np.uint8), iterations=2)

    is_simple = is_simple_background_region(bgr, shadow_mask)

    if is_simple:
        # 简单背景：直接 inpaint + 轻微模糊
        clean = cv2.inpaint(bgr, shadow_mask, INPAINT_RADIUS, cv2.INPAINT_TELEA)
    else:
        # 复杂背景：先用邻域像素粗填，再用小半径 NS 微调
        prefilled = fill_with_neighbor(bgr, shadow_mask)
        clean = cv2.inpaint(
            prefilled, shadow_mask, max(3, INPAINT_RADIUS // 2), cv2.INPAINT_NS
        )

    clean = cv2.GaussianBlur(clean, (3, 3), 0.5)

    # 只在 shadow_mask 区域应用 inpaint 结果
    result = bgr.copy()
    mask_3ch = cv2.cvtColor(shadow_mask, cv2.COLOR_GRAY2BGR) / 255.0
    result = (clean * mask_3ch + bgr * (1 - mask_3ch)).astype(np.uint8)

    return result


def ocr_images_to_ppt(
    image_paths: Sequence[str],
    output_pptx: str,
    add_background_image: bool = ADD_BACKGROUND_IMAGE,
    clean_background: bool = CLEAN_BACKGROUND,
    use_text_color: bool = EXTRACT_TEXT_COLOR,
) -> str:
    """
    将图片通过OCR转换为可编辑文字的PPT（优化版）

    注意：该函数为内部实现，推荐通过 images_to_pdf_and_ppt /
    convert_images_dir_to_pdf_and_ppt 间接调用。
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    for idx, img_path in enumerate(image_paths, start=1):
        print(f"Processing slide #{idx}: {os.path.basename(img_path)}")

        bgr = read_bgr(img_path)

        # 预处理：放大和锐化
        ocr_img, scale = preprocess_for_ocr(bgr)

        if idx <= DEBUG_DUMP_FIRST_N:
            debug_dump(bgr, f"before_ocr_raw_{idx}")
            debug_dump(ocr_img, f"before_ocr_up_{idx}")
            print(f"slide#{idx} upscale scale={scale:.3f}")

        h0, w0 = bgr.shape[:2]  # 原图尺寸
        h1, w1 = ocr_img.shape[:2]  # OCR输入尺寸

        # 创建幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # OCR识别（PaddleOCR 直接处理 BGR 图像）
        lines = paddle_ocr(ocr_img)

        # 把bbox从OCR图坐标缩回原图坐标
        if lines and (w1 != w0 or h1 != h0):
            sx = w0 / float(w1)
            sy = h0 / float(h1)
            lines = [
                ([b[0] * sx, b[1] * sy, b[2] * sx, b[3] * sy], t, c)
                for (b, t, c) in lines
            ]

        # 再合并一轮
        y_tol = max(12, int(h0 * 0.008))
        x_gap = max(18, int(w0 * 0.01))
        lines = merge_lines(lines, y_tol=y_tol, x_gap=x_gap)

        # 统计正文行高，用于后续字号估计
        body_h_px = analyze_line_heights(lines)

        if not lines:
            print(f"[WARN] slide#{idx} no text detected")
        else:
            print(f"slide#{idx} detected {len(lines)} text boxes")

        # 估计背景颜色（用于颜色提取）
        bg_color = None
        if use_text_color and lines:
            bg_color = estimate_background_color(bgr, lines)
            if bg_color:
                print(f"slide#{idx} estimated background color: RGB{bg_color}")

        # 底图处理：可选 inpaint 生成"干净底图"
        bg_for_slide = bgr
        if add_background_image:
            if clean_background and lines:
                print(f"slide#{idx} applying inpainting...")
                bg_for_slide = make_clean_background(bgr, lines)
                if idx <= DEBUG_DUMP_FIRST_N:
                    debug_dump(bg_for_slide, f"clean_bg_{idx}")
            tmp = f"__ppt_bg_{idx}.png"
            add_background(slide, bg_for_slide, slide_w_emu, slide_h_emu, tmp)

        scale_x = slide_w_emu / w0
        scale_y = slide_h_emu / h0

        for bbox, text, conf in lines:
            x1, y1, x2, y2 = bbox
            if (x2 - x1) < 6 or (y2 - y1) < 6:
                continue

            left = px_to_emu(x1, scale_x)
            top = px_to_emu(y1, scale_y)
            width = max(1, px_to_emu((x2 - x1), scale_x))
            height = max(1, px_to_emu((y2 - y1), scale_y))

            # 添加透明文本框
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = True

            # 设置文本框透明
            tb.fill.background()  # 无填充
            tb.line.fill.background()  # 无边框

            p = tf.paragraphs[0]
            p.text = text
            p.font.size = estimate_font_pt(bbox, img_h_px=h0, body_h_px=body_h_px)

            # 提取并设置文字颜色
            if use_text_color:
                text_color = extract_text_color(bgr, bbox, bg_color)
                p.font.color.rgb = RGBColor(*text_color)
            else:
                # 默认黑色
                p.font.color.rgb = RGBColor(0, 0, 0)

    prs.save(output_pptx)
    return output_pptx


# ----------------------------
# Public API
# ----------------------------


def images_to_pdf_and_ppt(
    image_paths: Sequence[str],
    output_pdf_path: Optional[str] = None,
    output_pptx_path: Optional[str] = None,
    add_background_image: bool = ADD_BACKGROUND_IMAGE,
    clean_background: bool = CLEAN_BACKGROUND,
    extract_text_color: bool = EXTRACT_TEXT_COLOR,
) -> Dict[str, Optional[str]]:
    """
    将给定的一组图片转换为 PDF 和可编辑 PPTX。

    参数:
        image_paths: 按页面顺序排列的图片路径列表。
        output_pdf_path: 输出 PDF 文件路径，若为 None 则不生成 PDF。
        output_pptx_path: 输出 PPTX 文件路径，若为 None 则不生成 PPT。
        add_background_image: 是否在 PPT 中加入整页背景图。
        clean_background: 是否对背景进行 inpaint 处理（在 add_background_image 为 True 时生效）。
        extract_text_color: 是否根据原图估计文字颜色，用于 PPT 文本着色。

    返回:
        包含已生成文件路径的字典，例如:
        {
            "pdf": "/path/to/output.pdf" 或 None,
            "pptx": "/path/to/output_editable.pptx" 或 None,
        }
    """
    result: Dict[str, Optional[str]] = {"pdf": None, "pptx": None}

    if output_pdf_path is not None:
        result["pdf"] = images_to_pdf(image_paths, output_pdf_path)

    if output_pptx_path is not None:
        result["pptx"] = ocr_images_to_ppt(
            image_paths=image_paths,
            output_pptx=output_pptx_path,
            add_background_image=add_background_image,
            clean_background=clean_background,
            use_text_color=extract_text_color,
        )

    return result


def convert_images_dir_to_pdf_and_ppt(
    input_dir: str,
    output_pdf_path: Optional[str] = None,
    output_pptx_path: Optional[str] = None,
    add_background_image: bool = ADD_BACKGROUND_IMAGE,
    clean_background: bool = CLEAN_BACKGROUND,
    extract_text_color: bool = EXTRACT_TEXT_COLOR,
) -> Dict[str, Optional[str]]:
    """
    给定图片目录，自动读取所有图片并生成 PDF + PPTX。

    参数:
        input_dir: 包含图片的目录，内部按文件名自然排序。
        其余参数同 images_to_pdf_and_ppt。

    返回:
        同 images_to_pdf_and_ppt。
    """
    image_paths = list_images_in_dir(input_dir)
    if not image_paths:
        raise ValueError(f"No images found in {input_dir!r}")

    return images_to_pdf_and_ppt(
        image_paths=image_paths,
        output_pdf_path=output_pdf_path,
        output_pptx_path=output_pptx_path,
        add_background_image=add_background_image,
        clean_background=clean_background,
        extract_text_color=extract_text_color,
    )
