from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest
from pydantic import ValidationError
from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from vendor.presentagent_runtime import EditablePPTRunRequest, run_from_pagecontent
from vendor.presentagent_runtime import cli, runner
from vendor.presentagent_runtime.contracts import EditablePPTRunArtifacts
from vendor.presentagent_runtime.planner import DeckIR
from vendor.presentagent_runtime.planner.pagecontent_adapter import (
    build_default_theme,
    normalize_image_slide,
    normalize_outline_slide,
    pagecontent_to_deck_ir,
)
from vendor.presentagent_runtime.planner.ir_models import SlideIR
from vendor.presentagent_runtime.renderer.artifact_store import build_artifacts


@pytest.fixture
def outline_pagecontent_slide() -> dict[str, object]:
    return {
        "title": "Research Highlights",
        "layout_description": "Title and bullets aligned on the left.",
        "key_points": ["Three datasets", "Deterministic pipeline", "Regression coverage"],
        "asset_ref": "assets/chart.png",
    }


@pytest.fixture
def direct_image_pagecontent_slide() -> dict[str, object]:
    return {
        "title": "Architecture Diagram",
        "layout_description": "A full-page diagram with no supporting bullets.",
        "ppt_img_path": "/tmp/architecture.png",
    }


def test_runtime_package_imports_cleanly() -> None:
    assert EditablePPTRunRequest is not None
    assert DeckIR is not None
    assert run_from_pagecontent is not None


def test_runtime_request_supports_agent_models() -> None:
    req = EditablePPTRunRequest(
        result_path="/tmp/runtime",
        pagecontent=[],
        model="gpt-5.1",
        api_url="https://llm.example.com/v1",
        api_key="secret",
        vlm_model="gpt-4.1-mini",
        vlm_api_url="https://vlm.example.com/v1",
        vlm_api_key="vlm-secret",
        image_model="gemini-3-pro-image-preview",
        image_api_url="https://image.example.com/v1",
        image_api_key="image-secret",
        enable_agent_planner=True,
        enable_material_resolution=True,
        enable_llm_codegen=False,
    )

    assert req.model == "gpt-5.1"
    assert req.vlm_model == "gpt-4.1-mini"
    assert req.image_model == "gemini-3-pro-image-preview"
    assert req.enable_agent_planner is True
    assert req.enable_material_resolution is True
    assert req.enable_llm_codegen is False


def test_pagecontent_deck_planner_uses_llm_response_to_enrich_ir() -> None:
    from vendor.presentagent_runtime.planner.brief_adapter import pagecontent_to_slide_briefs
    from vendor.presentagent_runtime.planner.llm_planner import PagecontentDeckPlanner

    class FakeClient:
        def chat(self, messages, temperature=0.2, response_format=None):
            return json.dumps(
                {
                    "deck_title": "LLM Planned Deck",
                    "deck_subtitle": "Built from pagecontent briefs",
                    "storyline_summary": "Cover -> evidence -> conclusion",
                    "planner_notes": ["llm planner engaged"],
                    "slides": [
                        {
                            "slide_id": "slide-001",
                            "title": "Migration Overview",
                            "core_message": "Lead with the migration summary.",
                            "objective": "Explain the new runtime direction.",
                            "layout_name": "cover",
                            "points": ["Phase 1 first", "Fallback renderer stays"],
                            "visual_intent": "hero diagram",
                        }
                    ],
                }
            )

    pagecontent = [
        {
            "title": "Intro",
            "layout_description": "Overview slide",
            "key_points": ["Phase 1", "Phase 2"],
        }
    ]
    deck_ir = pagecontent_to_deck_ir(pagecontent, language="en", style="clean")
    slide_briefs = pagecontent_to_slide_briefs(pagecontent, language="en", style="clean")

    planned = PagecontentDeckPlanner(FakeClient()).plan_deck(
        slide_briefs=slide_briefs,
        base_deck_ir=deck_ir,
        materials={"assets": []},
    )

    assert planned.title == "LLM Planned Deck"
    assert planned.subtitle == "Built from pagecontent briefs"
    assert planned.planner_notes == ["llm planner engaged"]
    assert planned.slides[0].title == "Migration Overview"
    assert planned.slides[0].core_message == "Lead with the migration summary."
    assert planned.slides[0].layout["name"] == "cover"
    assert planned.slides[0].points == ["Phase 1 first", "Fallback renderer stays"]


def test_material_resolver_uses_vlm_scores_to_select_best_asset(tmp_path: Path) -> None:
    from vendor.presentagent_runtime.materials.material_pipeline import collect_materials, resolve_materials

    asset_a = tmp_path / "asset_a.png"
    asset_b = tmp_path / "asset_b.png"
    image_bytes = (PROJECT_ROOT / "tests" / "test_02.png").read_bytes()
    asset_a.write_bytes(image_bytes)
    asset_b.write_bytes(image_bytes)

    deck_ir = pagecontent_to_deck_ir(
        [
            {
                "title": "Results",
                "layout_description": "Compare two candidate figures.",
                "key_points": ["Use the clearest chart"],
                "asset_paths": [str(asset_a), str(asset_b)],
            }
        ],
        language="en",
        style="clean",
        asset_base_dir=str(tmp_path),
    )
    material_manifest = collect_materials(deck_ir)

    class FakeDescriptor:
        def score_candidates(self, request, candidates):
            scored = []
            for candidate in candidates:
                score = 0.2 if candidate["path"] == str(asset_a) else 0.95
                scored.append({**candidate, "vlm_score": score, "score_reason": "fake"})
            return sorted(scored, key=lambda item: item["vlm_score"], reverse=True)

    resolved_deck, resolution = resolve_materials(deck_ir, material_manifest, descriptor=FakeDescriptor())

    assert resolved_deck.slides[0].selected_asset_path == str(asset_b)
    assert resolution["requests"][0]["resolved_candidate"]["path"] == str(asset_b)
    assert resolution["requests"][0]["candidate_pool"][0]["vlm_score"] == 0.95


def test_write_code_artifacts_uses_llm_generated_slide_code(tmp_path: Path) -> None:
    from vendor.presentagent_runtime.coder.pptx_code_agent import PPTXCodeAgent, write_code_artifacts

    deck_ir = pagecontent_to_deck_ir(
        [{"title": "Codegen", "key_points": ["One", "Two"]}],
        language="en",
        style="clean",
    )

    class FakeClient:
        def chat(self, messages, temperature=0.2, response_format=None):
            return """```python
def slide_001(prs, deck_ir, slide_ir, materials):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 3000000, 800000)
    textbox.text_frame.text = "Generated by LLM"
    return slide
```"""

    write_code_artifacts(
        tmp_path,
        deck_ir,
        code_agent=PPTXCodeAgent(FakeClient()),
    )

    slide_code = (tmp_path / "code" / "generated" / "slides" / "slide_001.py").read_text(encoding="utf-8")
    assert "Generated by LLM" in slide_code
    assert (tmp_path / "code" / "generated" / "build_deck.py").exists()


def test_code_agent_repairs_invalid_slide_code_before_persisting(tmp_path: Path) -> None:
    from vendor.presentagent_runtime.coder.pptx_code_agent import PPTXCodeAgent, write_code_artifacts

    deck_ir = pagecontent_to_deck_ir(
        [{"title": "Repair Slide", "key_points": ["Broken first", "Repair second"]}],
        language="en",
        style="clean",
    )
    observed_prompts: list[str] = []

    class FakeClient:
        def __init__(self) -> None:
            self.calls = 0

        def chat(self, messages, temperature=0.2, response_format=None):
            observed_prompts.append(messages[0]["content"])
            self.calls += 1
            if self.calls == 1:
                return """```python
def slide_001(prs, deck_ir, slide_ir, materials):
    return {
```"""
            return """```python
def slide_001(prs, deck_ir, slide_ir, materials):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 3000000, 800000)
    textbox.text_frame.text = "Repaired by LLM"
    return slide
```"""

    write_code_artifacts(
        tmp_path,
        deck_ir,
        code_agent=PPTXCodeAgent(FakeClient(), max_attempts=2),
    )

    slide_code = (tmp_path / "code" / "generated" / "slides" / "slide_001.py").read_text(encoding="utf-8")
    cache_payload = json.loads((tmp_path / "code" / "generated" / "cache" / "slide_001.json").read_text(encoding="utf-8"))

    assert "Repaired by LLM" in slide_code
    assert cache_payload["validated"] is True
    assert cache_payload["attempt_count"] == 2
    assert cache_payload["repair_count"] == 1
    assert cache_payload["function_name"] == "slide_001"
    assert observed_prompts[1].lower().count("repair") >= 1


def test_code_agent_repairs_runtime_invalid_slide_code_before_persisting(tmp_path: Path) -> None:
    from vendor.presentagent_runtime.coder.pptx_code_agent import PPTXCodeAgent, write_code_artifacts

    deck_ir = pagecontent_to_deck_ir(
        [{"title": "Runtime Repair Slide", "key_points": ["Broken first", "Repair second"]}],
        language="en",
        style="clean",
    )
    observed_prompts: list[str] = []

    class FakeClient:
        def __init__(self) -> None:
            self.calls = 0

        def chat(self, messages, temperature=0.2, response_format=None):
            observed_prompts.append(messages[0]["content"])
            self.calls += 1
            if self.calls == 1:
                return """```python
def slide_001(prs, deck_ir, slide_ir, materials):
    raise RuntimeError("first runtime failure")
```"""
            return """```python
def slide_001(prs, deck_ir, slide_ir, materials):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 3000000, 800000)
    textbox.text_frame.text = "Runtime repaired slide"
    return slide
```"""

    write_code_artifacts(
        tmp_path,
        deck_ir,
        code_agent=PPTXCodeAgent(FakeClient(), max_attempts=2),
    )

    validation_dir = tmp_path / "code" / "generated" / "validation" / "slide_001"
    slide_code = (tmp_path / "code" / "generated" / "slides" / "slide_001.py").read_text(encoding="utf-8")
    cache_payload = json.loads((tmp_path / "code" / "generated" / "cache" / "slide_001.json").read_text(encoding="utf-8"))

    assert "Runtime repaired slide" in slide_code
    assert cache_payload["validated"] is True
    assert cache_payload["attempt_count"] == 2
    assert cache_payload["repair_count"] == 1
    assert "first runtime failure" in observed_prompts[1]
    assert (validation_dir / "attempt_01.py").exists()
    assert (validation_dir / "attempt_02.py").exists()
    assert (validation_dir / "attempt_02.pptx").exists()
    assert not (validation_dir / "attempt_01.pptx").exists()


def test_write_code_artifacts_persists_slide_cache_metadata(tmp_path: Path) -> None:
    from vendor.presentagent_runtime.coder.pptx_code_agent import PPTXCodeAgent, write_code_artifacts

    deck_ir = pagecontent_to_deck_ir(
        [{"title": "Cache Slide", "key_points": ["One", "Two"]}],
        language="en",
        style="clean",
    )

    class FakeClient:
        def chat(self, messages, temperature=0.2, response_format=None):
            return """```python
def slide_001(prs, deck_ir, slide_ir, materials):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 3000000, 800000)
    textbox.text_frame.text = "Cache Metadata"
    return slide
```"""

    write_code_artifacts(
        tmp_path,
        deck_ir,
        code_agent=PPTXCodeAgent(FakeClient(), max_attempts=2),
    )

    cache_payload = json.loads((tmp_path / "code" / "generated" / "cache" / "slide_001.json").read_text(encoding="utf-8"))

    assert cache_payload["slide_id"] == "slide-001"
    assert cache_payload["function_name"] == "slide_001"
    assert cache_payload["validated"] is True
    assert cache_payload["attempt_count"] == 1
    assert cache_payload["repair_count"] == 0
    assert cache_payload["code_path"].endswith("code/generated/slides/slide_001.py")
    assert cache_payload["code_hash"]
    assert cache_payload["validation_output_path"].endswith("code/generated/validation/slide_001/attempt_01.pptx")


def test_write_code_artifacts_persists_validation_attempt_artifacts(tmp_path: Path) -> None:
    from vendor.presentagent_runtime.coder.pptx_code_agent import PPTXCodeAgent, write_code_artifacts

    deck_ir = pagecontent_to_deck_ir(
        [{"title": "Validation Slide", "key_points": ["One", "Two"]}],
        language="en",
        style="clean",
    )

    class FakeClient:
        def chat(self, messages, temperature=0.2, response_format=None):
            return """```python
def slide_001(prs, deck_ir, slide_ir, materials):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 3000000, 800000)
    textbox.text_frame.text = "Validated single slide"
    return slide
```"""

    write_code_artifacts(
        tmp_path,
        deck_ir,
        code_agent=PPTXCodeAgent(FakeClient(), max_attempts=2),
    )

    validation_dir = tmp_path / "code" / "generated" / "validation" / "slide_001"

    assert (validation_dir / "attempt_01.py").exists()
    assert (validation_dir / "attempt_01.pptx").exists()


def test_build_deck_script_imports_generated_slide_functions(tmp_path: Path) -> None:
    from vendor.presentagent_runtime.coder.pptx_code_agent import write_code_artifacts

    deck_ir = pagecontent_to_deck_ir(
        [
            {"title": "First", "key_points": ["One"]},
            {"title": "Second", "key_points": ["Two"]},
        ],
        language="en",
        style="clean",
    )

    write_code_artifacts(tmp_path, deck_ir)

    build_deck = (tmp_path / "code" / "generated" / "build_deck.py").read_text(encoding="utf-8")

    assert "from slides.slide_001 import slide_001" in build_deck
    assert "from slides.slide_002 import slide_002" in build_deck
    assert "def build_presentation(output_path):" in build_deck
    assert "slide_001(prs, deck_ir, deck_ir['slides'][0], materials)" in build_deck
    assert "slide_002(prs, deck_ir, deck_ir['slides'][1], materials)" in build_deck


def test_runner_uses_planner_vlm_and_code_agent_when_enabled(monkeypatch, tmp_path: Path) -> None:
    request = EditablePPTRunRequest(
        result_path=str(tmp_path / "runtime-agent"),
        pagecontent=[{"title": "Slide 1", "key_points": ["Point 1"]}],
        language="en",
        style="clean",
        model="gpt-test",
        api_url="https://llm.example.com/v1",
        api_key="llm-key",
        vlm_model="vlm-test",
        vlm_api_url="https://vlm.example.com/v1",
        vlm_api_key="vlm-key",
        enable_agent_planner=True,
        enable_material_resolution=True,
        enable_llm_codegen=True,
    )
    observed: dict[str, object] = {}

    class FakeClient:
        def __init__(self, *, api_key, api_base, model, client_type="llm"):
            observed.setdefault("clients", []).append(
                {
                    "api_key": api_key,
                    "api_base": api_base,
                    "model": model,
                    "client_type": client_type,
                }
            )
            self.client_type = client_type

    class FakePlanner:
        def __init__(self, client):
            observed["planner_client_type"] = client.client_type

        def plan_deck(self, slide_briefs, base_deck_ir, materials):
            observed["planner_called"] = True
            return base_deck_ir

    class FakeDescriptor:
        def __init__(self, client):
            observed["descriptor_client_type"] = client.client_type

        def score_candidates(self, request, candidates):
            return list(candidates)

    class FakeCodeAgent:
        def __init__(self, client):
            observed["code_agent_client_type"] = client.client_type

    def fake_write_code_artifacts(run_dir, deck_ir, code_agent=None):
        observed["code_artifacts_codegen"] = code_agent is not None
        generated_dir = Path(run_dir) / "code" / "generated"
        slides_dir = generated_dir / "slides"
        slides_dir.mkdir(parents=True, exist_ok=True)
        (generated_dir / "build_deck.py").write_text("# generated\n", encoding="utf-8")
        (slides_dir / "slide_001.py").write_text("# slide\n", encoding="utf-8")

    def fake_render_pptx(deck_ir, output_path):
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(b"pptx")

    monkeypatch.setattr(runner, "LLMClient", FakeClient)
    monkeypatch.setattr(runner, "PagecontentDeckPlanner", FakePlanner)
    monkeypatch.setattr(runner, "VLMDescriptor", FakeDescriptor)
    monkeypatch.setattr(runner, "PPTXCodeAgent", FakeCodeAgent)
    monkeypatch.setattr(runner, "write_code_artifacts", fake_write_code_artifacts)
    monkeypatch.setattr(runner, "render_pptx", fake_render_pptx)

    run_from_pagecontent(request)

    assert observed["planner_called"] is True
    assert observed["planner_client_type"] == "llm"
    assert observed["descriptor_client_type"] == "vlm"
    assert observed["code_agent_client_type"] == "llm"
    assert observed["code_artifacts_codegen"] is True
    assert observed["clients"] == [
        {
            "api_key": "llm-key",
            "api_base": "https://llm.example.com/v1",
            "model": "gpt-test",
            "client_type": "llm",
        },
        {
            "api_key": "vlm-key",
            "api_base": "https://vlm.example.com/v1",
            "model": "vlm-test",
            "client_type": "vlm",
        },
    ]


def test_run_from_pagecontent_falls_back_when_llm_planner_fails(monkeypatch, tmp_path: Path) -> None:
    request = EditablePPTRunRequest(
        result_path=str(tmp_path / "runtime-planner-fallback"),
        pagecontent=[{"title": "Slide 1", "key_points": ["Point 1"]}],
        language="en",
        style="clean",
        model="gpt-test",
        api_url="https://llm.example.com/v1",
        api_key="llm-key",
    )

    class FakeClient:
        def __init__(self, **kwargs):
            return None

    class ExplodingPlanner:
        def __init__(self, client):
            return None

        def plan_deck(self, slide_briefs, base_deck_ir, materials):
            raise RuntimeError("planner request failed")

    monkeypatch.setattr(runner, "LLMClient", FakeClient)
    monkeypatch.setattr(runner, "PagecontentDeckPlanner", ExplodingPlanner)

    artifacts = run_from_pagecontent(request)

    assert Path(artifacts.final_ir_path).exists()
    final_ir = json.loads(Path(artifacts.final_ir_path).read_text(encoding="utf-8"))
    log_text = Path(artifacts.log_path).read_text(encoding="utf-8")
    assert final_ir["slides"][0]["title"] == "Slide 1"
    assert "planner fallback" in log_text


def test_material_resolver_falls_back_when_vlm_descriptor_errors(tmp_path: Path) -> None:
    from vendor.presentagent_runtime.materials.material_pipeline import collect_materials, resolve_materials

    asset_a = tmp_path / "asset_a.png"
    asset_b = tmp_path / "asset_b.png"
    image_bytes = (PROJECT_ROOT / "tests" / "test_02.png").read_bytes()
    asset_a.write_bytes(image_bytes)
    asset_b.write_bytes(image_bytes)

    deck_ir = pagecontent_to_deck_ir(
        [
            {
                "title": "Results",
                "layout_description": "Compare two candidate figures.",
                "key_points": ["Use the clearest chart"],
                "asset_paths": [str(asset_a), str(asset_b)],
            }
        ],
        language="en",
        style="clean",
        asset_base_dir=str(tmp_path),
    )
    material_manifest = collect_materials(deck_ir)

    class ExplodingDescriptor:
        def score_candidates(self, request, candidates):
            raise RuntimeError("vlm scoring failed")

    resolved_deck, resolution = resolve_materials(deck_ir, material_manifest, descriptor=ExplodingDescriptor())

    assert resolved_deck.slides[0].selected_asset_path == str(asset_a)
    assert resolution["requests"][0]["attempt_log"][0]["descriptor_used"] is False
    assert resolution["requests"][0]["attempt_log"][0]["descriptor_error"] == "vlm scoring failed"


def test_write_code_artifacts_falls_back_when_llm_codegen_fails(tmp_path: Path) -> None:
    from vendor.presentagent_runtime.coder.pptx_code_agent import PPTXCodeAgent, write_code_artifacts

    deck_ir = pagecontent_to_deck_ir(
        [{"title": "Codegen", "key_points": ["One", "Two"]}],
        language="en",
        style="clean",
    )

    class ExplodingClient:
        def chat(self, messages, temperature=0.2, response_format=None):
            raise RuntimeError("codegen request failed")

    write_code_artifacts(
        tmp_path,
        deck_ir,
        code_agent=PPTXCodeAgent(ExplodingClient()),
    )

    slide_code = (tmp_path / "code" / "generated" / "slides" / "slide_001.py").read_text(encoding="utf-8")
    assert "Generated by LLM" not in slide_code
    assert "selected_asset_path" in slide_code


def test_runner_persists_same_refined_ir_used_for_rendering(monkeypatch) -> None:
    raw_deck = DeckIR(
        title="Raw Deck",
        language="en",
        style="clean",
        slides=[SlideIR(slide_id="slide-1", page_num=1, title="Raw Slide")],
    )
    refined_deck = DeckIR(
        title="Refined Deck",
        language="en",
        style="clean",
        slides=[SlideIR(slide_id="slide-1", page_num=1, title="Refined Slide")],
    )
    request = EditablePPTRunRequest(
        result_path="/tmp/paper2ppt-code-runtime",
        pagecontent=[{"title": "Slide 1", "key_points": ["Point 1"]}],
        language="en",
        style="clean",
    )
    observed: dict[str, object] = {}

    def fake_pagecontent_to_deck_ir(pagecontent, *, language, style, asset_base_dir=""):
        observed["pagecontent"] = pagecontent
        observed["language"] = language
        observed["style"] = style
        observed["asset_base_dir"] = asset_base_dir
        return raw_deck

    def fake_refine_ir(deck_ir):
        observed["refine_input"] = deck_ir
        return refined_deck

    def fake_build_artifacts(run_dir):
        observed["artifact_run_dir"] = run_dir
        return EditablePPTRunArtifacts(
            run_dir=f"{run_dir}/code_runtime",
            ir_path=f"{run_dir}/code_runtime/ir/deck_ir.json",
            recipe_path=f"{run_dir}/code_runtime/recipes/render_recipe.json",
            pptx_path=f"{run_dir}/code_runtime/exports/paper2ppt_code_editable.pptx",
            pdf_path=f"{run_dir}/code_runtime/exports/paper2ppt_code_preview.pdf",
            log_path=f"{run_dir}/code_runtime/logs/run.log",
        )

    def fake_render_recipe(deck_ir):
        observed["recipe_deck_ir"] = deck_ir
        return {"deck_title": deck_ir.title, "slide_count": len(deck_ir.slides)}

    def fake_write_json(path, payload):
        observed.setdefault("json_writes", []).append((str(path), payload))

    def fake_render_pptx(deck_ir, output_path):
        observed["render_pptx"] = (deck_ir, str(output_path))

    def fake_export_pdf_preview(pptx_path, pdf_path, log_lines):
        observed["export_pdf_preview"] = (str(pptx_path), str(pdf_path))
        return str(pdf_path)

    def fake_write_log(path, lines):
        observed["log_write"] = (str(path), list(lines))

    monkeypatch.setattr(runner, "pagecontent_to_deck_ir", fake_pagecontent_to_deck_ir)
    monkeypatch.setattr(runner, "refine_ir", fake_refine_ir)
    monkeypatch.setattr(runner, "build_artifacts", fake_build_artifacts)
    monkeypatch.setattr(runner, "render_recipe", fake_render_recipe)
    monkeypatch.setattr(runner, "write_json", fake_write_json)
    monkeypatch.setattr(runner, "render_pptx", fake_render_pptx)
    monkeypatch.setattr(runner, "_export_pdf_preview", fake_export_pdf_preview)
    monkeypatch.setattr(runner, "write_log", fake_write_log)

    artifacts = run_from_pagecontent(request)

    assert observed["pagecontent"] == request.pagecontent
    assert observed["language"] == "en"
    assert observed["style"] == "clean"
    assert observed["asset_base_dir"] == request.result_path
    assert observed["refine_input"] is raw_deck
    assert observed["artifact_run_dir"] == request.result_path
    assert observed["recipe_deck_ir"] is refined_deck
    assert observed["render_pptx"] == (
        refined_deck,
        "/tmp/paper2ppt-code-runtime/code_runtime/exports/paper2ppt_code_editable.pptx",
    )
    assert observed["export_pdf_preview"] == (
        "/tmp/paper2ppt-code-runtime/code_runtime/exports/paper2ppt_code_editable.pptx",
        "/tmp/paper2ppt-code-runtime/code_runtime/exports/paper2ppt_code_preview.pdf",
    )
    assert observed["json_writes"] == [
        (
            "/tmp/paper2ppt-code-runtime/code_runtime/ir/deck_ir.json",
            refined_deck.model_dump(),
        ),
        (
            "/tmp/paper2ppt-code-runtime/code_runtime/recipes/render_recipe.json",
            {"deck_title": "Refined Deck", "slide_count": 1},
        ),
    ]
    assert observed["log_write"][0] == "/tmp/paper2ppt-code-runtime/code_runtime/logs/run.log"
    assert any("wrote editable pptx" in line for line in observed["log_write"][1])
    assert artifacts == EditablePPTRunArtifacts(
        run_dir="/tmp/paper2ppt-code-runtime/code_runtime",
        ir_path="/tmp/paper2ppt-code-runtime/code_runtime/ir/deck_ir.json",
        recipe_path="/tmp/paper2ppt-code-runtime/code_runtime/recipes/render_recipe.json",
        pptx_path="/tmp/paper2ppt-code-runtime/code_runtime/exports/paper2ppt_code_editable.pptx",
        pdf_path="/tmp/paper2ppt-code-runtime/code_runtime/exports/paper2ppt_code_preview.pdf",
        log_path="/tmp/paper2ppt-code-runtime/code_runtime/logs/run.log",
    )


@pytest.mark.parametrize(
    ("page_num", "item", "expected_layout"),
    [
        (
            1,
            {"title": "Cover Slide", "key_points": ["First point", "Second point"]},
            "cover",
        ),
        (
            2,
            {"title": "Key Takeaways", "key_points": ["One", "Two", "Three"]},
            "bullets",
        ),
        (
            2,
            {
                "title": "Results",
                "key_points": ["Accuracy gains", "Latency held"],
                "asset_ref": "assets/chart.png",
            },
            "two_column",
        ),
        (
            2,
            {"title": "System Diagram", "asset_ref": "assets/diagram.png"},
            "image_focus",
        ),
    ],
)
def test_normalize_outline_slide_maps_supported_layouts(
    page_num: int,
    item: dict[str, object],
    expected_layout: str,
) -> None:
    slide = normalize_outline_slide(item, page_num)

    assert slide.layout_type == expected_layout


def test_normalize_image_slide_maps_full_bleed_image_layout(
    direct_image_pagecontent_slide: dict[str, object],
) -> None:
    slide = normalize_image_slide(direct_image_pagecontent_slide, page_num=2)

    assert slide.layout_type == "full_bleed_image"
    assert slide.asset_paths == ["/tmp/architecture.png"]


def test_pagecontent_to_deck_ir_routes_canonical_asset_paths_to_full_bleed_image() -> None:
    deck = pagecontent_to_deck_ir(
        [
            {
                "title": "Canonical Diagram",
                "layout_description": "Full-page architecture visual.",
                "asset_paths": [" /tmp/canonical-diagram.png "],
            }
        ],
        language="en",
        style="clean",
    )

    assert len(deck.slides) == 1
    assert deck.slides[0].layout_type == "full_bleed_image"
    assert deck.slides[0].asset_paths == ["/tmp/canonical-diagram.png"]
    assert deck.slides[0].bullets == []


def test_pagecontent_to_deck_ir_normalizes_renderer_friendly_slide_fields(
    outline_pagecontent_slide: dict[str, object],
    direct_image_pagecontent_slide: dict[str, object],
) -> None:
    outline_pagecontent_slide = {
        **outline_pagecontent_slide,
        "title": "  Research Highlights  ",
        "layout_description": "  Title and bullets aligned on the left.  ",
        "key_points": [" Three datasets ", "", "  ", None, " Regression coverage "],
        "asset_ref": " assets/chart.png ",
    }
    direct_image_pagecontent_slide = {
        **direct_image_pagecontent_slide,
        "title": "  Architecture Diagram  ",
        "layout_description": "  A full-page diagram with no supporting bullets.  ",
        "ppt_img_path": " /tmp/architecture.png ",
    }

    deck = pagecontent_to_deck_ir(
        [outline_pagecontent_slide, direct_image_pagecontent_slide],
        language="en",
        style="clean",
    )

    assert len(deck.slides) == 2

    outline_slide = deck.slides[0]
    assert outline_slide.title == "Research Highlights"
    assert outline_slide.summary == "Title and bullets aligned on the left."
    assert outline_slide.bullets == [
        "Three datasets",
        "Regression coverage",
    ]
    assert outline_slide.asset_paths == ["assets/chart.png"]

    image_slide = deck.slides[1]
    assert image_slide.title == "Architecture Diagram"
    assert image_slide.layout_type == "full_bleed_image"
    assert image_slide.summary == "A full-page diagram with no supporting bullets."
    assert image_slide.bullets == []
    assert image_slide.asset_paths == ["/tmp/architecture.png"]


def test_slide_ir_rejects_unknown_layout_types() -> None:
    with pytest.raises(ValidationError):
        SlideIR(slide_id="slide-1", page_num=1, layout_type="poster")


@pytest.mark.parametrize(
    ("legacy_field", "value"),
    [
        ("asset_ref", "assets/chart.png"),
        ("ppt_img_path", "/tmp/slide.png"),
    ],
)
def test_slide_ir_rejects_legacy_source_specific_asset_fields(
    legacy_field: str,
    value: str,
) -> None:
    with pytest.raises(ValidationError):
        SlideIR(slide_id="slide-1", page_num=1, **{legacy_field: value})


def test_build_default_theme_returns_usable_theme_for_empty_style() -> None:
    theme = build_default_theme(style="", language="zh")

    assert theme.primary
    assert theme.accent
    assert theme.background
    assert theme.title_font
    assert theme.body_font


def test_cli_accepts_pagecontent_json_payload(monkeypatch) -> None:
    captured: dict[str, object] = {}

    def fake_run_from_pagecontent(req):
        captured["request"] = req
        return EditablePPTRunArtifacts(
            run_dir=f"{req.result_path}/code_runtime",
            ir_path=f"{req.result_path}/code_runtime/ir/deck_ir.json",
            recipe_path=f"{req.result_path}/code_runtime/recipes/render_recipe.json",
            pptx_path=f"{req.result_path}/code_runtime/exports/paper2ppt_code_editable.pptx",
            pdf_path=f"{req.result_path}/code_runtime/exports/paper2ppt_code_preview.pdf",
            log_path=f"{req.result_path}/code_runtime/logs/run.log",
        )

    monkeypatch.setattr(cli, "run_from_pagecontent", fake_run_from_pagecontent)
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "presentagent-runtime",
            "--result-path",
            "/tmp/paper2ppt-code-runtime-cli",
            "--language",
            "en",
            "--style",
            "clean",
            "--pagecontent-json",
            '[{"title":"CLI Slide","ppt_img_path":"/tmp/slide.png"}]',
        ],
    )

    exit_code = cli.main()

    assert exit_code == 0
    assert captured["request"] == EditablePPTRunRequest(
        result_path="/tmp/paper2ppt-code-runtime-cli",
        pagecontent=[{"title": "CLI Slide", "ppt_img_path": "/tmp/slide.png"}],
        language="en",
        style="clean",
    )


def test_build_artifacts_initializes_pdf_path_as_blank(tmp_path: Path) -> None:
    artifacts = build_artifacts(str(tmp_path / "runtime-artifacts"))

    assert artifacts.run_dir.endswith("code_runtime")
    assert artifacts.pdf_path == ""
    assert Path(artifacts.materials_manifest_path).parent.exists()
    assert Path(artifacts.planned_ir_path).parent.exists()
    assert Path(artifacts.final_ir_path).parent.exists()
    assert artifacts.ir_path == artifacts.final_ir_path


def test_run_from_pagecontent_persists_slide_briefs_materials_and_final_ir(
    tmp_path: Path,
) -> None:
    result_dir = tmp_path / "runtime-rich-result"
    request = EditablePPTRunRequest(
        result_path=str(result_dir),
        pagecontent=[
            {
                "title": "Editable Runtime",
                "layout_description": "Subtitle for the cover slide.",
                "key_points": ["Deterministic", "Editable output"],
            },
            {
                "title": "Image Focus",
                "layout_description": "Diagram and supporting notes.",
                "key_points": ["One asset", "Short bullets"],
                "asset_ref": str(PROJECT_ROOT / "tests" / "test_02.png"),
            },
        ],
        language="en",
        style="clean",
        api_url="https://llm.example.com/v1",
        api_key="llm-key",
        model="gpt-test",
        vlm_api_url="https://vlm.example.com/v1",
        vlm_api_key="vlm-key",
        vlm_model="gpt-4.1-mini",
        image_api_url="https://image.example.com/v1",
        image_api_key="image-key",
        image_model="gemini-3-pro-image-preview",
    )

    artifacts = run_from_pagecontent(request)

    run_dir = Path(artifacts.run_dir)
    assert run_dir == result_dir / "code_runtime"
    assert Path(artifacts.materials_manifest_path).exists()
    assert Path(artifacts.material_resolution_path).exists()
    assert Path(artifacts.planned_ir_path).exists()
    assert Path(artifacts.final_ir_path).exists()
    assert (run_dir / "ir" / "planned" / "slide_briefs.json").exists()
    assert (run_dir / "ir" / "planned" / "deck_stage.json").exists()
    assert (run_dir / "ir" / "planned" / "slides" / "slide_001.json").exists()
    assert Path(artifacts.pptx_path).exists()

    planned_ir = json.loads(Path(artifacts.planned_ir_path).read_text(encoding="utf-8"))
    final_ir = json.loads(Path(artifacts.final_ir_path).read_text(encoding="utf-8"))
    materials_manifest = json.loads(Path(artifacts.materials_manifest_path).read_text(encoding="utf-8"))
    material_resolution = json.loads(Path(artifacts.material_resolution_path).read_text(encoding="utf-8"))

    assert planned_ir["slides"]
    assert final_ir["slides"]
    assert materials_manifest["assets"]
    assert "requests" in material_resolution
    assert final_ir["slides"][1]["selected_asset_path"]


def test_enable_llm_codegen_writes_code_artifacts(tmp_path: Path) -> None:
    result_dir = tmp_path / "runtime-codegen"
    request = EditablePPTRunRequest(
        result_path=str(result_dir),
        pagecontent=[
            {
                "title": "Codegen Slide",
                "layout_description": "A single generated slide.",
                "key_points": ["One", "Two"],
            }
        ],
        language="en",
        style="clean",
        enable_llm_codegen=True,
    )

    artifacts = run_from_pagecontent(request)

    generated_root = Path(artifacts.run_dir) / "code" / "generated"
    assert (generated_root / "build_deck.py").exists()
    assert (generated_root / "slides").exists()


def test_enable_llm_codegen_persists_validated_slide_cache(monkeypatch, tmp_path: Path) -> None:
    class FakeClient:
        def __init__(self, *, api_key, api_base, model, client_type="llm"):
            self.client_type = client_type

        def chat(self, messages, temperature=0.2, response_format=None):
            return """```python
def slide_001(prs, deck_ir, slide_ir, materials):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 3000000, 800000)
    textbox.text_frame.text = "Validated Codegen"
    return slide
```"""

    def fake_render_pptx(deck_ir, output_path):
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(b"pptx")

    monkeypatch.setattr(runner, "LLMClient", FakeClient)
    monkeypatch.setattr(runner, "render_pptx", fake_render_pptx)

    result_dir = tmp_path / "runtime-codegen-agent"
    request = EditablePPTRunRequest(
        result_path=str(result_dir),
        pagecontent=[
            {
                "title": "Codegen Agent Slide",
                "layout_description": "A generated slide with validated cache.",
                "key_points": ["One", "Two"],
            }
        ],
        language="en",
        style="clean",
        model="gpt-test",
        api_url="https://llm.example.com/v1",
        api_key="llm-key",
        enable_agent_planner=False,
        enable_material_resolution=False,
        enable_llm_codegen=True,
        include_pdf_preview=False,
    )

    artifacts = run_from_pagecontent(request)

    cache_payload = json.loads(
        (Path(artifacts.run_dir) / "code" / "generated" / "cache" / "slide_001.json").read_text(encoding="utf-8")
    )

    assert cache_payload["validated"] is True
    assert cache_payload["attempt_count"] == 1
    assert cache_payload["function_name"] == "slide_001"


def test_enable_llm_codegen_executes_generated_build_deck_for_final_pptx(monkeypatch, tmp_path: Path) -> None:
    observed: dict[str, object] = {}

    class FakeClient:
        def __init__(self, *, api_key, api_base, model, client_type="llm"):
            self.client_type = client_type

        def chat(self, messages, temperature=0.2, response_format=None):
            return """```python
def slide_001(prs, deck_ir, slide_ir, materials):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 3000000, 800000)
    textbox.text_frame.text = "LLM Runtime Slide"
    return slide
```"""

    def fake_render_pptx(deck_ir, output_path):
        observed["renderer_called"] = True
        output_path.parent.mkdir(parents=True, exist_ok=True)
        Presentation().save(output_path)

    monkeypatch.setattr(runner, "LLMClient", FakeClient)
    monkeypatch.setattr(runner, "render_pptx", fake_render_pptx)

    result_dir = tmp_path / "runtime-codegen-exec"
    request = EditablePPTRunRequest(
        result_path=str(result_dir),
        pagecontent=[{"title": "Generated deck", "key_points": ["One", "Two"]}],
        language="en",
        style="clean",
        model="gpt-test",
        api_url="https://llm.example.com/v1",
        api_key="llm-key",
        enable_agent_planner=False,
        enable_material_resolution=False,
        enable_llm_codegen=True,
        include_pdf_preview=False,
    )

    artifacts = run_from_pagecontent(request)

    assert observed.get("renderer_called") is not True
    assert Path(artifacts.pptx_path).exists()
    presentation = Presentation(artifacts.pptx_path)
    assert len(presentation.slides) == 1
    slide_texts = [shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")]
    assert any("LLM Runtime Slide" in text for text in slide_texts)


def test_enable_llm_codegen_falls_back_to_renderer_when_build_deck_execution_fails(monkeypatch, tmp_path: Path) -> None:
    observed: dict[str, object] = {}

    class FakeClient:
        def __init__(self, *, api_key, api_base, model, client_type="llm"):
            self.client_type = client_type

        def chat(self, messages, temperature=0.2, response_format=None):
            return """```python
def slide_001(prs, deck_ir, slide_ir, materials):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 3000000, 800000)
    textbox.text_frame.text = "Codegen Slide Before Deck Failure"
    return slide
```"""

    def fake_render_pptx(deck_ir, output_path):
        observed["renderer_called"] = True
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(0, 0, 3000000, 800000)
        textbox.text_frame.text = "Fallback Renderer Slide"
        presentation.save(output_path)

    def fake_execute_build_deck(build_deck_path, output_path):
        raise RuntimeError("build deck execution failed")

    monkeypatch.setattr(runner, "LLMClient", FakeClient)
    monkeypatch.setattr(runner, "render_pptx", fake_render_pptx)
    monkeypatch.setattr(runner, "execute_build_deck", fake_execute_build_deck)

    result_dir = tmp_path / "runtime-codegen-fallback"
    request = EditablePPTRunRequest(
        result_path=str(result_dir),
        pagecontent=[{"title": "Fallback deck", "key_points": ["One", "Two"]}],
        language="en",
        style="clean",
        model="gpt-test",
        api_url="https://llm.example.com/v1",
        api_key="llm-key",
        enable_agent_planner=False,
        enable_material_resolution=False,
        enable_llm_codegen=True,
        include_pdf_preview=False,
    )

    artifacts = run_from_pagecontent(request)

    assert observed["renderer_called"] is True
    presentation = Presentation(artifacts.pptx_path)
    slide_texts = [shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")]
    assert any("Fallback Renderer Slide" in text for text in slide_texts)
    log_text = Path(artifacts.log_path).read_text(encoding="utf-8")
    assert "codegen execution fallback" in log_text


def test_run_from_pagecontent_writes_log_when_ir_build_fails(tmp_path: Path, monkeypatch) -> None:
    result_dir = tmp_path / "runtime-failure"
    request = EditablePPTRunRequest(
        result_path=str(result_dir),
        pagecontent=[{"title": "Broken Slide"}],
        language="en",
        style="clean",
    )

    def fake_pagecontent_to_deck_ir(pagecontent, *, language, style, asset_base_dir=""):
        raise RuntimeError("planner exploded")

    monkeypatch.setattr(runner, "pagecontent_to_deck_ir", fake_pagecontent_to_deck_ir)

    with pytest.raises(RuntimeError, match="planner exploded"):
        run_from_pagecontent(request)

    log_path = result_dir / "code_runtime" / "logs" / "run.log"
    assert log_path.exists()
    log_text = log_path.read_text(encoding="utf-8")
    assert "run_from_pagecontent started" in log_text
    assert "run failed: RuntimeError: planner exploded" in log_text


def test_run_from_pagecontent_resolves_relative_image_assets_from_result_path(tmp_path: Path) -> None:
    result_dir = tmp_path / "runtime-relative-assets"
    image_dir = result_dir / "images"
    image_dir.mkdir(parents=True, exist_ok=True)
    relative_image_path = image_dir / "diagram.png"
    relative_image_path.write_bytes((PROJECT_ROOT / "tests" / "test_02.png").read_bytes())

    request = EditablePPTRunRequest(
        result_path=str(result_dir),
        pagecontent=[
            {
                "title": "Relative Asset Slide",
                "ppt_img_path": "images/diagram.png",
            }
        ],
        language="en",
        style="clean",
    )

    artifacts = run_from_pagecontent(request)
    presentation = Presentation(artifacts.pptx_path)

    assert len(presentation.slides) == 1
    assert any(hasattr(shape, "image") for shape in presentation.slides[0].shapes)


def test_run_from_pagecontent_persists_runtime_artifact_contents(tmp_path: Path, monkeypatch) -> None:
    result_dir = tmp_path / "runtime-result"
    request = EditablePPTRunRequest(
        result_path=str(result_dir),
        pagecontent=[
            {
                "title": "Editable Runtime",
                "layout_description": "Subtitle for the cover slide.",
                "key_points": ["Deterministic", "Editable output"],
            },
            {
                "title": "Image Focus",
                "layout_description": "Diagram and supporting notes.",
                "key_points": ["One asset", "Short bullets"],
                "asset_ref": str(PROJECT_ROOT / "tests" / "test_02.png"),
            },
        ],
        language="en",
        style="clean",
    )

    def fake_refine_ir(deck_ir: DeckIR) -> DeckIR:
        return deck_ir.model_copy(
            update={
                "title": "Refined Runtime Deck",
                "slides": [
                    deck_ir.slides[0].model_copy(
                        update={
                            "title": "Refined Runtime Cover",
                            "summary": "Refined subtitle for persisted IR assertions.",
                        }
                    )
                ],
            }
        )

    monkeypatch.setattr(runner, "refine_ir", fake_refine_ir)

    artifacts = run_from_pagecontent(request)

    assert Path(artifacts.run_dir) == result_dir / "code_runtime"
    assert Path(artifacts.ir_path).exists()
    assert Path(artifacts.recipe_path).exists()
    assert Path(artifacts.pptx_path).exists()
    assert artifacts.pdf_path == "" or Path(artifacts.pdf_path).exists()
    assert Path(artifacts.log_path).exists()

    ir_payload = json.loads(Path(artifacts.ir_path).read_text(encoding="utf-8"))
    recipe_payload = json.loads(Path(artifacts.recipe_path).read_text(encoding="utf-8"))
    log_text = Path(artifacts.log_path).read_text(encoding="utf-8")

    assert ir_payload["title"] == "Refined Runtime Deck"
    assert len(ir_payload["slides"]) == 1
    assert ir_payload["slides"][0]["title"] == "Refined Runtime Cover"
    assert recipe_payload["deck_title"] == "Refined Runtime Deck"
    assert recipe_payload["slide_count"] == 1
    assert recipe_payload["slides"][0]["title"] == "Refined Runtime Cover"
    assert "wrote deck ir to" in log_text
    assert "wrote render recipe to" in log_text
    assert "wrote editable pptx to" in log_text

    presentation = Presentation(artifacts.pptx_path)
    assert len(presentation.slides) == 1
    slide_texts = [
        shape.text
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text")
    ]
    assert any("Refined Runtime Cover" in text for text in slide_texts)


def test_run_from_pagecontent_without_callback_matches_noop_callback(tmp_path: Path, monkeypatch):
    """callback=None vs noop callback must produce identical artifacts."""
    monkeypatch.setattr(runner, "_build_llm_client", lambda req: None)
    monkeypatch.setattr(runner, "_build_vlm_client", lambda req: None)
    monkeypatch.setattr(runner, "_export_pdf_preview", lambda *a, **k: "")

    # Run A and B use the same result_path so path-derived fields (asset_base_dir, artifact
    # paths) don't diverge. We capture A's final_ir contents before B's run overwrites them.
    result_path = str(tmp_path / "run")

    req_a = EditablePPTRunRequest(
        result_path=result_path,
        pagecontent=[{"title": "P1", "key_points": ["k"], "layout_description": ""}],
    )
    artifacts_a = runner.run_from_pagecontent(req_a)
    import json as _json
    final_ir_a = _json.loads(Path(artifacts_a.final_ir_path).read_text(encoding="utf-8"))
    artifacts_a_paths = artifacts_a.model_dump()
    pptx_a_exists = Path(artifacts_a.pptx_path).exists()

    req_b = EditablePPTRunRequest(
        result_path=result_path,
        pagecontent=[{"title": "P1", "key_points": ["k"], "layout_description": ""}],
    )
    artifacts_b = runner.run_from_pagecontent(req_b, progress_callback=lambda name, payload: None)
    final_ir_b = _json.loads(Path(artifacts_b.final_ir_path).read_text(encoding="utf-8"))

    assert artifacts_a_paths == artifacts_b.model_dump()
    assert final_ir_a == final_ir_b
    assert pptx_a_exists
    assert Path(artifacts_b.pptx_path).exists()


def test_run_from_pagecontent_emits_callback_sequence(tmp_path: Path, monkeypatch):
    monkeypatch.setattr(runner, "_build_llm_client", lambda req: None)
    monkeypatch.setattr(runner, "_build_vlm_client", lambda req: None)
    monkeypatch.setattr(runner, "_export_pdf_preview", lambda *a, **k: "")

    from vendor.presentagent_runtime.renderer.artifact_store import build_artifacts

    expected_artifacts = build_artifacts(str(tmp_path))

    events: list[tuple[str, dict]] = []
    req = EditablePPTRunRequest(
        result_path=str(tmp_path),
        pagecontent=[{"title": "P1"}, {"title": "P2"}],
    )
    runner.run_from_pagecontent(req, progress_callback=lambda name, payload: events.append((name, payload)))

    event_names = [name for name, _ in events]
    assert event_names[:2] == ["planning_done", "final_ir_done"], event_names

    planning_payload = events[0][1]
    assert planning_payload["planned_ir_path"] == expected_artifacts.planned_ir_path
    assert planning_payload["materials_manifest_path"] == expected_artifacts.materials_manifest_path
    assert planning_payload["material_resolution_path"] == expected_artifacts.material_resolution_path

    final_payload = events[1][1]
    assert final_payload["final_ir_path"] == expected_artifacts.final_ir_path


def test_render_pptx_per_slide_writes_one_file_per_slide(tmp_path: Path, monkeypatch):
    monkeypatch.setattr(runner, "_build_llm_client", lambda req: None)
    monkeypatch.setattr(runner, "_build_vlm_client", lambda req: None)
    monkeypatch.setattr(runner, "_export_pdf_preview", lambda *a, **k: "")

    req = EditablePPTRunRequest(
        result_path=str(tmp_path),
        pagecontent=[{"title": "P1"}, {"title": "P2"}, {"title": "P3"}],
    )
    artifacts = runner.run_from_pagecontent(req)

    slides_dir = Path(artifacts.slides_dir)
    assert slides_dir.exists()
    assert (slides_dir / "slide_000.pptx").exists()
    assert (slides_dir / "slide_001.pptx").exists()
    assert (slides_dir / "slide_002.pptx").exists()
    assert len(artifacts.slide_artifacts) == 3
    assert artifacts.slide_artifacts[0].index == 0
    assert artifacts.slide_artifacts[0].pptx_path.endswith("slide_000.pptx")


def test_run_from_pagecontent_emits_slide_rendered_events(tmp_path: Path, monkeypatch):
    monkeypatch.setattr(runner, "_build_llm_client", lambda req: None)
    monkeypatch.setattr(runner, "_build_vlm_client", lambda req: None)
    monkeypatch.setattr(runner, "_export_pdf_preview", lambda *a, **k: "")

    events: list[tuple[str, dict]] = []
    req = EditablePPTRunRequest(
        result_path=str(tmp_path),
        pagecontent=[{"title": "P1"}, {"title": "P2"}],
    )
    runner.run_from_pagecontent(req, progress_callback=lambda name, payload: events.append((name, payload)))

    slide_events = [e for e in events if e[0] == "slide_rendered"]
    assert len(slide_events) == 2
    assert slide_events[0][1]["index"] == 0
    assert slide_events[0][1]["total"] == 2
    assert slide_events[0][1]["pptx_path"].endswith("slide_000.pptx")
    rendering_done = [e for e in events if e[0] == "rendering_done"]
    assert len(rendering_done) == 1
    assert rendering_done[0][1]["total"] == 2
    assert len(rendering_done[0][1]["slide_artifacts"]) == 2

    # Invariance: rendering_done.slide_artifacts[i] must match slide_rendered[i]'s payload for the same index
    final_list = rendering_done[0][1]["slide_artifacts"]
    slide_rendered_payloads = [e[1] for e in slide_events]
    for i, final_item in enumerate(final_list):
        assert final_item["index"] == slide_rendered_payloads[i]["index"]
        assert final_item["pptx_path"] == slide_rendered_payloads[i]["pptx_path"]
        assert final_item["slide_id"] == slide_rendered_payloads[i]["slide_id"]
        assert final_item["title"] == slide_rendered_payloads[i]["title"]


def test_render_slide_preview_returns_empty_when_libreoffice_missing(tmp_path: Path, monkeypatch):
    monkeypatch.setattr(runner.shutil, "which", lambda name: None)
    result = runner._render_slide_preview(
        slide_pptx=tmp_path / "slide_000.pptx",
        preview_png=tmp_path / "slide_000.png",
    )
    assert result == ""


def test_first_slide_preview_arrives_in_slide_rendered_event(tmp_path: Path, monkeypatch):
    monkeypatch.setattr(runner, "_build_llm_client", lambda req: None)
    monkeypatch.setattr(runner, "_build_vlm_client", lambda req: None)
    monkeypatch.setattr(runner, "_export_pdf_preview", lambda *a, **k: "")
    # Fake preview: write a placeholder png file and return its path
    def fake_preview(slide_pptx, preview_png):
        preview_png.write_bytes(b"PNG")
        return str(preview_png)
    monkeypatch.setattr(runner, "_render_slide_preview", fake_preview)

    events: list[tuple[str, dict]] = []
    req = EditablePPTRunRequest(
        result_path=str(tmp_path),
        pagecontent=[{"title": "P1"}, {"title": "P2"}, {"title": "P3"}],
    )
    runner.run_from_pagecontent(req, progress_callback=lambda name, payload: events.append((name, payload)))

    slide_events = [e for e in events if e[0] == "slide_rendered"]
    # First slide must carry preview png path; others are empty on their slide_rendered event
    assert slide_events[0][1]["preview_png_path"].endswith("slide_000.png")
    assert slide_events[1][1]["preview_png_path"] == ""
    assert slide_events[2][1]["preview_png_path"] == ""
    # rendering_done should have all 3 preview_png_path populated (batch step)
    rendering_done = [e for e in events if e[0] == "rendering_done"][0]
    previews = [s["preview_png_path"] for s in rendering_done[1]["slide_artifacts"]]
    assert len(previews) == 3
    assert all(p.endswith(".png") for p in previews)


def test_render_slide_preview_returns_empty_on_soffice_failure(tmp_path: Path, monkeypatch):
    monkeypatch.setattr(runner.shutil, "which", lambda name: "/usr/bin/soffice")
    class _FakeResult:
        returncode = 1
        stderr = ""
        stdout = ""
    monkeypatch.setattr(runner.subprocess, "run", lambda *a, **k: _FakeResult())
    result = runner._render_slide_preview(
        slide_pptx=tmp_path / "slide_000.pptx",
        preview_png=tmp_path / "slide_000.png",
    )
    assert result == ""


def test_assemble_deck_from_final_ir_produces_full_pptx(tmp_path: Path, monkeypatch):
    monkeypatch.setattr(runner, "_build_llm_client", lambda req: None)
    monkeypatch.setattr(runner, "_build_vlm_client", lambda req: None)
    monkeypatch.setattr(runner, "_export_pdf_preview", lambda *a, **k: "")

    req = EditablePPTRunRequest(
        result_path=str(tmp_path),
        pagecontent=[{"title": "P1"}, {"title": "P2"}],
    )
    run_artifacts = runner.run_from_pagecontent(req)

    output_pptx = tmp_path / "assembled" / "final.pptx"
    events: list[tuple[str, dict]] = []
    result = runner.assemble_deck_from_final_ir(
        final_ir_path=Path(run_artifacts.final_ir_path),
        output_pptx=output_pptx,
        include_pdf_preview=False,
        progress_callback=lambda name, payload: events.append((name, payload)),
    )
    assert output_pptx.exists()
    from pptx import Presentation
    assert len(Presentation(str(output_pptx)).slides) == 2
    assert result["pptx_path"] == str(output_pptx)
    assert [name for name, _ in events] == ["exporting_done"]
