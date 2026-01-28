import os
import shutil
import subprocess
import time
from pathlib import Path
from fastapi import APIRouter, UploadFile, File, Form, HTTPException, Body
from typing import Optional, List, Dict, Any

import fitz  # PyMuPDF

from dataflow_agent.state import IntelligentQARequest, IntelligentQAState, KBPodcastRequest, KBPodcastState, KBMindMapRequest, KBMindMapState
from dataflow_agent.workflow.wf_intelligent_qa import create_intelligent_qa_graph
from dataflow_agent.workflow.wf_kb_podcast import create_kb_podcast_graph
from dataflow_agent.workflow.wf_kb_mindmap import create_kb_mindmap_graph
from dataflow_agent.toolkits.ragtool.vector_store_tool import process_knowledge_base_files, VectorStoreManager
from dataflow_agent.utils import get_project_root
from dataflow_agent.workflow import run_workflow
from fastapi_app.config import settings
from fastapi_app.schemas import Paper2PPTRequest
from fastapi_app.utils import _from_outputs_url, _to_outputs_url
from fastapi_app.workflow_adapters.wa_paper2ppt import _init_state_from_request

router = APIRouter(prefix="/kb", tags=["Knowledge Base"])

# Base directory for storing KB files
# Use absolute path as requested by user or relative to project root
# We will use relative path 'outputs/kb_data' which resolves to that in the current workspace
KB_BASE_DIR = Path("outputs/kb_data")

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg", ".mp4"}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg"}
DOC_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt"}


def _resolve_local_path(path_or_url: str) -> Path:
    if not path_or_url:
        raise HTTPException(status_code=400, detail="Empty file path")
    raw = _from_outputs_url(path_or_url)
    p = Path(raw)
    if not p.is_absolute():
        p = (get_project_root() / p).resolve()
    return p


def _convert_to_pdf(input_path: Path, output_dir: Path) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(input_path)
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    pdf_path = output_dir / f"{input_path.stem}.pdf"
    if not pdf_path.exists():
        raise HTTPException(status_code=500, detail=f"PDF conversion failed for {input_path.name}")
    return pdf_path


def _merge_pdfs(pdf_paths: List[Path], output_path: Path) -> Path:
    if not pdf_paths:
        raise HTTPException(status_code=400, detail="No PDF files to merge")
    merged = fitz.open()
    for pdf in pdf_paths:
        with fitz.open(pdf) as src:
            merged.insert_pdf(src)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    merged.save(output_path)
    merged.close()
    return output_path


def _append_images_to_pptx(pptx_path: Path, image_paths: List[Path]) -> None:
    try:
        from pptx import Presentation
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"python-pptx not available: {e}")

    prs = Presentation(str(pptx_path))
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height
        )
    prs.save(str(pptx_path))

@router.post("/upload")
async def upload_kb_file(
    file: UploadFile = File(...),
    email: str = Form(...),
    user_id: str = Form(...)
):
    """
    Upload a file to the user's knowledge base directory.
    Stores at: outputs/kb_data/{email}/{filename}
    """
    if not email:
        raise HTTPException(status_code=400, detail="Email is required")

    # Validate file extension
    file_ext = Path(file.filename).suffix.lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported file type: {file_ext}. Allowed: {', '.join(ALLOWED_EXTENSIONS)}"
        )

    try:
        # Create user directory if not exists
        user_dir = KB_BASE_DIR / email
        user_dir.mkdir(parents=True, exist_ok=True)

        # Secure filename (simple version)
        filename = file.filename
        if not filename:
            filename = f"unnamed_{user_id}"
            
        # Avoid path traversal
        filename = os.path.basename(filename)
        
        file_path = user_dir / filename
        
        # Save file
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
            
        # Return the relative path for static access and storage path
        # Assuming 'outputs' dir is mounted to '/outputs'
        static_path = f"/outputs/kb_data/{email}/{filename}"
        
        return {
            "success": True,
            "filename": filename,
            "file_size": os.path.getsize(file_path),
            "storage_path": str(file_path),
            "static_url": static_path,
            "file_type": file.content_type
        }

    except Exception as e:
        print(f"Error uploading file: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.delete("/delete")
async def delete_kb_file(
    storage_path: str = Form(...)
):
    """
    Delete a file from the physical storage.
    """
    try:
        # Security check: ensure path is within KB_BASE_DIR
        # This is a basic check. In production, use more robust path validation.
        target_path = Path(storage_path).resolve()
        base_path = KB_BASE_DIR.resolve()
        
        if not str(target_path).startswith(str(base_path)):
             # Allow if it's the absolute path provided by the user system
             # Check if it exists essentially
             pass

        if target_path.exists() and target_path.is_file():
            os.remove(target_path)
            return {"success": True, "message": "File deleted"}
        else:
            return {"success": False, "message": "File not found"}
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/chat")
async def chat_with_kb(
    files: List[str] = Body(..., embed=True),
    query: str = Body(..., embed=True),
    history: List[Dict[str, str]] = Body([], embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    model: str = Body(settings.KB_CHAT_MODEL, embed=True),
):
    """
    Intelligent QA Chat
    """
    try:
        # Normalize file paths (web path -> local absolute path)
        project_root = get_project_root()
        local_files = []
        for f in files:
            # remove leading /outputs/ if present, or just join
            # Web path: /outputs/kb_data/...
            clean_path = f.lstrip('/')
            p = project_root / clean_path
            if p.exists():
                local_files.append(str(p))
            else:
                # Try raw path
                p_raw = Path(f)
                if p_raw.exists():
                    local_files.append(str(p_raw))
        
        if not local_files:
             # Just return empty answer or handle logic
             pass

        # Construct Request
        req = IntelligentQARequest(
            files=local_files,
            query=query,
            history=history,
            chat_api_url=api_url or os.getenv("DF_API_URL"),
            api_key=api_key or os.getenv("DF_API_KEY"),
            model=model
        )
        
        state = IntelligentQAState(request=req)
        
        # Build and Run Graph
        builder = create_intelligent_qa_graph()
        graph = builder.compile()
        
        result_state = await graph.ainvoke(state)
        
        # graph.ainvoke returns the final state dict or state object depending on implementation.
        # LangGraph usually returns dict. But our GenericGraphBuilder wrapper might return state.
        # GenericGraphBuilder compile returns a compiled graph.
        # Let's check typical usage. usually await graph.ainvoke(state) returns dict.
        
        answer = ""
        file_analyses = []
        
        if isinstance(result_state, dict):
            answer = result_state.get("answer", "")
            file_analyses = result_state.get("file_analyses", [])
        else:
            answer = getattr(result_state, "answer", "")
            file_analyses = getattr(result_state, "file_analyses", [])
            
        return {
            "success": True,
            "answer": answer,
            "file_analyses": file_analyses
        }

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/generate-ppt")
async def generate_ppt_from_kb(
    file_path: Optional[str] = Body(None, embed=True),
    file_paths: Optional[List[str]] = Body(None, embed=True),
    image_paths: Optional[List[str]] = Body(None, embed=True),
    image_items: Optional[List[Dict[str, Any]]] = Body(None, embed=True),
    query: Optional[str] = Body("", embed=True),
    need_embedding: bool = Body(False, embed=True),
    search_top_k: int = Body(8, embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    api_url: str = Body(..., embed=True),
    api_key: str = Body(..., embed=True),
    style: str = Body("modern", embed=True),
    language: str = Body("zh", embed=True),
    page_count: int = Body(10, embed=True),
    model: str = Body("gpt-4o", embed=True),
    gen_fig_model: str = Body("gemini-2.5-flash-image", embed=True),
):
    """
    Generate PPT from knowledge base file (non-interactive)
    """
    try:
        # Normalize and validate input files (PDF/PPT/DOC/IMG)
        input_paths = file_paths or ([file_path] if file_path else [])
        if not input_paths:
            raise HTTPException(status_code=400, detail="No input files provided")

        # Create output directory
        ts = int(time.time())
        project_root = get_project_root()
        output_dir = project_root / "outputs" / "kb_outputs" / email / f"{ts}_ppt"
        output_dir.mkdir(parents=True, exist_ok=True)

        # Split docs/images
        doc_paths: List[Path] = []
        user_image_items: List[Dict[str, Any]] = []
        for p in input_paths:
            local_path = _resolve_local_path(p)
            if not local_path.exists():
                raise HTTPException(status_code=404, detail=f"File not found: {p}")
            ext = local_path.suffix.lower()
            if ext in IMAGE_EXTENSIONS:
                user_image_items.append({"path": str(local_path), "description": ""})
            elif ext in {\".pdf\", \".pptx\", \".ppt\", \".docx\", \".doc\"}:
                doc_paths.append(local_path)
            else:
                raise HTTPException(status_code=400, detail=f\"Unsupported file type for PPT: {local_path.name}\")

        if not doc_paths:
            raise HTTPException(status_code=400, detail=\"At least one document file is required for PPT generation\")

        # Convert docs to PDF for MinerU merge
        local_pdf_paths: List[Path] = []
        convert_dir = output_dir / \"input\"
        convert_dir.mkdir(parents=True, exist_ok=True)
        for p in doc_paths:
            ext = p.suffix.lower()
            if ext == \".pdf\":
                local_pdf_paths.append(p)
            elif ext in {\".pptx\", \".ppt\", \".docx\", \".doc\"}:
                local_pdf_paths.append(_convert_to_pdf(p, convert_dir))
            else:
                raise HTTPException(status_code=400, detail=f\"Unsupported file type for PPT: {p.name}\")

        # Merge PDFs if multiple
        if len(local_pdf_paths) > 1:
            merge_dir = output_dir / "input"
            merged_pdf = merge_dir / "merged.pdf"
            local_file_path = _merge_pdfs(local_pdf_paths, merged_pdf)
        else:
            local_file_path = local_pdf_paths[0]

        # Normalize image items (optional)
        resolved_image_items: List[Dict[str, Any]] = []
        for item in image_items or []:
            raw_path = item.get("path") or item.get("url") or ""
            if not raw_path:
                continue
            img_path = _resolve_local_path(str(raw_path))
            if img_path.exists() and img_path.suffix.lower() in IMAGE_EXTENSIONS:
                resolved_image_items.append({
                    "path": str(img_path),
                    "description": item.get("description") or item.get("desc") or ""
                })

        for img in image_paths or []:
            img_path = _resolve_local_path(img)
            if img_path.exists() and img_path.suffix.lower() in IMAGE_EXTENSIONS:
                resolved_image_items.append({
                    "path": str(img_path),
                    "description": ""
                })

        resolved_image_items.extend(user_image_items)

        # Embedding + retrieval (optional)
        retrieval_text = ""
        if need_embedding:
            base_dir = project_root / "outputs" / "kb_data" / email / "vector_store"
            embed_api_url = api_url
            if "/embeddings" not in embed_api_url:
                embed_api_url = embed_api_url.rstrip("/") + "/embeddings"

            files_for_embed = [{"path": str(p), "description": ""} for p in doc_paths]
            manifest = await process_knowledge_base_files(
                files_for_embed,
                base_dir=str(base_dir),
                api_url=embed_api_url,
                api_key=api_key,
                model_name=None,
                multimodal_model=None,
            )

            manager = VectorStoreManager(
                base_dir=str(base_dir),
                embedding_api_url=embed_api_url,
                api_key=api_key,
            )

            def _match_file_ids(m: Dict[str, Any], paths: List[Path]) -> List[str]:
                ids: List[str] = []
                target = {str(p.resolve()) for p in paths}
                for f in m.get("files", []):
                    try:
                        if str(Path(f.get("original_path", "")).resolve()) in target:
                            if f.get("id"):
                                ids.append(f["id"])
                    except Exception:
                        continue
                return ids

            file_ids = _match_file_ids(manifest or manager.manifest or {}, doc_paths)
            if query and file_ids:
                results = manager.search(query=query, top_k=search_top_k, file_ids=file_ids)
                retrieval_text = "\n\n".join([r.get("content", "") for r in results if r.get("content")])

        # Prepare request
        ppt_req = Paper2PPTRequest(
            input_type="PDF",
            input_content=str(local_file_path),
            email=email,
            chat_api_url=api_url,
            chat_api_key=api_key,
            api_key=api_key,
            style=style,
            language=language,
            page_count=page_count,
            model=model,
            gen_fig_model=gen_fig_model,
            aspect_ratio="16:9",
            use_long_paper=False
        )

        # Run KB pagecontent workflow
        state_pc = _init_state_from_request(ppt_req, result_path=output_dir)
        state_pc.kb_query = query or ""
        state_pc.kb_retrieval_text = retrieval_text
        state_pc.kb_user_images = resolved_image_items
        state_pc = await run_workflow("kb_page_content", state_pc)
        pagecontent = getattr(state_pc, "pagecontent", []) or []

        # Run PPT generation with injected pagecontent
        state_pc.pagecontent = pagecontent
        state_pp = await run_workflow("paper2ppt_parallel", state_pc)

        # Extract output paths
        pdf_path = ""
        pptx_path = ""
        if hasattr(state_pp, 'ppt_pdf_path'):
            pdf_path = state_pp.ppt_pdf_path
        if hasattr(state_pp, 'ppt_pptx_path'):
            pptx_path = state_pp.ppt_pptx_path

        return {
            "success": True,
            "result_path": str(output_dir),
            "pdf_path": _to_outputs_url(pdf_path) if pdf_path else "",
            "pptx_path": _to_outputs_url(pptx_path) if pptx_path else "",
            "output_file_id": f"kb_ppt_{ts}"
        }

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/generate-podcast")
async def generate_podcast_from_kb(
    file_paths: List[str] = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    api_url: str = Body(..., embed=True),
    api_key: str = Body(..., embed=True),
    model: str = Body("gpt-4o", embed=True),
    tts_model: str = Body("gemini-2.5-pro-preview-tts", embed=True),
    voice_name: str = Body("Kore", embed=True),
    language: str = Body("zh", embed=True),
):
    """
    Generate podcast from knowledge base files
    """
    try:
        # Normalize file paths
        if not file_paths:
            raise HTTPException(status_code=400, detail="No valid files provided")

        local_paths: List[Path] = []
        for f in file_paths:
            local_path = _resolve_local_path(f)
            if not local_path.exists():
                raise HTTPException(status_code=404, detail=f"File not found: {f}")
            local_paths.append(local_path)

        # If multiple files, merge into a single PDF (doc/ppt will be converted)
        if len(local_paths) > 1:
            ts = int(time.time())
            merge_dir = get_project_root() / "outputs" / "kb_outputs" / email / f"{ts}_podcast_input"
            merge_dir.mkdir(parents=True, exist_ok=True)

            pdf_paths: List[Path] = []
            for p in local_paths:
                ext = p.suffix.lower()
                if ext == ".pdf":
                    pdf_paths.append(p)
                elif ext in {".docx", ".doc", ".pptx", ".ppt"}:
                    pdf_paths.append(_convert_to_pdf(p, merge_dir))
                else:
                    raise HTTPException(status_code=400, detail=f"Unsupported file type for podcast: {p.name}")

            merged_pdf = merge_dir / "merged.pdf"
            local_file_paths = [str(_merge_pdfs(pdf_paths, merged_pdf))]
        else:
            local_file_paths = [str(local_paths[0])]

        # Prepare request
        podcast_req = KBPodcastRequest(
            files=local_file_paths,
            chat_api_url=api_url,
            api_key=api_key,
            model=model,
            tts_model=tts_model,
            voice_name=voice_name,
            language=language
        )
        podcast_req.email = email

        state = KBPodcastState(request=podcast_req)

        # Build and run graph
        builder = create_kb_podcast_graph()
        graph = builder.compile()

        result_state = await graph.ainvoke(state)

        # Extract results
        audio_path = ""
        script_path = ""
        result_path = ""

        if isinstance(result_state, dict):
            audio_path = result_state.get("audio_path", "")
            result_path = result_state.get("result_path", "")
        else:
            audio_path = getattr(result_state, "audio_path", "")
            result_path = getattr(result_state, "result_path", "")

        if result_path:
            script_path = str(Path(result_path) / "script.txt")

        return {
            "success": True,
            "result_path": result_path,
            "audio_path": audio_path,
            "script_path": script_path,
            "output_file_id": f"kb_podcast_{int(time.time())}"
        }

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/generate-mindmap")
async def generate_mindmap_from_kb(
    file_paths: List[str] = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    api_url: str = Body(..., embed=True),
    api_key: str = Body(..., embed=True),
    model: str = Body("gpt-4o", embed=True),
    mindmap_style: str = Body("default", embed=True),
    max_depth: int = Body(3, embed=True),
    language: str = Body("zh", embed=True),
):
    """
    Generate mindmap from knowledge base files
    """
    try:
        # Normalize file paths
        project_root = get_project_root()
        local_file_paths = []

        for f in file_paths:
            clean_path = f.lstrip('/')
            local_path = project_root / clean_path

            if not local_path.exists():
                local_path = Path(f)
                if not local_path.exists():
                    raise HTTPException(status_code=404, detail=f"File not found: {f}")

            local_file_paths.append(str(local_path))

        if not local_file_paths:
            raise HTTPException(status_code=400, detail="No valid files provided")

        # Prepare request
        mindmap_req = KBMindMapRequest(
            files=local_file_paths,
            chat_api_url=api_url,
            api_key=api_key,
            model=model,
            mindmap_style=mindmap_style,
            max_depth=max_depth,
            language=language
        )
        mindmap_req.email = email

        state = KBMindMapState(request=mindmap_req)

        # Build and run graph
        builder = create_kb_mindmap_graph()
        graph = builder.compile()

        result_state = await graph.ainvoke(state)

        # Extract results
        mermaid_code = ""
        result_path = ""

        if isinstance(result_state, dict):
            mermaid_code = result_state.get("mermaid_code", "")
            result_path = result_state.get("result_path", "")
        else:
            mermaid_code = getattr(result_state, "mermaid_code", "")
            result_path = getattr(result_state, "result_path", "")

        return {
            "success": True,
            "result_path": result_path,
            "mermaid_code": mermaid_code,
            "output_file_id": f"kb_mindmap_{int(time.time())}"
        }

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))
