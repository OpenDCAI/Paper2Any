from __future__ import annotations

import asyncio
from datetime import datetime
import uuid
from pathlib import Path
from typing import Optional

from fastapi import APIRouter, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse

# 简单的邀请码校验：从本地文本文件加载白名单
INVITE_CODES_FILE = Path(__file__).resolve().parent.parent / "invite_codes.txt"


def load_invite_codes() -> set[str]:
    """
    从 invite_codes.txt 中加载邀请码列表。

    文件格式：每行一个邀请码，忽略空行和以 # 开头的注释行。
    """
    codes: set[str] = set()
    if not INVITE_CODES_FILE.exists():
        return codes
    for line in INVITE_CODES_FILE.read_text(encoding="utf-8").splitlines():
        code = line.strip()
        if not code or code.startswith("#"):
            continue
        codes.add(code)
    return codes


VALID_INVITE_CODES = load_invite_codes()


def validate_invite_code(code: str | None) -> None:
    """
    校验邀请码是否有效。无效则抛出 403。
    """
    if not code:
        raise HTTPException(status_code=403, detail="invite_code is required")
    if code not in VALID_INVITE_CODES:
        raise HTTPException(status_code=403, detail="Invalid invite_code")


# 全局信号量：控制重任务并发度（排队机制）
# 目前设为 1，即串行执行；如需并行可调大此值
task_semaphore = asyncio.Semaphore(1)

# 输出根目录：按任务类型 / 时间戳+UUID 组织
BASE_OUTPUT_DIR = Path("outputs")


router = APIRouter()


def create_run_dir(task_type: str) -> Path:
    """
    为一次请求创建独立目录：
        outputs/{task_type}/{timestamp}_{short_uuid}/
    并在其中创建 input/ 与 output/ 子目录。
    """
    ts = datetime.utcnow().strftime("%Y%m%dT%H%M%S")
    rid = uuid.uuid4().hex[:6]
    run_dir = BASE_OUTPUT_DIR / task_type / f"{ts}_{rid}"

    (run_dir / "input").mkdir(parents=True, exist_ok=True)
    (run_dir / "output").mkdir(parents=True, exist_ok=True)

    return run_dir


def create_dummy_pptx(output_path: Path, title: str, content: str) -> None:
    """
    生成一个非常简单的 PPTX 文件，作为占位 / Demo。

    若未来想使用真正的 workflow，可以在此处替换为实际调用逻辑。
    """
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        # 如果未安装 python-pptx，则写入一个占位的二进制文件，保证前端能正常下载
        output_path.write_bytes(b"Dummy PPTX content - please install python-pptx for real PPTX generation.")
        return

    prs = Presentation()
    # 使用标题 + 内容布局（通常为 1）
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

    prs.save(output_path)


@router.post("/paper2graph/generate")
async def generate_paper2graph(
    model_name: str = Form(...),
    chat_api_url: str = Form(...),
    api_key: str = Form(...),
    input_type: str = Form(...),  # 'file' | 'text' | 'image'
    invite_code: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
    file_kind: Optional[str] = Form(None),  # 'pdf' | 'image'
    text: Optional[str] = Form(None),
):
    """
    Paper2Graph 假接口（带邀请码校验）：

    - 需要前端在 FormData 中传入 invite_code，并在本地白名单文件中验证；
    - 接收前端上传的 PDF / 图片 / 文本；
    - 为每次请求在 outputs/paper2graph 下创建独立目录；
    - 使用全局信号量控制重任务串行执行；
    - 返回一个简单的 PPTX 文件，供前端下载测试。

    - 图片比例，16:9
    """
    # 0. 邀请码校验
    validate_invite_code(invite_code)

    # 1. 基础参数校验
    if input_type in ("file", "image"):
        if file is None:
            raise HTTPException(status_code=400, detail="file is required when input_type is 'file' or 'image'")
        if file_kind not in ("pdf", "image"):
            raise HTTPException(status_code=400, detail="file_kind must be 'pdf' or 'image'")
    elif input_type == "text":
        if not text:
            raise HTTPException(status_code=400, detail="text is required when input_type is 'text'")
    else:
        raise HTTPException(status_code=400, detail="invalid input_type, must be one of: file, text, image")

    # 2. 创建本次请求的独立目录
    run_dir = create_run_dir("paper2graph")
    input_dir = run_dir / "input"
    output_dir = run_dir / "output"

    # 3. 保存输入内容到 input/ 目录
    saved_input_name = "unknown"
    if input_type in ("file", "image"):
        original_name = file.filename or "uploaded"
        ext = Path(original_name).suffix or ""
        input_path = input_dir / f"input{ext}"
        content_bytes = await file.read()
        input_path.write_bytes(content_bytes)
        saved_input_name = input_path.name
    else:  # text
        input_path = input_dir / "input.txt"
        input_path.write_text(text or "", encoding="utf-8")
        saved_input_name = input_path.name

    # 4. 重任务段：受信号量保护，确保排队执行
    async with task_semaphore:
        output_pptx = output_dir / "paper2graph.pptx"
        # 这里未来可以替换为真正的 LLM + PDF/图片解析逻辑
        demo_title = "Paper2Graph Demo"
        demo_content = (
            f"model_name: {model_name}\n"
            f"chat_api_url: {chat_api_url}\n"
            f"input_type: {input_type}\n"
            f"file_kind: {file_kind or 'N/A'}\n"
            f"saved_input: {saved_input_name}\n"
        )
        create_dummy_pptx(output_pptx, demo_title, demo_content)

    # 5. 返回 PPTX 文件
    return FileResponse(
        path=output_pptx,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="paper2graph.pptx",
    )


@router.post("/paper2ppt/generate")
async def generate_paper2ppt(
    model_name: str = Form(...),
    chat_api_url: str = Form(...),
    api_key: str = Form(...),
    input_type: str = Form(...),  # 当前前端固定为 'file'
    invite_code: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
    file_kind: Optional[str] = Form(None),  # 当前前端固定为 'pdf'
):
    """
    Paper2PPT 假接口（带邀请码校验）：

    - 需要前端在 FormData 中传入 invite_code，并在本地白名单文件中验证；
    - 接收前端上传的 PDF；
    - 为每次请求在 outputs/paper2ppt 下创建独立目录；
    - 使用全局信号量控制重任务串行执行；
    - 返回一个简单的 PPTX 文件，供前端下载测试。
    """
    # 0. 邀请码校验
    validate_invite_code(invite_code)

    if input_type != "file":
        raise HTTPException(status_code=400, detail="Paper2PPT currently only supports input_type='file'")

    if file is None:
        raise HTTPException(status_code=400, detail="file is required for Paper2PPT")

    if file_kind not in ("pdf", None):
        # 允许 None（前端若未传），否则校验必须为 pdf
        raise HTTPException(status_code=400, detail="file_kind must be 'pdf' for Paper2PPT")

    # 2. 创建本次请求的独立目录
    run_dir = create_run_dir("paper2ppt")
    input_dir = run_dir / "input"
    output_dir = run_dir / "output"

    # 3. 保存输入 PDF
    original_name = file.filename or "uploaded.pdf"
    ext = Path(original_name).suffix or ".pdf"
    input_path = input_dir / f"input{ext}"
    content_bytes = await file.read()
    input_path.write_bytes(content_bytes)
    saved_input_name = input_path.name

    # 4. 重任务段：受信号量保护，确保排队执行
    async with task_semaphore:
        output_pptx = output_dir / "paper2ppt.pptx"
        demo_title = "Paper2PPT Demo"
        demo_content = (
            f"model_name: {model_name}\n"
            f"chat_api_url: {chat_api_url}\n"
            f"input_type: {input_type}\n"
            f"file_kind: {file_kind or 'pdf'}\n"
            f"saved_input: {saved_input_name}\n"
        )
        create_dummy_pptx(output_pptx, demo_title, demo_content)

    # 5. 返回 PPTX 文件
    return FileResponse(
        path=output_pptx,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="paper2ppt.pptx",
    )
